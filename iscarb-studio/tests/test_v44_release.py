"""Regression tests for the reported 20-unit, missing-detail and PDF defects."""
from pathlib import Path
import subprocess
import sys
import tempfile
from unittest.mock import patch

import pymupdf
import pytest
from pptx import Presentation

from app.deterministic_blueprint_fallback import build_deterministic_blueprint
from app.gate_v15 import deterministic_gate
from app.presenter_v44 import (export_presenter_pdf, export_presenter_pptx, teaching_items,
                               text_layout, wrap, title_block, rubric_layout)
from app.source_bundle import SourceBundle, SourceItem
from app.source_profile_fallback import build_deterministic_source_profile
from app.source_visuals import VisualAsset, VisualRegistry, anchor_slides, load_registry
from app.source_visuals_v42 import _looks_like_title_only, plan_for_unit_v42

ROOT=Path(__file__).resolve().parents[1]
LECTURE=ROOT.parent / "lectures/cimt/CPIT455-class2-NooR.pdf"


@pytest.fixture(scope="module")
def source():
    bundle=SourceBundle(items=[SourceItem("primary","P1",LECTURE.name,LECTURE,LECTURE.name)],lecture_focus="",session_minutes=90)
    profile=build_deterministic_source_profile(bundle)
    return profile,build_deterministic_blueprint(profile)


@pytest.mark.parametrize("anchor,expected",[
    ("[P1]",[]),("[P1] chapter 10",[]),("[P1] PAGE 4",[4]),
    ("[P1] SLIDE 40",[40]),("[P1] PAGE 4; [P1] PAGE 7",[4,7]),
    ("[P1] pp. 4–6",[4,5,6]),("[P1] PAGES 5–4",[4,5]),
])
def test_source_id_is_not_a_page(anchor,expected):
    assert anchor_slides(anchor)==expected
    from app.source_visual_patch_v2 import robust_anchor_slides
    assert robust_anchor_slides(anchor)==expected


def test_explicit_missing_page_does_not_borrow_another(source):
    _,bp=source
    u=bp.units[5].model_copy(update={"source_anchor":"[P1] PAGE 99"})
    asset=VisualAsset(4,alt_text="What is dependability? A dependable system provides trusted service. Hardware, software and operational failures all matter.")
    registry=VisualRegistry("local:test.pdf","test",(asset,))
    assert plan_for_unit_v42(bp,u,registry).reuse_mode=="REDRAW"


def test_diagram_only_pages_are_not_mistaken_for_covers():
    assert not _looks_like_title_only(VisualAsset(13,alt_text="AI enhancement",visual_area_ratio=.6))
    assert _looks_like_title_only(VisualAsset(1,alt_text="Chapter 10 Dependable Systems Adeeb Noor PhD IT Department Faculty of Computing King Abdulaziz University Fall 2025"))


def test_scoped_missing_upload_never_falls_back_to_another_job(source,tmp_path):
    _,bp=source
    with patch("app.source_visuals._discover_local_primary_pdf") as discover:
        assert load_registry(bp,source_root=tmp_path / "expired-job") is None
        discover.assert_not_called()


def test_source_extraction_keeps_late_list_members(source):
    profile,bp=source
    formal=next(x for x in profile.coverage_items if x.source_anchor=="[P1] PAGE 11")
    assert "Refinement-based" in formal.why_important
    assert "unfamiliar" in formal.why_important
    full=" ".join(x for u in bp.units[5:15] for x in u.core_content)
    for term in ["Resilience","Standardized","Refinement-based","unfamiliar","Documentable"]:
        assert term in full
    assert any("PAGE 12" in u.source_anchor for u in bp.units[5:15])
    assert "PAGE 13" in bp.units[14].source_anchor


def test_model_does_not_silently_drop_the_ninth_source_fact(source):
    _,bp=source
    payload=bp.units[5].model_dump()
    payload["core_content"]=[f"Source fact number {i}." for i in range(12)]
    rebuilt=type(bp.units[5]).model_validate(payload)
    assert len(rebuilt.core_content)==12


def test_twenty_records_are_not_twenty_teaching_jobs(source):
    profile,bp=source
    broken=bp.model_copy(deep=True)
    for u in broken.units:
        u.title="A lecture slide"
        u.engineering_question="What is this topic?"
        u.pedagogy_content=[]
        u.student_action="Read the slide."
        u.takeaway="A broad topic."
        u.evidence=""
    broken.rubric_criteria[0].distinguished=""
    checks=deterministic_gate(broken,profile)
    assert not checks["v15_complete_20_unit_grammar"]
    assert not checks["v15_unit05_predict_constraint_derive_name"]
    assert not checks["v15_unit19_job_is_visible"]
    assert not checks["v15_unit20_job_is_visible"]


def test_no_silent_text_suffix_loss(source):
    _,bp=source
    sentence="Intolerable risks must be eliminated regardless of the cost required to do so."
    assert " ".join(wrap(sentence,130,18))==sentence
    u=bp.units[5].model_copy(update={"core_content":[sentence],"pedagogy_content":[]})
    items=teaching_items(bp,u)
    assert any(sentence==body for _,body in items)
    blocks,_,fits=text_layout(items)
    assert fits and sentence in " ".join(" ".join(b.lines) for b in blocks)


def test_titles_and_rubric_stay_in_their_bands(source):
    _,bp=source
    title=title_block("Dependability and Security Specification for Critical Sociotechnical Systems")
    assert title.y+len(title.lines)*title.size*1.22<=111
    blocks,_,_=rubric_layout(bp)
    assert max(b.y+len(b.lines)*b.size*1.22 for b in blocks)<460


def test_pdf_uses_exact_source_pages_and_embeds_fonts(source,tmp_path):
    _,bp=source
    path=export_presenter_pdf(bp,tmp_path/"lecture.pdf",source_root=LECTURE)
    with pymupdf.open(path) as doc:
        assert len(doc)==20
        assert "ENGINEERING CRISIS" in doc[0].get_text()
        rubric=doc[18].get_text()
        for level in ["DISTINGUISHED","READY","DEVELOPING","NOT YET READY"]:
            assert level in rubric
        # Core pages from 4 through 13 survive, including the image-only AI case.
        assert all(doc[n].get_images() for n in range(5,15))
        fonts=doc[0].get_fonts()
        assert any("DejaVuSans" in f[3] or "Vera" in f[3] for f in fonts)
        assert "AI MUST NOT BE TRUSTED AUTONOMOUSLY" in " ".join(doc[14].get_text().split())


def test_editable_pptx_retains_full_core_in_source_notes(source,tmp_path):
    _,bp=source
    path=export_presenter_pptx(bp,tmp_path/"lecture.pptx",source_root=LECTURE)
    deck=Presentation(path)
    assert len(deck.slides)==20
    assert "Refinement-based" in " ".join(slide.notes_slide.notes_text_frame.text for slide in deck.slides)
    rubric=" ".join(shape.text for shape in deck.slides[18].shapes if hasattr(shape,"text"))
    assert all(level in rubric for level in ["DISTINGUISHED","DEVELOPING","NOT YET READY"])


def test_production_entry_serves_v44_and_scoped_exports():
    # Version modules intentionally share one FastAPI app. Test the production
    # bootstrap in its own interpreter so historical route tests remain valid.
    code="""from fastapi.testclient import TestClient
from app.start_v440 import app
c=TestClient(app)
h=c.get('/api/health').json()
assert h['version']=='4.5.5'
assert h['generation_batch_size']==4
assert h['targeted_unit_repair'] is True
assert 'v15' in h['deterministic_gate']
p=c.get('/')
assert p.status_code==200 and 'studio_v440.js' in p.text
assert 'unitGrid' in p.text and 'coverageBody' in p.text
assert c.get('/static/studio_v440.js').status_code==200
assert c.get('/api/jobs/not-a-job/presenter').status_code==404
"""
    subprocess.run([sys.executable,"-c",code],cwd=ROOT,check=True,capture_output=True,text=True,timeout=60)


def test_ui_never_uses_vacuous_success_for_missing_checks():
    js=(ROOT/"app/static/studio_v440.js").read_text()
    assert "entries.length>0&&failed.length===0" in js
    assert "NOT CHECKED" in js
    assert "if(error.status===404||failures>=4)" in js
    assert "card.classList.toggle('selected'" in js
    assert "slideshare.net" not in js
