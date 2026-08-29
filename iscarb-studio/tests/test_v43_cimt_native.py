from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from fastapi.testclient import TestClient
from pptx import Presentation
from pypdf import PdfReader

from app.cimt_native_v43 import (
    export_cimt_presenter_pdf_v43,
    export_cimt_presenter_pptx_v43,
    presenter_text,
    render_cimt_presenter_preview_v43,
)
from app.gate_v14 import deterministic_gate
from app.start_v430 import app
from tests.v43_fixture import make_blueprint


class TestCimtNativeV43(unittest.TestCase):
    def setUp(self):
        self.bp = make_blueprint()

    def test_presenter_text_uses_no_ellipsis(self):
        text = 'This is a complete sentence. This second sentence is intentionally very long and should never be hard-truncated with visible ellipsis artifacts in the classroom presenter.'
        out = presenter_text(text, 55)
        self.assertEqual(out, 'This is a complete sentence.')
        self.assertNotIn('...', out)
        self.assertNotIn('…', out)

    def test_gate_v14_presenter_sentinels(self):
        checks = deterministic_gate(self.bp, None, '')
        self.assertTrue(checks['v14_reserved_iscarb_scaffolds_not_mislabeled_as_p1_core'])
        self.assertTrue(checks['v14_no_legacy_security_template_residue'])
        self.assertTrue(checks['v14_exactly_20_presenter_jobs'])
        self.assertTrue(checks['v14_exactly_90_live_minutes'])
        self.assertTrue(checks['v14_no_authored_hard_truncation_tokens'])

    def test_preview_is_cimt_native_and_20_units(self):
        html = render_cimt_presenter_preview_v43(self.bp, 'READY')
        up = html.upper()
        self.assertEqual(html.count('class="slide'), 20)
        self.assertIn('CIMT-NATIVE PRESENTER', up)
        self.assertIn('DEPENDABLE SYSTEMS', up)
        self.assertNotIn('SECURITY ENGINEERING', up)
        self.assertNotIn('PLATFORM PROTECTION', up)
        self.assertNotIn('...', html)
        self.assertNotIn('…', html)
        self.assertIn('#005634', html)
        self.assertIn('#C49A27', html)

    def test_pptx_20_slides_and_large_cimt_title(self):
        with tempfile.TemporaryDirectory() as td:
            path = export_cimt_presenter_pptx_v43(self.bp, Path(td) / 'presenter.pptx')
            prs = Presentation(str(path))
            self.assertEqual(len(prs.slides), 20)
            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not getattr(shape, 'has_text_frame', False):
                        continue
                    txt = shape.text.strip()
                    if txt:
                        all_text.append(txt)
                # The first content title box uses the CIMT large serif hierarchy.
                title_shapes = [s for s in slide.shapes if getattr(s, 'has_text_frame', False) and s.text.strip() == slide.shapes[1].text.strip()]
                # Robust font check: at least one Georgia run >= 28pt on every slide.
                large_georgia = False
                for s in slide.shapes:
                    if not getattr(s, 'has_text_frame', False):
                        continue
                    for p in s.text_frame.paragraphs:
                        for r in p.runs:
                            if (r.font.name or '').lower() == 'georgia' and r.font.size and r.font.size.pt >= 28:
                                large_georgia = True
                self.assertTrue(large_georgia)
            joined = ' '.join(all_text).upper()
            self.assertNotIn('SECURITY ENGINEERING', joined)
            self.assertNotIn('PLATFORM PROTECTION', joined)
            self.assertNotIn('...', joined)
            self.assertNotIn('…', joined)
            self.assertGreater(path.stat().st_size, 10000)

    def test_pdf_has_20_pages_and_no_security_residue(self):
        with tempfile.TemporaryDirectory() as td:
            path = export_cimt_presenter_pdf_v43(self.bp, Path(td) / 'presenter.pdf')
            reader = PdfReader(str(path))
            self.assertEqual(len(reader.pages), 20)
            txt = ' '.join((p.extract_text() or '') for p in reader.pages).upper()
            self.assertIn('DEPENDABLE', txt)
            self.assertNotIn('SECURITY ENGINEERING', txt)
            self.assertNotIn('PLATFORM PROTECTION', txt)
            self.assertNotIn('...', txt)
            self.assertGreater(path.stat().st_size, 10000)

    def test_local_output_lab_and_all_exports(self):
        client = TestClient(app)
        raw = self.bp.model_dump_json(by_alias=True).encode('utf-8')
        r = client.post('/api/render-blueprint', files={'blueprint_file': ('bp.json', raw, 'application/json')})
        self.assertEqual(r.status_code, 200, r.text)
        job_id = r.json()['job_id']
        p = client.get(f'/api/jobs/{job_id}/presenter')
        self.assertEqual(p.status_code, 200)
        self.assertIn('CIMT-native Presenter', p.text)
        for fmt in ('pptx','presenter-pdf','pdf','docx','student','json'):
            out = client.get(f'/api/jobs/{job_id}/export/{fmt}')
            self.assertEqual(out.status_code, 200, f'{fmt}: {out.text[:300]}')
            self.assertGreater(len(out.content), 500, fmt)

    def test_health_advertises_v43_contract(self):
        client = TestClient(app)
        r = client.get('/api/health')
        self.assertEqual(r.status_code, 200)
        data = r.json()
        self.assertEqual(data['version'], '4.3.0')
        self.assertEqual(data['deterministic_gate'], 'v14-provenance-presenter-on-v13')
        self.assertEqual(data['presenter_exact_units'], 20)
        self.assertEqual(data['session_minutes'], 90)


if __name__ == '__main__':
    unittest.main()
