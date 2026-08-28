from __future__ import annotations

from pathlib import Path


def extract_source_text(path: Path, limit: int = 600_000) -> str:
    """Best-effort local text extraction for deterministic gates.

    This is NOT used to replace Gemini's reading of the source. It gives hard gates
    a local corpus for conservative checks such as ETEC SLO atomicity and obvious
    unsupported-term detection.
    """
    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md"}:
            return path.read_text(encoding="utf-8", errors="replace")[:limit]
        if suffix == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            chunks = [(p.extract_text() or "") for p in reader.pages]
            return "\n".join(chunks)[:limit]
        if suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            chunks: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
            return "\n".join(chunks)[:limit]
        if suffix == ".docx":
            from docx import Document
            doc = Document(str(path))
            chunks = [p.text for p in doc.paragraphs if p.text]
            for table in doc.tables:
                for row in table.rows:
                    chunks.append(" | ".join(cell.text for cell in row.cells))
            return "\n".join(chunks)[:limit]
    except Exception:
        return ""
    return ""
