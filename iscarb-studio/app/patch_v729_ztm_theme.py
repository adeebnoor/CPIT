from __future__ import annotations

"""v7.2.9 — ISCARB ZTM-inspired high-contrast visual system.

This patch changes the visual projection, not the Golden v6.6 learning grammar.
It keeps the 20 core jobs, semantic P1 expansions, time-boxing, Rule 11 scaffold,
peer-review quick card, and source provenance intact while replacing the dark
presenter surface with a clean floating-card system.

Progressive disclosure is real in the HTML presenter (three stages per slide).
PPTX/PDF remain portable static final-state exports because python-pptx/reportlab
cannot author reliable cross-viewer PowerPoint animation timelines without
vendor-specific OOXML.
"""

import html
import json
import re
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib.colors import HexColor
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from . import main as engine
from . import start_v440 as base
from . import start_v670_prod as prod
from . import v670_contract as contract
from . import presenter_v67_prod as presenter
from . import patch_v725_golden_v660 as golden

# ---------------------------------------------------------------------------
# ZTM design tokens — user-approved specification
# ---------------------------------------------------------------------------
BG_BASE = "#FFFFFF"
BG_SURFACE = "#F8FAFC"
TEXT_HEADING = "#0F172A"
TEXT_BODY = "#475569"
ACCENT_PRIMARY = "#4F46E5"
ACCENT_CYAN = "#06B6D4"
ALERT_URGENT = "#F43F5E"
BORDER = "#E2E8F0"
PILL_BG = "#EEF2FF"
ROSE_BG = "#FFF1F2"
CYAN_BG = "#ECFEFF"
SHADOW = "#E2E8F0"
ACCENT_GOLD = "#D97706"
ACCENT_GREEN = "#059669"
ACCENT_BLUE = "#2563EB"

_PATCHED = False
_PREVIOUS_DRAFT = None


def _clean(v) -> str:
    return re.sub(r"\s+", " ", str(v or "")).strip()


def _timebox_parts(text: str, default: str = "2 min") -> tuple[str, str]:
    s = _clean(text)
    m = re.match(r"^TIMEBOX:\s*(.*?)\s+-\s+(.*)$", s, flags=re.I)
    if m:
        return m.group(1).strip(), m.group(2).strip()
    return default, s


def _label_value(rows, startswith: str) -> str:
    key = startswith.lower()
    for raw in rows or []:
        s = _clean(raw)
        if s.lower().startswith(key):
            parts = re.split(r"\s*(?:—|:|-)\s*", s, maxsplit=1)
            return parts[1].strip() if len(parts) == 2 else s
    return ""


def _expansion_id(spec: dict, fallback: int) -> int:
    m = re.search(r"(\d+)", str(spec.get("expansion_id", "")))
    return int(m.group(1)) if m else fallback


class ZTMTokens:
    """Compatibility object for the existing /api/design-tokens endpoint."""
    bg = BG_BASE
    panel = BG_SURFACE
    panel_soft = BG_SURFACE
    text = TEXT_HEADING
    muted = TEXT_BODY
    cyan = ACCENT_CYAN
    magenta = ACCENT_PRIMARY
    gold = ACCENT_GOLD
    green = ACCENT_GREEN
    blue = ACCENT_BLUE
    danger = ALERT_URGENT
    footer_bg = BG_SURFACE
    primary = ACCENT_PRIMARY
    secondary = ACCENT_CYAN
    heritage = ACCENT_GOLD

    def model_dump(self):
        return {
            "bg": self.bg, "panel": self.panel, "panel_soft": self.panel_soft,
            "text": self.text, "muted": self.muted, "cyan": self.cyan,
            "magenta": self.magenta, "gold": self.gold, "green": self.green,
            "blue": self.blue, "danger": self.danger, "footer_bg": self.footer_bg,
            "primary": self.primary, "secondary": self.secondary, "heritage": self.heritage,
        }

    def css_variables(self):
        # New public tokens plus old compatibility names used by existing clients.
        out = {
            "--bg-base": BG_BASE,
            "--bg-surface": BG_SURFACE,
            "--text-heading": TEXT_HEADING,
            "--text-body": TEXT_BODY,
            "--accent-primary": ACCENT_PRIMARY,
            "--accent-cyan": ACCENT_CYAN,
            "--alert-urgent": ALERT_URGENT,
        }
        out.update({f"--iscarb-{k.replace('_', '-')}": v for k, v in self.model_dump().items()})
        return out

    @staticmethod
    def _lum(h):
        vals = []
        for i in (1, 3, 5):
            c = int(h[i:i+2], 16) / 255
            vals.append(c/12.92 if c <= .04045 else ((c+.055)/1.055)**2.4)
        return .2126*vals[0] + .7152*vals[1] + .0722*vals[2]

    @classmethod
    def _ratio(cls, a, b):
        x, y = cls._lum(a), cls._lum(b)
        return (max(x, y)+.05)/(min(x, y)+.05)

    def contrast_checks(self):
        # Cyan is decorative/iconic by design; body text never relies on cyan.
        return {
            "heading_on_base": self._ratio(TEXT_HEADING, BG_BASE) >= 4.5,
            "body_on_base": self._ratio(TEXT_BODY, BG_BASE) >= 4.5,
            "heading_on_surface": self._ratio(TEXT_HEADING, BG_SURFACE) >= 4.5,
            "body_on_surface": self._ratio(TEXT_BODY, BG_SURFACE) >= 4.5,
            "primary_large_on_base": self._ratio(ACCENT_PRIMARY, BG_BASE) >= 3.0,
            "urgent_large_on_base": self._ratio(ALERT_URGENT, BG_BASE) >= 3.0,
        }


def ztm_chapter_design_tokens(title="", preferred=""):
    return ZTMTokens()


# ---------------------------------------------------------------------------
# Physical plan: current semantic expansions, not old overflow-only presenter.
# ---------------------------------------------------------------------------
def _ztm_contract_plan(bp, target=30, strict=True):
    if strict and not contract.verdict_eligible(bp):
        raise ValueError("Bounded Verdict blocked until Rule 18 assurance and 6x4 rubric are complete")
    specs = list(contract.plan_expansions(bp, target=target) or [])[:8]
    by = {}
    for spec in specs:
        by.setdefault(int(spec.get("after_unit", 15)), []).append(spec)
    plan = [{"kind": "COVER"}]
    for u in bp.units:
        plan.append({"kind": "CORE", "unit_number": u.number})
        for spec in by.get(u.number, []):
            plan.append({"kind": "SOURCE_EXPANSION", **spec})
    plan.append({"kind": "CLOSE"})
    if len(plan) > 30:
        raise ValueError(f"ZTM/Golden physical overflow: {len(plan)} slides")
    for i, row in enumerate(plan, 1):
        row.update(physical_index=i, physical_total=len(plan))
    return plan


def _physical_plan(bp):
    rows = _ztm_contract_plan(bp, target=30, strict=False)
    units = {u.number: u for u in bp.units}
    out = []
    exp_counter = 0
    for row in rows:
        kind = row.get("kind")
        if kind == "COVER":
            out.append(("cover", None, None))
        elif kind == "CLOSE":
            out.append(("close", None, None))
        elif kind == "CORE":
            out.append(("unit", units[int(row["unit_number"])], None))
        else:
            exp_counter += 1
            after = int(row.get("after_unit", 15))
            out.append(("expansion", units[after], (_expansion_id(row, exp_counter), row)))
    return out


# ---------------------------------------------------------------------------
# PPTX primitives
# ---------------------------------------------------------------------------
def _ppt_text(slide, x, y, w, h, text, size=16, color=TEXT_BODY, bold=False,
              align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Inter"):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = Inches(.02)
    tf.margin_top = tf.margin_bottom = Inches(.01)
    p = tf.paragraphs[0]
    p.text = _clean(text)
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = presenter._rgb(color)
    p.alignment = align
    return sh


def _ppt_badge(slide, x, y, w, text, bg=PILL_BG, fg=ACCENT_PRIMARY, border=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(.34))
    sh.fill.solid(); sh.fill.fore_color.rgb = presenter._rgb(bg)
    sh.line.color.rgb = presenter._rgb(border or bg)
    tf = sh.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(.05); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.text = _clean(text).upper(); p.alignment = PP_ALIGN.CENTER
    p.font.name = "Inter"; p.font.size = Pt(7.5); p.font.bold = True; p.font.color.rgb = presenter._rgb(fg)
    return sh


def _ppt_box(slide, x, y, w, h, title, body="", accent=ACCENT_PRIMARY, fill=BG_SURFACE,
             body_size=10.7, title_size=10.5):
    # Subtle offset shadow; no transparency dependency.
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+.035), Inches(y+.045), Inches(w), Inches(h))
    shadow.fill.solid(); shadow.fill.fore_color.rgb = presenter._rgb(SHADOW); shadow.line.fill.background()
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = presenter._rgb(fill)
    sh.line.color.rgb = presenter._rgb(BORDER); sh.line.width = Pt(.8)
    # Slim accent rule on the left gives semantics without noisy borders.
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+.03), Inches(y+.12), Inches(.045), Inches(max(.22, h-.24)))
    bar.fill.solid(); bar.fill.fore_color.rgb = presenter._rgb(accent); bar.line.fill.background()
    tf = sh.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(.18); tf.margin_right = Inches(.14); tf.margin_top = Inches(.12); tf.margin_bottom = Inches(.09)
    if title:
        p = tf.paragraphs[0]; p.text = _clean(title).upper(); p.font.name = "Inter"; p.font.size = Pt(title_size)
        p.font.bold = True; p.font.color.rgb = presenter._rgb(accent)
    else:
        p = tf.paragraphs[0]; p.text = ""
    if body:
        q = tf.add_paragraph() if title else p
        q.text = _clean(body); q.font.name = "Inter"; q.font.size = Pt(body_size); q.font.color.rgb = presenter._rgb(TEXT_BODY)
        if title: q.space_before = Pt(4)
    return sh


def _ppt_header(slide, u, page_idx, total):
    is_crisis = u.number == 1
    _ppt_badge(slide, .38, .18, 1.02, f"RULE {u.number:02d}", ROSE_BG if is_crisis else PILL_BG,
               ALERT_URGENT if is_crisis else ACCENT_PRIMARY)
    phase = presenter.PHASE_LABEL.get(u.phase, u.phase)
    kind = presenter.RULE_KIND.get(u.number, "CONCEPT")
    _ppt_badge(slide, 1.52, .18, 2.45, f"{phase} · {kind}", CYAN_BG, ACCENT_CYAN)
    _ppt_text(slide, 11.02, .22, 1.92, .22, f"U{u.number:02d}/20 · {page_idx:02d}/{total:02d}", 6.9, TEXT_BODY, False, PP_ALIGN.RIGHT)
    _ppt_text(slide, .38, .67, 12.15, .47, presenter.RULE_NAMES.get(u.number, u.title), 21.5, TEXT_HEADING, True)
    _ppt_text(slide, .38, 1.16, 12.05, .35, presenter._short(u.engineering_question, 28), 9.7, TEXT_BODY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.38), Inches(1.54), Inches(12.55), Inches(.012))
    line.fill.solid(); line.fill.fore_color.rgb = presenter._rgb(BORDER); line.line.fill.background()


def _ppt_taskbar(slide, task: str, anchor: str = "", timebox: str = "2 min"):
    y, h = 6.62, .79
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.32), Inches(y+.035), Inches(12.68), Inches(h))
    shadow.fill.solid(); shadow.fill.fore_color.rgb = presenter._rgb(SHADOW); shadow.line.fill.background()
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.29), Inches(y), Inches(12.68), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = presenter._rgb(BG_SURFACE); card.line.color.rgb = presenter._rgb(BORDER); card.line.width = Pt(.8)
    # 4px-ish two-color top rule (gradient surrogate, portable in PPTX/PDF).
    for x, w, col in [(.34, 6.28, ACCENT_PRIMARY), (6.62, 6.28, ACCENT_CYAN)]:
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y+.025), Inches(w), Inches(.045))
        r.fill.solid(); r.fill.fore_color.rgb = presenter._rgb(col); r.line.fill.background()
    _ppt_text(slide, .50, y+.18, 1.05, .20, "YOUR TASK", 7.3, ACCENT_PRIMARY, True)
    _ppt_text(slide, 1.55, y+.15, 8.45, .36, presenter._short(task, 34), 8.5, TEXT_HEADING, False, valign=MSO_ANCHOR.MIDDLE)
    _ppt_badge(slide, 10.20, y+.16, 1.68, f"TIMEBOX · {timebox}", ROSE_BG, ALERT_URGENT, "#FECDD3")
    if anchor:
        _ppt_text(slide, 11.95, y+.20, .80, .18, presenter._short(anchor, 7), 5.6, TEXT_BODY, False, PP_ALIGN.RIGHT)


def _ppt_footer(slide, u):
    tb, task = _timebox_parts(u.student_action)
    _ppt_taskbar(slide, task, presenter._anchor(u), tb)


def _ppt_source_split(slide, u, asset_path, accent):
    # Left source card.
    _ppt_box(slide, .48, 1.82, 6.03, 4.48, "PRIMARY SOURCE", "", ACCENT_CYAN, BG_SURFACE, 8, 8)
    presenter._ppt_add_image(slide, asset_path, .63, 2.12, 5.72, 3.85)
    _ppt_badge(slide, .75, 5.78, 1.62, "P1 SOURCE", PILL_BG, ACCENT_PRIMARY)
    items = presenter._items(u, 4)
    _ppt_text(slide, 6.82, 1.84, 5.55, .24, "FROM THE PRIMARY SOURCE", 8.2, ACCENT_PRIMARY, True)
    y = 2.18
    for item in items[:4]:
        _ppt_box(slide, 6.82, y, 5.62, .78, "", presenter._short(item, 22), accent, BG_SURFACE, 9.3, 1)
        y += .91
    if u.takeaway:
        _ppt_box(slide, 6.82, 5.75, 5.62, .63, "ENGINEERING TAKEAWAY", presenter._short(u.takeaway, 18), ACCENT_PRIMARY, PILL_BG, 7.8, 7.4)


def _ppt_semantic(slide, bp, u, accent):
    # Operational Rule 11: scaffold before local transfer.
    if u.number == 11:
        micro = _label_value(u.pedagogy_content, "MICRO-CASE") or "Solve one tiny mechanism-first case before adding local complexity."
        transfer = _label_value(u.pedagogy_content, "TRANSFER RULE") or "Reuse the same mechanism → evidence → decision-boundary chain on the Saudi/local case."
        _ppt_badge(slide, .72, 1.86, 1.44, "STEP 1", PILL_BG, ACCENT_PRIMARY)
        _ppt_box(slide, .72, 2.24, 5.72, 2.75, "MICRO-CASE", presenter._short(micro, 48), ACCENT_PRIMARY, BG_SURFACE, 11.0, 10.0)
        _ppt_badge(slide, 6.90, 1.86, 1.44, "STEP 2", CYAN_BG, ACCENT_CYAN)
        _ppt_box(slide, 6.90, 2.24, 5.72, 2.75, "SAUDI / LOCAL TRANSFER", presenter._short(transfer, 48), ACCENT_CYAN, BG_SURFACE, 11.0, 10.0)
        _ppt_box(slide, 2.18, 5.33, 8.88, .72, "THINKING CHAIN", "Mechanism → Evidence → Decision Boundary", ACCENT_PRIMARY, PILL_BG, 8.8, 8.0)
        return
    # Operational Rule 19: only two peer-review questions on the learner slide.
    if u.number == 19:
        q1 = u.core_content[0] if len(u.core_content) > 0 else "Is the evidence independently inspectable?"
        q2 = u.core_content[1] if len(u.core_content) > 1 else "What would falsify the claim?"
        _ppt_badge(slide, .72, 1.87, 2.10, "PEER REVIEW · 2 QUESTIONS", PILL_BG, ACCENT_PRIMARY)
        _ppt_box(slide, .72, 2.36, 5.72, 2.80, "Q1 · INSPECTABILITY", presenter._short(q1, 42), ACCENT_PRIMARY, BG_SURFACE, 12.0, 10.2)
        _ppt_box(slide, 6.90, 2.36, 5.72, 2.80, "Q2 · FALSIFIER", presenter._short(q2, 42), ACCENT_CYAN, BG_SURFACE, 12.0, 10.2)
        _ppt_text(slide, 2.0, 5.52, 9.3, .34, "The full 6×4 rubric remains in the instructor / blueprint layer.", 9.3, TEXT_BODY, False, PP_ALIGN.CENTER)
        return
    return _ORIGINAL_PPT_SEMANTIC(slide, bp, u, accent)


def _decision_box_text(spec: dict) -> str:
    for row in spec.get("content", []) or []:
        if str(row).upper().startswith("DECISION EVIDENCE BOX"):
            return re.sub(r"^DECISION EVIDENCE BOX\s*[-—:]\s*", "", str(row), flags=re.I).strip()
    return str(spec.get("visual_evidence_role", "Convert theory into a decision and inspectable evidence."))


def _source_rows(spec: dict):
    return [str(x) for x in spec.get("content", []) or [] if not str(x).upper().startswith("DECISION EVIDENCE BOX")]


def _ppt_expansion(slide, u, idx, spec, page_idx, total):
    presenter._ppt_bg(slide, BG_BASE)
    title = str(spec.get("title") or f"{u.title} — source detail")
    source_anchor = str(spec.get("source_anchor") or presenter._anchor(u))
    _ppt_badge(slide, .38, .18, 1.82, f"SOURCE EXPANSION · X{idx:02d}", PILL_BG, ACCENT_PRIMARY)
    _ppt_text(slide, 11.0, .22, 1.92, .22, f"{page_idx:02d}/{total:02d}", 6.9, TEXT_BODY, False, PP_ALIGN.RIGHT)
    _ppt_text(slide, .38, .70, 8.35, .55, title, 20.5, TEXT_HEADING, True)
    _ppt_text(slide, .38, 1.27, 8.25, .32, "Primary-source detail preserved at readable density.", 9.2, TEXT_BODY)
    _ppt_box(slide, 9.00, .68, 3.56, 1.25, "DECISION → EVIDENCE", presenter._short(_decision_box_text(spec), 30), ACCENT_PRIMARY, PILL_BG, 8.2, 7.9)
    rows = _source_rows(spec)[:6]
    cols = 2; w = 5.85
    for i, item in enumerate(rows):
        r, c = divmod(i, cols)
        _ppt_box(slide, .55+c*6.15, 2.10+r*1.34, w, 1.08, f"P1 · {i+1}", presenter._short(item, 24), ACCENT_CYAN if c else ACCENT_PRIMARY, BG_SURFACE, 9.0, 7.3)
    task = str(spec.get("student_task") or "Use one source detail to strengthen or challenge the current decision.")
    _ppt_taskbar(slide, task, source_anchor, "2 min")


def _ppt_cover(slide, bp, total):
    presenter._ppt_bg(slide, BG_BASE)
    # Minimal ZTM cover: whitespace, energetic accent geometry, no decorative stock image.
    _ppt_badge(slide, .72, .58, 1.58, "ISCARB", PILL_BG, ACCENT_PRIMARY)
    _ppt_badge(slide, 2.42, .58, 2.22, "ENGINEERING LECTURE", CYAN_BG, ACCENT_CYAN)
    _ppt_text(slide, .72, 1.42, 9.85, 1.55, bp.lecture_title, 30, TEXT_HEADING, True)
    _ppt_text(slide, .75, 3.15, 8.35, .78, "One engineering decision. Source-backed mechanisms. Evidence before verdict.", 13.8, TEXT_BODY)
    _ppt_box(slide, .75, 4.36, 5.45, 1.08, "PRIMARY SOURCE PRESERVED", f"20 core units · {max(0,total-22)} semantic source expansions · ≤30 physical slides", ACCENT_PRIMARY, BG_SURFACE, 9.0, 8.2)
    _ppt_box(slide, 6.55, 4.36, 5.55, 1.08, "DECISION JOURNEY", "CRISIS → MAP → MECHANISM → TRADE-OFF → EVIDENCE → VERDICT", ACCENT_CYAN, BG_SURFACE, 8.7, 8.2)
    # Accent ribbons.
    for x, y, w, h, col in [(10.55, .35, 2.2, .13, ACCENT_PRIMARY), (11.05, .62, 1.7, .13, ACCENT_CYAN), (11.55, .89, 1.2, .13, ACCENT_GOLD)]:
        r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        r.fill.solid(); r.fill.fore_color.rgb = presenter._rgb(col); r.line.fill.background()
    _ppt_text(slide, .75, 6.46, 8.9, .28, "High contrast · floating cards · visible tasks · progressive presenter layers", 8.2, TEXT_BODY)


def _ppt_close(slide, bp, total):
    presenter._ppt_bg(slide, BG_BASE)
    _ppt_badge(slide, .72, .58, 1.72, "ISCARB · CLOSE", PILL_BG, ACCENT_PRIMARY)
    _ppt_text(slide, .72, 1.28, 9.6, .62, "Bounded engineering verdict", 27, TEXT_HEADING, True)
    _ppt_text(slide, .72, 2.02, 9.8, .45, "The lecture closes only after the assurance chain is inspectable.", 12.2, TEXT_BODY)
    labs = ["CLAIM", "EVIDENCE", "WARRANT", "COUNTER-EVIDENCE", "RESIDUAL UNCERTAINTY", "VERDICT"]
    cols = [ACCENT_PRIMARY, ACCENT_CYAN, ACCENT_BLUE, ACCENT_GOLD, ACCENT_GREEN, ACCENT_PRIMARY]
    for i, lab in enumerate(labs):
        _ppt_box(slide, .45+i*2.08, 3.03, 1.82, 1.42, lab, "", cols[i], BG_SURFACE, 8, 7.5)
    _ppt_taskbar(slide, "State the final verdict and the one piece of evidence that would make you revisit it.", "ASSURANCE", "3 min")


# ---------------------------------------------------------------------------
# PDF primitives
# ---------------------------------------------------------------------------
def _pdf_box(c, x, y, w, h, title, body="", accent=ACCENT_PRIMARY, fill=BG_SURFACE,
             body_size=11, title_size=9):
    c.setFillColor(HexColor(SHADOW)); c.setStrokeColor(HexColor(SHADOW)); c.roundRect(x+2.5, y-2.5, w, h, 9, fill=1, stroke=0)
    c.setFillColor(HexColor(fill)); c.setStrokeColor(HexColor(BORDER)); c.setLineWidth(.7); c.roundRect(x, y, w, h, 9, fill=1, stroke=1)
    c.setFillColor(HexColor(accent)); c.roundRect(x+4, y+8, 3.5, max(10, h-16), 1.5, fill=1, stroke=0)
    if title:
        presenter._pdf_text(c, x+13, y+h-27, w-24, 18, title.upper(), title_size, accent, True, max_lines=1)
    if body:
        presenter._pdf_text(c, x+13, y+10, w-24, h-(38 if title else 18), body, body_size, TEXT_BODY, False,
                            max_lines=max(1, int((h-(42 if title else 20))/(body_size*1.18))))


def _pdf_badge(c, x, y, w, text, bg=PILL_BG, fg=ACCENT_PRIMARY):
    c.setFillColor(HexColor(bg)); c.setStrokeColor(HexColor(bg)); c.roundRect(x, y, w, 24, 12, fill=1, stroke=0)
    presenter._pdf_text(c, x+5, y+5, w-10, 14, text.upper(), 6.3, fg, True, "center", 1)


def _pdf_header(c, u, page_idx, total):
    is_crisis = u.number == 1
    _pdf_badge(c, 26, 501, 76, f"RULE {u.number:02d}", ROSE_BG if is_crisis else PILL_BG, ALERT_URGENT if is_crisis else ACCENT_PRIMARY)
    _pdf_badge(c, 112, 501, 162, f"{presenter.PHASE_LABEL.get(u.phase,u.phase)} · {presenter.RULE_KIND.get(u.number,'CONCEPT')}", CYAN_BG, ACCENT_CYAN)
    presenter._pdf_text(c, 792, 505, 142, 14, f"U{u.number:02d}/20 · {page_idx:02d}/{total:02d}", 6.2, TEXT_BODY, False, "right", 1)
    presenter._pdf_text(c, 28, 455, 900, 38, presenter.RULE_NAMES.get(u.number, u.title), 19.5, TEXT_HEADING, True, max_lines=1)
    presenter._pdf_text(c, 28, 427, 900, 23, presenter._short(u.engineering_question, 28), 8.2, TEXT_BODY, max_lines=2)
    c.setStrokeColor(HexColor(BORDER)); c.setLineWidth(.7); c.line(28, 420, 932, 420)


def _pdf_taskbar(c, task: str, anchor: str = "", timebox: str = "2 min"):
    c.setFillColor(HexColor(SHADOW)); c.roundRect(25, 7, 910, 54, 9, fill=1, stroke=0)
    c.setFillColor(HexColor(BG_SURFACE)); c.setStrokeColor(HexColor(BORDER)); c.setLineWidth(.7); c.roundRect(23, 9, 910, 54, 9, fill=1, stroke=1)
    c.setFillColor(HexColor(ACCENT_PRIMARY)); c.rect(28, 58, 450, 3, fill=1, stroke=0)
    c.setFillColor(HexColor(ACCENT_CYAN)); c.rect(478, 58, 450, 3, fill=1, stroke=0)
    presenter._pdf_text(c, 38, 31, 75, 15, "YOUR TASK", 6.4, ACCENT_PRIMARY, True, max_lines=1)
    presenter._pdf_text(c, 118, 23, 570, 28, presenter._short(task, 36), 7.3, TEXT_HEADING, False, max_lines=2)
    _pdf_badge(c, 706, 25, 142, f"TIMEBOX · {timebox}", ROSE_BG, ALERT_URGENT)
    if anchor:
        presenter._pdf_text(c, 852, 28, 70, 14, presenter._short(anchor, 7), 4.9, TEXT_BODY, False, "right", 1)


def _pdf_footer(c, u):
    tb, task = _timebox_parts(u.student_action)
    _pdf_taskbar(c, task, presenter._anchor(u), tb)


def _pdf_semantic(c, bp, u, accent):
    if u.number == 11:
        micro = _label_value(u.pedagogy_content, "MICRO-CASE") or "Solve one tiny mechanism-first case before local complexity."
        transfer = _label_value(u.pedagogy_content, "TRANSFER RULE") or "Reuse mechanism → evidence → decision boundary on the Saudi/local case."
        _pdf_badge(c, 54, 385, 92, "STEP 1", PILL_BG, ACCENT_PRIMARY)
        _pdf_box(c, 54, 185, 405, 185, "MICRO-CASE", presenter._short(micro, 48), ACCENT_PRIMARY, BG_SURFACE, 9.5, 8.0)
        _pdf_badge(c, 505, 385, 92, "STEP 2", CYAN_BG, ACCENT_CYAN)
        _pdf_box(c, 505, 185, 405, 185, "SAUDI / LOCAL TRANSFER", presenter._short(transfer, 48), ACCENT_CYAN, BG_SURFACE, 9.5, 8.0)
        _pdf_box(c, 180, 98, 600, 52, "THINKING CHAIN", "Mechanism → Evidence → Decision Boundary", ACCENT_PRIMARY, PILL_BG, 7.5, 6.8)
        return
    if u.number == 19:
        q1 = u.core_content[0] if len(u.core_content) > 0 else "Is the evidence independently inspectable?"
        q2 = u.core_content[1] if len(u.core_content) > 1 else "What would falsify the claim?"
        _pdf_badge(c, 52, 385, 155, "PEER REVIEW · 2 QUESTIONS", PILL_BG, ACCENT_PRIMARY)
        _pdf_box(c, 52, 180, 405, 185, "Q1 · INSPECTABILITY", presenter._short(q1, 42), ACCENT_PRIMARY, BG_SURFACE, 10.2, 8.2)
        _pdf_box(c, 503, 180, 405, 185, "Q2 · FALSIFIER", presenter._short(q2, 42), ACCENT_CYAN, BG_SURFACE, 10.2, 8.2)
        presenter._pdf_text(c, 180, 120, 600, 25, "Full 6×4 rubric remains in the instructor / blueprint layer.", 7.8, TEXT_BODY, False, "center", 2)
        return
    return _ORIGINAL_PDF_SEMANTIC(c, bp, u, accent)


def _pdf_expansion(c, u, idx, spec, page_idx, total):
    c.setFillColor(HexColor(BG_BASE)); c.rect(0, 0, 960, 540, fill=1, stroke=0)
    title = str(spec.get("title") or f"{u.title} — source detail")
    anchor = str(spec.get("source_anchor") or presenter._anchor(u))
    _pdf_badge(c, 28, 501, 142, f"SOURCE EXPANSION · X{idx:02d}", PILL_BG, ACCENT_PRIMARY)
    presenter._pdf_text(c, 790, 505, 142, 14, f"{page_idx:02d}/{total:02d}", 6.2, TEXT_BODY, False, "right", 1)
    presenter._pdf_text(c, 28, 458, 610, 34, title, 18.5, TEXT_HEADING, True, max_lines=1)
    presenter._pdf_text(c, 28, 432, 600, 20, "Primary-source detail preserved at readable density.", 7.8, TEXT_BODY, max_lines=1)
    _pdf_box(c, 665, 423, 265, 72, "DECISION → EVIDENCE", presenter._short(_decision_box_text(spec), 30), ACCENT_PRIMARY, PILL_BG, 6.5, 6.2)
    rows = _source_rows(spec)[:6]
    for i, item in enumerate(rows):
        r, col = divmod(i, 2)
        _pdf_box(c, 45+col*455, 315-r*82, 410, 65, f"P1 · {i+1}", presenter._short(item, 24), ACCENT_PRIMARY if col == 0 else ACCENT_CYAN, BG_SURFACE, 7.5, 6.2)
    _pdf_taskbar(c, str(spec.get("student_task") or "Use one source detail to strengthen or challenge the current decision."), anchor, "2 min")


def _pdf_cover(c, bp, total):
    c.setFillColor(HexColor(BG_BASE)); c.rect(0, 0, 960, 540, fill=1, stroke=0)
    _pdf_badge(c, 50, 485, 92, "ISCARB", PILL_BG, ACCENT_PRIMARY)
    _pdf_badge(c, 152, 485, 150, "ENGINEERING LECTURE", CYAN_BG, ACCENT_CYAN)
    presenter._pdf_text(c, 50, 350, 720, 110, bp.lecture_title, 27, TEXT_HEADING, True, max_lines=3)
    presenter._pdf_text(c, 52, 300, 650, 42, "One engineering decision. Source-backed mechanisms. Evidence before verdict.", 11.5, TEXT_BODY, max_lines=2)
    _pdf_box(c, 52, 190, 390, 75, "PRIMARY SOURCE PRESERVED", f"20 core units · {max(0,total-22)} semantic source expansions · ≤30 physical slides", ACCENT_PRIMARY, BG_SURFACE, 7.5, 6.8)
    _pdf_box(c, 470, 190, 420, 75, "DECISION JOURNEY", "CRISIS → MAP → MECHANISM → TRADE-OFF → EVIDENCE → VERDICT", ACCENT_CYAN, BG_SURFACE, 7.2, 6.8)
    c.setFillColor(HexColor(ACCENT_PRIMARY)); c.roundRect(760, 465, 150, 8, 4, fill=1, stroke=0)
    c.setFillColor(HexColor(ACCENT_CYAN)); c.roundRect(795, 445, 115, 8, 4, fill=1, stroke=0)
    c.setFillColor(HexColor(ACCENT_GOLD)); c.roundRect(830, 425, 80, 8, 4, fill=1, stroke=0)
    presenter._pdf_text(c, 52, 118, 700, 22, "High contrast · floating cards · visible tasks · progressive presenter layers", 7.2, TEXT_BODY, max_lines=1)


def _pdf_close(c, bp, total):
    c.setFillColor(HexColor(BG_BASE)); c.rect(0, 0, 960, 540, fill=1, stroke=0)
    _pdf_badge(c, 50, 485, 110, "ISCARB · CLOSE", PILL_BG, ACCENT_PRIMARY)
    presenter._pdf_text(c, 50, 420, 650, 45, "Bounded engineering verdict", 24, TEXT_HEADING, True, max_lines=1)
    presenter._pdf_text(c, 50, 375, 700, 30, "The lecture closes only after the assurance chain is inspectable.", 10.3, TEXT_BODY, max_lines=2)
    labs = ["CLAIM", "EVIDENCE", "WARRANT", "COUNTER-EVIDENCE", "RESIDUAL UNCERTAINTY", "VERDICT"]
    cols = [ACCENT_PRIMARY, ACCENT_CYAN, ACCENT_BLUE, ACCENT_GOLD, ACCENT_GREEN, ACCENT_PRIMARY]
    for i, lab in enumerate(labs):
        _pdf_box(c, 24+i*155, 205, 140, 90, lab, "", cols[i], BG_SURFACE, 7, 6.0)
    _pdf_taskbar(c, "State the final verdict and the one piece of evidence that would make you revisit it.", "ASSURANCE", "3 min")


# ---------------------------------------------------------------------------
# HTML progressive disclosure presenter
# ---------------------------------------------------------------------------
def _html_source_body(u, source):
    if source and u.number in {6, 7, 9, 10, 12, 13, 14}:
        uri = presenter._html_img(source[1])
        cards = "".join(f"<div class='card compact'>{html.escape(presenter._short(x,22))}</div>" for x in presenter._items(u,4))
        return f"<div class='source-card'><img class='source-img' src='{uri}' alt='Primary source visual'/><span class='pill'>P1 SOURCE</span></div><div class='stack'>{cards}</div>"
    return "".join(f"<div class='card'><span class='mini'>{i+1:02d}</span>{html.escape(presenter._short(x,24))}</div>" for i, x in enumerate(presenter._items(u,6) or [u.takeaway]))


def render_presenter_preview_ztm(bp, release_state="REVIEW", source_root=None) -> str:
    plan = _physical_plan(bp); assets = presenter._source_assets(bp, source_root); total = len(plan); pages = []
    hourglass = """<svg viewBox='0 0 24 24' aria-hidden='true'><path d='M7 3h10M7 21h10M8 4c0 4 2 5 4 7-2 2-4 3-4 7m8-14c0 4-2 5-4 7 2 2 4 3 4 7' fill='none' stroke='currentColor' stroke-width='1.8' stroke-linecap='round'/></svg>"""
    for page_idx, (kind, u, extra) in enumerate(plan, 1):
        if kind == "cover":
            pages.append(f"<section class='slide cover' data-stage='3'><div class='cover-accent'></div><div class='badges'><span class='pill'>ISCARB</span><span class='pill cyan'>ENGINEERING LECTURE</span></div><h1>{html.escape(bp.lecture_title)}</h1><p>One engineering decision. Source-backed mechanisms. Evidence before verdict.</p><div class='cover-grid'><div class='card'><b>PRIMARY SOURCE PRESERVED</b>20 core units · {max(0,total-22)} semantic source expansions · ≤30 physical slides</div><div class='card'><b>DECISION JOURNEY</b>CRISIS → MAP → MECHANISM → TRADE-OFF → EVIDENCE → VERDICT</div></div></section>")
            continue
        if kind == "close":
            pages.append("<section class='slide cover' data-stage='3'><div class='badges'><span class='pill'>ISCARB · CLOSE</span></div><h1>Bounded engineering verdict</h1><p>CLAIM · EVIDENCE · WARRANT · COUNTER-EVIDENCE · RESIDUAL UNCERTAINTY · VERDICT</p></section>")
            continue
        if kind == "expansion":
            idx, spec = extra
            title = str(spec.get("title") or f"{u.title} — source detail")
            rows = _source_rows(spec)[:6]
            cards = "".join(f"<div class='card'><span class='mini'>P1 · {i+1}</span>{html.escape(presenter._short(x,24))}</div>" for i, x in enumerate(rows))
            task = str(spec.get("student_task") or "Use one source detail to strengthen or challenge the current decision.")
            pages.append(f"<section class='slide' data-stage='1'><div class='stage stage1'><div class='top'><span class='pill'>SOURCE EXPANSION · X{idx:02d}</span><span class='counter'>{page_idx:02d}/{total:02d}</span></div><h2>{html.escape(title)}</h2><p class='question'>Primary-source detail preserved at readable density.</p><div class='decision-corner'><b>DECISION → EVIDENCE</b>{html.escape(presenter._short(_decision_box_text(spec),30))}</div></div><div class='stage stage2 body'>{cards}</div><div class='stage stage3 taskbar'><div><b>YOUR TASK</b><span>{html.escape(task)}</span></div><span class='timebox'>{hourglass}<b>TIMEBOX · 2 MIN</b></span></div><span class='stage-indicator'>1/3</span></section>")
            continue
        tb, task = _timebox_parts(u.student_action)
        crisis = " crisis" if u.number == 1 else ""
        if u.number == 11:
            micro = _label_value(u.pedagogy_content, "MICRO-CASE") or "Solve one tiny mechanism-first case before local complexity."
            transfer = _label_value(u.pedagogy_content, "TRANSFER RULE") or "Reuse mechanism → evidence → decision boundary on the Saudi/local case."
            body = f"<div class='card'><span class='mini'>STEP 1 · MICRO-CASE</span>{html.escape(presenter._short(micro,48))}</div><div class='card'><span class='mini cyan-text'>STEP 2 · SAUDI / LOCAL TRANSFER</span>{html.escape(presenter._short(transfer,48))}</div><div class='card chain'><b>Mechanism → Evidence → Decision Boundary</b></div>"
        elif u.number == 19:
            q1 = u.core_content[0] if len(u.core_content) > 0 else "Is the evidence independently inspectable?"
            q2 = u.core_content[1] if len(u.core_content) > 1 else "What would falsify the claim?"
            body = f"<div class='card'><span class='mini'>Q1 · INSPECTABILITY</span>{html.escape(q1)}</div><div class='card'><span class='mini cyan-text'>Q2 · FALSIFIER</span>{html.escape(q2)}</div><div class='rubric-note'>Full 6×4 rubric stays in the instructor / blueprint layer.</div>"
        else:
            body = _html_source_body(u, assets.get(u.number))
        phase = presenter.PHASE_LABEL.get(u.phase, u.phase)
        kind_label = presenter.RULE_KIND.get(u.number, "CONCEPT")
        pages.append(f"<section class='slide{crisis}' data-stage='1'><div class='stage stage1'><div class='top'><div class='badges'><span class='pill rule'>RULE {u.number:02d}</span><span class='pill cyan'>{html.escape(phase)} · {html.escape(kind_label)}</span></div><span class='counter'>U{u.number:02d}/20 · {page_idx:02d}/{total:02d}</span></div><h2>{html.escape(presenter.RULE_NAMES.get(u.number,u.title))}</h2><p class='question'>{html.escape(presenter._short(u.engineering_question,28))}</p></div><div class='stage stage2 body'>{body}</div><div class='stage stage3 taskbar'><div><b>YOUR TASK</b><span>{html.escape(task)}</span></div><span class='timebox'>{hourglass}<b>TIMEBOX · {html.escape(tb.upper())}</b></span></div><span class='stage-indicator'>1/3</span></section>")

    css = f"""
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;800&display=swap');
:root{{--bg-base:{BG_BASE};--bg-surface:{BG_SURFACE};--text-heading:{TEXT_HEADING};--text-body:{TEXT_BODY};--accent-primary:{ACCENT_PRIMARY};--accent-cyan:{ACCENT_CYAN};--alert-urgent:{ALERT_URGENT};--border:{BORDER};}}
*{{box-sizing:border-box}} body{{margin:0;background:#EEF2F7;color:var(--text-body);font-family:'Inter',Arial,sans-serif}} .deck{{padding:24px}}
.slide{{width:min(1280px,96vw);aspect-ratio:16/9;margin:0 auto 28px;background:var(--bg-base);border:1px solid var(--border);border-radius:18px;position:relative;overflow:hidden;box-shadow:0 18px 45px rgba(15,23,42,.10)}}
.stage1{{position:absolute;left:3.6%;right:3.6%;top:4%;height:24%}} .top{{display:flex;justify-content:space-between;align-items:center}} .badges{{display:flex;gap:10px;align-items:center}}
.pill{{display:inline-flex;align-items:center;background:#EEF2FF;color:var(--accent-primary);border-radius:9999px;padding:5px 12px;font-size:12px;font-weight:600;letter-spacing:.02em}} .pill.cyan{{background:#ECFEFF;color:var(--accent-cyan)}}
.crisis .pill.rule{{background:#FFF1F2;color:var(--alert-urgent)}} .counter{{font-size:11px;color:var(--text-body)}} h2{{margin:18px 0 5px;color:var(--text-heading);font-size:34px;line-height:1.12;font-weight:800;letter-spacing:-.02em}} .question{{margin:0;font-size:16px;line-height:1.6}}
.body{{position:absolute;left:4%;right:4%;top:30%;bottom:16%;display:grid;grid-template-columns:repeat(2,minmax(0,1fr));gap:16px;align-content:center}} .card{{position:relative;background:var(--bg-surface);border:1px solid var(--border);border-radius:12px;box-shadow:0 10px 25px -5px rgba(0,0,0,.05);padding:20px 22px;font-size:17px;line-height:1.5;min-height:88px}} .card:before{{content:'';position:absolute;left:0;top:16px;bottom:16px;width:4px;background:var(--accent-primary);border-radius:8px}} .card:nth-child(even):before{{background:var(--accent-cyan)}} .card b{{display:block;color:var(--text-heading);margin-bottom:8px}} .mini{{display:block;color:var(--accent-primary);font-size:12px;font-weight:800;text-transform:uppercase;letter-spacing:.05em;margin-bottom:8px}} .cyan-text{{color:var(--accent-cyan)}} .compact{{font-size:15px;min-height:72px}} .stack{{display:grid;gap:10px}} .source-card{{position:relative;background:var(--bg-surface);border:1px solid var(--border);border-radius:12px;padding:12px;box-shadow:0 10px 25px -5px rgba(0,0,0,.05)}} .source-img{{width:100%;height:100%;max-height:300px;object-fit:contain}} .source-card .pill{{position:absolute;left:18px;bottom:18px}}
.chain{{grid-column:1/-1;text-align:center;min-height:auto}} .rubric-note{{grid-column:1/-1;text-align:center;font-size:14px;color:var(--text-body)}} .decision-corner{{position:absolute;right:0;top:0;width:29%;background:#EEF2FF;border:1px solid #C7D2FE;border-radius:12px;padding:13px 16px;font-size:12px;line-height:1.4}} .decision-corner b{{display:block;color:var(--accent-primary);font-size:11px;margin-bottom:5px}}
.taskbar{{position:absolute;left:2.8%;right:2.8%;bottom:2.2%;min-height:82px;background:var(--bg-surface);border:1px solid var(--border);border-radius:12px;box-shadow:0 10px 25px -5px rgba(0,0,0,.06);padding:18px 20px 14px;display:flex;justify-content:space-between;align-items:center;gap:20px}} .taskbar:before{{content:'';position:absolute;left:0;right:0;top:0;height:4px;border-radius:12px 12px 0 0;background:linear-gradient(90deg,var(--accent-primary),var(--accent-cyan))}} .taskbar>div{{display:flex;gap:16px;align-items:center;min-width:0}} .taskbar>div>b{{color:var(--accent-primary);font-size:12px;white-space:nowrap}} .taskbar>div>span{{color:var(--text-heading);font-size:15px;line-height:1.4}} .timebox{{display:inline-flex;align-items:center;gap:7px;background:#FFF1F2;color:var(--alert-urgent);border:1px solid #FECDD3;border-radius:9999px;padding:7px 11px;font-size:11px;white-space:nowrap}} .timebox svg{{width:17px;height:17px}} .stage-indicator{{position:absolute;right:12px;bottom:7px;font-size:9px;color:#94A3B8}}
.stage2,.stage3{{opacity:0;transform:translateY(16px);pointer-events:none;transition:opacity .28s ease,transform .28s ease}} .slide[data-stage='2'] .stage2,.slide[data-stage='3'] .stage2,.slide[data-stage='3'] .stage3{{opacity:1;transform:translateY(0);pointer-events:auto}} .slide[data-stage='3'] .timebox{{animation:pulse 1.45s ease-in-out infinite}} @keyframes pulse{{0%,100%{{box-shadow:0 0 0 0 rgba(244,63,94,.05)}}50%{{box-shadow:0 0 0 6px rgba(244,63,94,.10)}}}}
.cover{{padding:6.2%;display:flex;flex-direction:column;justify-content:center}} .cover h1{{max-width:78%;font-size:52px;line-height:1.05;color:var(--text-heading);margin:22px 0 14px;font-weight:800;letter-spacing:-.03em}} .cover p{{font-size:19px;max-width:70%;line-height:1.6}} .cover-grid{{display:grid;grid-template-columns:1fr 1fr;gap:18px;margin-top:28px;max-width:82%}} .cover-accent{{position:absolute;right:5%;top:7%;width:18%;height:10px;background:linear-gradient(90deg,var(--accent-primary),var(--accent-cyan));border-radius:999px}}
@media (prefers-reduced-motion:reduce){{.stage2,.stage3{{transition:none}}.slide[data-stage='3'] .timebox{{animation:none}}}} @media print{{body{{background:white}}.deck{{padding:0}}.slide{{margin:0;border:0;border-radius:0;box-shadow:none;break-after:page}}.stage2,.stage3{{opacity:1!important;transform:none!important}}.stage-indicator{{display:none}}@page{{size:16in 9in;margin:0}}}}
"""
    js = """
<script>
const slides=[...document.querySelectorAll('.slide')];let active=0;
function setStage(s,n){n=Math.max(1,Math.min(3,n));s.dataset.stage=String(n);const i=s.querySelector('.stage-indicator');if(i)i.textContent=n+'/3';}
function advance(s){const n=Number(s.dataset.stage||1);if(n<3)setStage(s,n+1);else{const i=slides.indexOf(s);if(i<slides.length-1){active=i+1;slides[active].scrollIntoView({behavior:'smooth',block:'center'});}}}
slides.forEach((s,i)=>{s.addEventListener('click',e=>{if(e.target.closest('a,button'))return;active=i;advance(s);});});
const io=new IntersectionObserver(es=>{const v=es.filter(e=>e.isIntersecting).sort((a,b)=>b.intersectionRatio-a.intersectionRatio)[0];if(v)active=slides.indexOf(v.target);},{threshold:[.35,.6,.8]});slides.forEach(s=>io.observe(s));
document.addEventListener('keydown',e=>{const s=slides[active];if(!s)return;if([' ','ArrowRight','Enter'].includes(e.key)){e.preventDefault();advance(s);}else if(e.key==='ArrowLeft'){e.preventDefault();setStage(s,Number(s.dataset.stage||1)-1);}else if(['ArrowDown','PageDown'].includes(e.key)&&active<slides.length-1){e.preventDefault();active++;slides[active].scrollIntoView({behavior:'smooth',block:'center'});}else if(['ArrowUp','PageUp'].includes(e.key)&&active>0){e.preventDefault();active--;slides[active].scrollIntoView({behavior:'smooth',block:'center'});}});
</script>
"""
    return "<!doctype html><html><head><meta charset='utf-8'><meta name='viewport' content='width=device-width,initial-scale=1'><title>ISCARB ZTM Presenter</title><style>" + css + "</style></head><body><main class='deck'>" + "".join(pages) + "</main>" + js + "</body></html>"


# Keep originals for non-overridden unit semantics.
_ORIGINAL_PPT_SEMANTIC = presenter._ppt_semantic
_ORIGINAL_PDF_SEMANTIC = presenter._pdf_semantic


def _install_inter_pdf_font():
    try:
        roots = list(Path("/usr/share/fonts").rglob("Inter-Regular.ttf"))
        bolds = list(Path("/usr/share/fonts").rglob("Inter-Bold.ttf"))
        if roots:
            pdfmetrics.registerFont(TTFont("ISCARB-Inter", str(roots[0]))); presenter._FONT = "ISCARB-Inter"
        if bolds:
            pdfmetrics.registerFont(TTFont("ISCARB-Inter-Bold", str(bolds[0]))); presenter._FONT_BOLD = "ISCARB-Inter-Bold"
    except Exception:
        pass


def _stabilize_golden_specs(bp, profile):
    specs = golden._decode_specs(bp)
    if not specs:
        specs = golden._ch10_expansions(profile) if golden._ch10_signature(profile) else golden._generic_expansions(profile)
    if specs:
        notes = [str(x) for x in list(getattr(bp, "release_notes", []) or [])
                 if not str(x).startswith(golden._NOTE_PREFIX) and "ZTM theme" not in str(x)]
        encoded = golden._NOTE_PREFIX + json.dumps(specs, ensure_ascii=False, separators=(",", ":"))
        bp.release_notes = notes[:18] + [encoded, "ZTM theme v7.2.9: white high-contrast floating-card presenter with progressive disclosure in HTML."]
    return bp


def apply_v729_ztm_theme_patch(app):
    global _PATCHED, _PREVIOUS_DRAFT
    if _PATCHED:
        return
    _PATCHED = True
    _PREVIOUS_DRAFT = engine._source_preserving_draft

    def draft(profile, bundle):
        return _stabilize_golden_specs(_PREVIOUS_DRAFT(profile, bundle), profile)

    engine._source_preserving_draft = draft
    base.engine._source_preserving_draft = draft

    # Public design-token contract.
    contract.chapter_design_tokens = ztm_chapter_design_tokens
    prod.chapter_design_tokens = ztm_chapter_design_tokens
    contract.physical_slide_plan = _ztm_contract_plan
    prod.physical_slide_plan = _ztm_contract_plan

    # Re-theme the existing portable renderer. Its export functions resolve these
    # globals/functions at call time, so no duplicate export pipeline is introduced.
    presenter.BG = BG_BASE
    presenter.PANEL = BG_SURFACE
    presenter.PANEL2 = BG_SURFACE
    presenter.TEXT = TEXT_HEADING
    presenter.MUTED = TEXT_BODY
    presenter.CYAN = ACCENT_CYAN
    presenter.MAGENTA = ACCENT_PRIMARY
    presenter.GOLD = ACCENT_GOLD
    presenter.GREEN = ACCENT_GREEN
    presenter.BLUE = ACCENT_BLUE
    presenter.DANGER = ACCENT_PRIMARY  # urgent rose is reserved for CRISIS/TIMEBOX only.
    presenter.LINE = BORDER
    presenter.PHASE_ACCENT = {"IFHAM": ACCENT_CYAN, "MARIS": ACCENT_PRIMARY, "ATQAN": ACCENT_GOLD, "MAYYIZ": ACCENT_BLUE}
    presenter.RULE_NAMES[11] = "Micro-case → Saudi/local transfer"
    presenter.RULE_NAMES[19] = "Peer-review quick card"
    presenter._physical_plan = _physical_plan
    presenter._ppt_text = _ppt_text
    presenter._ppt_box = _ppt_box
    presenter._ppt_header = _ppt_header
    presenter._ppt_footer = _ppt_footer
    presenter._ppt_source_split = _ppt_source_split
    presenter._ppt_semantic = _ppt_semantic
    presenter._ppt_expansion = _ppt_expansion
    presenter._ppt_cover = _ppt_cover
    presenter._ppt_close = _ppt_close
    presenter._pdf_box = _pdf_box
    presenter._pdf_header = _pdf_header
    presenter._pdf_footer = _pdf_footer
    presenter._pdf_semantic = _pdf_semantic
    presenter._pdf_expansion = _pdf_expansion
    presenter._pdf_cover = _pdf_cover
    presenter._pdf_close = _pdf_close
    presenter.render_presenter_preview = render_presenter_preview_ztm
    _install_inter_pdf_font()

    # The v4.4 routes look up these globals at request time.
    base.render_presenter_preview = render_presenter_preview_ztm
    base.export_presenter_pptx = presenter.export_presenter_pptx
    base.export_presenter_pdf = presenter.export_presenter_pdf

    previous_health = base._health_v440
    def health():
        data = dict(previous_health())
        data.update({
            "version": "7.2.9",
            "ztm_theme_version": "v7.2.9",
            "golden_theme": "ZTM-inspired high-contrast white surface over Golden v6.6 grammar",
            "ztm_tokens": {
                "bg_base": BG_BASE, "bg_surface": BG_SURFACE,
                "text_heading": TEXT_HEADING, "text_body": TEXT_BODY,
                "accent_primary": ACCENT_PRIMARY, "accent_cyan": ACCENT_CYAN,
                "alert_urgent": ALERT_URGENT,
            },
            "ztm_typography": "Inter preferred for HTML/PPTX/PDF; embedded local Inter when available, DejaVu fallback for PDF portability.",
            "ztm_cards": "12px rounded floating cards, subtle shadow, neutral border, semantic accent rail.",
            "ztm_taskbar": "Fixed lower YOUR TASK card with primary→cyan top rule and urgent timebox pill.",
            "ztm_urgent_usage": "#F43F5E reserved for CRISIS and TIMEBOX only.",
            "progressive_disclosure": "HTML presenter: stage 1 title/question, stage 2 theory/source, stage 3 task/timebox. PPTX/PDF export the final state for portability.",
            "source_expansion_ergonomics": "Every semantic P1 expansion carries a visible Decision → Evidence callout plus readable source cards.",
            "rule19_learner_surface": "Two-question peer-review quick card; full 6x4 rubric retained in instructor/blueprint layer.",
            "rule11_learner_surface": "Micro-case first, then Saudi/local transfer using the same mechanism→evidence→decision-boundary chain.",
        })
        return data
    base._health_v440 = health
    base.engine.health = health
