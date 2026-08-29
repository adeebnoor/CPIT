from __future__ import annotations

import html
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .models import Blueprint, LectureUnit
from . import visual_engine as ve


def _clean(s: str, n: int = 82) -> str:
    s = " ".join(str(s or "").split())
    return s if len(s) <= n else s[: n - 1].rstrip() + "…"


def _items(u: LectureUnit, n: int = 6) -> list[str]:
    src = list(u.core_content) or list(u.pedagogy_content)
    return [_clean(x, 92) for x in src[:n]]


def _primary_type(u: LectureUnit) -> str:
    if u.knowledge_types:
        return u.knowledge_types[0]
    return "CONCEPT"


def cimt_visual_html(bp: Blueprint, u: LectureUnit, fallback) -> str:
    """Browser visual for source-native technical knowledge.

    Reserved ISCARB Units retain the established visual grammar. Technical Units
    6-10 adapt to the knowledge type so computing subjects do not all look alike.
    """
    if u.number not in range(6, 11):
        return fallback(bp, u)
    t = _primary_type(u)
    xs = _items(u, 6)
    q = html.escape(_clean(u.engineering_question, 110))
    e = lambda s: html.escape(_clean(s, 82))

    if t == "ALGORITHM":
        steps = xs[:4] or ["Problem", "Invariant", "Trace", "Complexity"]
        return '<div style="display:grid;grid-template-columns:repeat(4,1fr);gap:14px;align-items:center">' + ''.join(
            f'<div style="text-align:center"><div style="width:54px;height:54px;border-radius:50%;display:grid;place-items:center;margin:auto;background:#e9f6ef;border:2px solid #1d8b56;font-weight:950">{i+1}</div><b style="display:block;margin:10px 0 5px">{e(s)}</b><small style="color:#69756e">TRACE STEP</small></div>'
            for i,s in enumerate(steps)
        ) + f'</div><div style="margin-top:16px;border-top:1px solid #d8e0da;padding-top:11px;font-size:12px;color:#516159"><b>DECISION TEST ·</b> {q}</div>'

    if t == "CODE":
        left = xs[:3] or ["Source fragment", "State change", "Observed output"]
        right = xs[3:6] or ["Input", "Execution state", "Failure / mutation"]
        return '<div style="display:grid;grid-template-columns:1.15fr .85fr;gap:18px;height:100%"><div style="background:#102526;color:#eaf6f2;border-radius:14px;padding:18px;font-family:ui-monospace,SFMono-Regular,Consolas,monospace">' + ''.join(f'<div style="padding:9px 0;border-bottom:1px solid #294447"><span style="color:#5ed8dc">{i+1:02d}</span> &nbsp; {e(x)}</div>' for i,x in enumerate(left)) + '</div><div style="display:grid;gap:10px">' + ''.join(f'<div style="border-left:4px solid #d2a753;background:#fbf7ed;padding:13px"><b>STATE {i+1}</b><div style="font-size:12px;margin-top:4px;color:#59645e">{e(x)}</div></div>' for i,x in enumerate(right)) + '</div></div>'

    if t == "EQUATION":
        vals = xs[:3] or ["Known quantities", "Derive relationship", "Interpret sensitivity"]
        return '<div style="display:grid;grid-template-columns:1fr auto 1fr auto 1fr;gap:12px;align-items:center;text-align:center">' + f'<div><small>KNOWN</small><strong style="display:block;font-size:19px;margin-top:9px">{e(vals[0])}</strong></div><div style="font-size:30px;color:#a38b5c">→</div><div style="border:2px solid #0a353e;border-radius:16px;padding:24px 12px;background:#edf6f5"><small>DERIVE</small><strong style="display:block;font-size:22px;margin-top:8px">{e(vals[1])}</strong></div><div style="font-size:30px;color:#a38b5c">→</div><div><small>INTERPRET</small><strong style="display:block;font-size:19px;margin-top:9px">{e(vals[2])}</strong></div></div>'

    if t == "PROTOCOL":
        msgs = xs[:4] or ["Request", "Validate", "Respond", "Failure path"]
        return '<div style="display:grid;grid-template-columns:150px 1fr 150px;gap:16px;align-items:stretch"><div style="border:2px solid #563c7d;border-radius:14px;display:grid;place-items:center;font-weight:900">ACTOR / LAYER A</div><div style="display:grid;gap:10px">' + ''.join(f'<div style="position:relative;border-top:2px solid #1d8b56;padding-top:7px;text-align:center;font-size:12px"><b>{i+1:02d}</b> · {e(m)} <span style="float:right">→</span></div>' for i,m in enumerate(msgs)) + '</div><div style="border:2px solid #0a353e;border-radius:14px;display:grid;place-items:center;font-weight:900">ACTOR / LAYER B</div></div>'

    if t == "DATA_MODEL":
        ents = xs[:4] or ["Entity A", "Entity B", "Constraint", "Query/use"]
        return '<div style="display:grid;grid-template-columns:1fr 120px 1fr;grid-template-rows:1fr 1fr;gap:12px;align-items:center">' + f'<div style="border:2px solid #563c7d;border-radius:12px;padding:18px"><b>ENTITY / TABLE</b><div>{e(ents[0])}</div></div><div style="text-align:center;font-size:28px">⇄</div><div style="border:2px solid #1d8b56;border-radius:12px;padding:18px"><b>ENTITY / TABLE</b><div>{e(ents[1])}</div></div><div style="border-left:4px solid #d2a753;padding:13px"><b>CONSTRAINT</b><div>{e(ents[2])}</div></div><div style="text-align:center;color:#6e756f">SCHEMA</div><div style="border-left:4px solid #0a353e;padding:13px"><b>QUERY / USE</b><div>{e(ents[3])}</div></div></div>'

    if t == "SYSTEM_BEHAVIOR":
        states = xs[:4] or ["State A", "Event", "State B", "Observable consequence"]
        return '<div style="display:flex;align-items:center;justify-content:space-around;gap:10px">' + ''.join(f'<div style="display:flex;align-items:center;gap:10px"><div style="width:145px;height:145px;border-radius:50%;border:2px solid #0a353e;background:#edf6f5;display:grid;place-items:center;text-align:center;padding:16px;font-weight:850">{e(s)}</div>{"<div style=\"font-size:30px;color:#d2a753\">→</div>" if i < len(states)-1 else ""}</div>' for i,s in enumerate(states)) + '</div>'

    if t == "TRADE_OFF":
        a = xs[0] if xs else "Alternative A"
        b = xs[1] if len(xs)>1 else "Alternative B"
        crit = xs[2:5] or ["Evidence", "Cost", "Risk"]
        return f'<div style="display:grid;grid-template-columns:1fr 1.1fr 1fr;gap:16px;align-items:stretch"><div style="border:2px solid #563c7d;border-radius:15px;padding:20px"><small>ALTERNATIVE A</small><b style="display:block;font-size:18px;margin-top:10px">{e(a)}</b></div><div style="display:grid;gap:8px">' + ''.join(f'<div style="display:grid;grid-template-columns:1fr auto 1fr;gap:8px;align-items:center;border-bottom:1px solid #ddd;padding:7px"><span>◀</span><b>{e(c)}</b><span>▶</span></div>' for c in crit) + f'</div><div style="border:2px solid #1d8b56;border-radius:15px;padding:20px"><small>ALTERNATIVE B</small><b style="display:block;font-size:18px;margin-top:10px">{e(b)}</b></div></div>'

    if t == "EMPIRICAL_RESULT":
        rows = xs[:4] or ["Setup", "Measure", "Observed result", "Uncertainty"]
        return '<div style="display:grid;grid-template-columns:220px 1fr;gap:14px;align-items:center">' + ''.join(f'<><div></div></>' for _ in []) + ''.join(f'<div style="font-size:11px;font-weight:900;color:#5f6962">{["SETUP","MEASURE","RESULT","UNCERTAINTY"][i]}</div><div style="height:48px;background:linear-gradient(90deg,#e9f6ef,#f8f4e9);border-left:5px solid {["#563c7d","#0a353e","#1d8b56","#d2a753"][i]};padding:13px;font-size:12px">{e(x)}</div>' for i,x in enumerate(rows)) + '</div>'

    if t in {"ARCHITECTURE", "PROCESS", "DESIGN_PRINCIPLE"}:
        nodes = xs[:4] or ["Input", "Mechanism", "Boundary", "Decision"]
        label = {"ARCHITECTURE":"COMPONENT / FLOW","PROCESS":"STAGE / HANDOFF","DESIGN_PRINCIPLE":"PRESSURE → PRINCIPLE"}[t]
        return '<div style="display:grid;grid-template-columns:repeat(4,1fr);gap:14px;align-items:center">' + ''.join(f'<div style="min-height:150px;border-top:5px solid {"#563c7d" if i%2==0 else "#1d8b56"};background:#fff;border-radius:10px;padding:18px;box-shadow:0 10px 28px rgba(0,0,0,.06)"><small>{label} {i+1}</small><b style="display:block;font-size:16px;margin-top:12px">{e(x)}</b></div>' for i,x in enumerate(nodes)) + '</div>'

    return fallback(bp, u)


def _circle(slide, x, y, d, title, body, line):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(250,249,246)
    sh.line.color.rgb = line; sh.line.width = Pt(1.8)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=_clean(title,28); p.font.name="Aptos"; p.font.size=Pt(11); p.font.bold=True; p.font.color.rgb=line; p.alignment=PP_ALIGN.CENTER
    p=tf.add_paragraph(); p.text=_clean(body,62); p.font.name="Aptos"; p.font.size=Pt(8); p.font.color.rgb=ve.MUTED; p.alignment=PP_ALIGN.CENTER
    return sh


def _render_native(slide, bp: Blueprint, u: LectureUnit):
    t=_primary_type(u); xs=_items(u,6); c=ve.PHASE_COLOR[u.phase]; soft=ve.PHASE_SOFT[u.phase]
    if t=="ALGORITHM":
        vals=xs[:4] or ["Problem","Invariant","Trace","Complexity"]
        for i,v in enumerate(vals):
            _circle(slide,.75+i*3.05,2.55,1.55,str(i+1),v,c)
            if i<3: ve._arrow(slide,2.36+i*3.05,3.2,.78,.22,c)
        ve._text(slide,.75,5.25,11.8,.55,"TRACE THE CHANGE · then justify time/space or decision trade-offs from the source.",size=12,color=ve.MUTED,bold=True,align=PP_ALIGN.CENTER); return
    if t=="EQUATION":
        vals=xs[:3] or ["Known quantities","Derive relationship","Interpret sensitivity"]
        ve._text(slide,.8,2.45,3.35,1.1,vals[0],size=18,bold=True,align=PP_ALIGN.CENTER)
        ve._text(slide,4.25,2.64,.7,.5,"→",size=28,color=c,bold=True,align=PP_ALIGN.CENTER)
        ve._box(slide,5.0,2.2,3.35,1.75,"DERIVE",vals[1],fill=soft,line=c,title_color=c,title_size=13,body_size=12)
        ve._text(slide,8.42,2.64,.7,.5,"→",size=28,color=c,bold=True,align=PP_ALIGN.CENTER)
        ve._text(slide,9.2,2.45,3.15,1.1,vals[2],size=18,bold=True,align=PP_ALIGN.CENTER)
        ve._box(slide,2.3,4.7,8.7,1.0,"SENSITIVITY / INTERPRETATION",u.student_action,fill=ve.WHITE,line=ve.AMBER,title_color=ve.AMBER,title_size=11,body_size=9.5); return
    if t=="PROTOCOL":
        ve._box(slide,.7,2.25,2.0,3.7,"ACTOR / LAYER A","Start state",fill=soft,line=c,title_color=c,title_size=12,body_size=10)
        ve._box(slide,10.65,2.25,2.0,3.7,"ACTOR / LAYER B","Receiving state",fill=soft,line=c,title_color=c,title_size=12,body_size=10)
        msgs=xs[:4] or ["Request","Validate","Respond","Failure path"]
        for i,m in enumerate(msgs):
            y=2.65+i*.72; ve._text(slide,3.0,y,7.2,.35,f"{i+1:02d} · {m}",size=10.5,bold=True,align=PP_ALIGN.CENTER); ve._arrow(slide,9.5,y+.03,.65,.18,c)
        return
    if t=="SYSTEM_BEHAVIOR":
        vals=xs[:4] or ["State A","Event","State B","Consequence"]
        for i,v in enumerate(vals):
            _circle(slide,.65+i*3.08,2.45,1.75,"STATE" if i!=1 else "EVENT",v,c)
            if i<3: ve._arrow(slide,2.45+i*3.08,3.2,.72,.2,c)
        return
    if t=="TRADE_OFF":
        a=xs[0] if xs else "Alternative A"; b=xs[1] if len(xs)>1 else "Alternative B"
        ve._box(slide,.7,2.25,3.4,3.3,"ALTERNATIVE A",a,fill=ve.SOFT_VIOLET,line=ve.VIOLET,title_color=ve.VIOLET,title_size=15,body_size=12)
        ve._box(slide,9.2,2.25,3.4,3.3,"ALTERNATIVE B",b,fill=ve.SOFT_GREEN,line=ve.GREEN,title_color=ve.GREEN,title_size=15,body_size=12)
        crit=xs[2:5] or ["Evidence","Cost","Risk"]
        for i,x in enumerate(crit): ve._box(slide,4.45,2.35+i*1.05,4.4,.8,f"CRITERION {i+1}",x,fill=ve.WHITE,line=ve.AMBER,title_color=ve.AMBER,title_size=10,body_size=8.5)
        return
    if t=="CODE":
        left=xs[:3] or ["Source fragment","State change","Observed output"]; right=xs[3:6] or ["Input","Execution state","Failure / mutation"]
        sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(.7),Inches(2.2),Inches(7.1),Inches(3.8)); sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(16,37,38); sh.line.fill.background()
        for i,x in enumerate(left): ve._text(slide,1.0,2.55+i*.82,6.5,.55,f"{i+1:02d}   {_clean(x,78)}",size=11,color=RGBColor(234,246,242),bold=False)
        for i,x in enumerate(right): ve._box(slide,8.15,2.2+i*1.25,4.45,1.0,f"STATE {i+1}",x,fill=soft,line=c,title_color=c,title_size=10.5,body_size=9)
        return
    # Architecture, process, data model, design principle, empirical result and concept
    labels={"ARCHITECTURE":"COMPONENT","PROCESS":"STAGE","DATA_MODEL":"ENTITY / CONSTRAINT","DESIGN_PRINCIPLE":"PRESSURE / PRINCIPLE","EMPIRICAL_RESULT":"EVIDENCE","CONCEPT":"CONCEPT"}
    vals=xs[:4] or ["Input","Mechanism","Boundary","Decision"]
    for i,v in enumerate(vals):
        x=.65+i*3.08; ve._box(slide,x,2.45,2.7,2.75,f"{labels.get(t,'ELEMENT')} {i+1}",v,fill=soft if i%2==0 else ve.WHITE,line=c,title_color=c,title_size=11.5,body_size=10.5)
        if i<3: ve._arrow(slide,x+2.74,3.65,.28,.2,c)
    ve._text(slide,.75,5.68,11.7,.5,u.visual_plan.visual_evidence_role if u.visual_plan else u.takeaway,size=10.5,color=ve.MUTED,bold=True,align=PP_ALIGN.CENTER)


def export_cimt_presenter_pptx(bp: Blueprint, out: Path) -> Path:
    """CIMT+ PPTX: established ISCARB reserved visuals + source-native technical visuals."""
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    for u in bp.units:
        slide=ve._blank(prs); ve._base(slide,u)
        if u.number in range(6,11):
            _render_native(slide,bp,u)
        else:
            renderer=ve._RENDERERS.get(u.number)
            if renderer: renderer(slide,bp,u)
            else: ve._generic_cards(slide,u)
    prs.save(out); return out
