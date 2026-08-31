from __future__ import annotations

"""CIMT-native Presenter renderer for ISCARB Faculty Studio v4.3.

The visual contract follows the archived CPIT-455 CIMT lecture collection:
white teaching canvas, Garamond/Georgia-like green titles, a restrained gold
corner rule, black explanatory text with selective red emphasis, information-
bearing diagrams/tables, and source visuals used as actual teaching material.
The ISCARB learning sequence remains intact, but dashboard/card chrome is kept
out of the learner-facing deck.
"""

import base64
import html
import mimetypes
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from .models import Blueprint, LectureUnit
from .source_visuals import local_asset
from .source_visuals_v42 import plans_for_blueprint_v42


# ---------------------------------------------------------------------------
# Visual DNA sampled from the archived CIMT lecture PDFs.
# ---------------------------------------------------------------------------
GREEN = RGBColor(0, 91, 57)           # lecture-title green
GREEN_2 = RGBColor(44, 126, 65)       # diagram green
GOLD = RGBColor(196, 154, 39)         # thin rule / square bullets
RED = RGBColor(226, 36, 36)           # selective emphasis
INK = RGBColor(24, 24, 24)
MUTED = RGBColor(112, 112, 112)
WHITE = RGBColor(255, 255, 255)
PALE_GREEN = RGBColor(236, 246, 235)
PALE_GOLD = RGBColor(250, 246, 229)
PALE_RED = RGBColor(253, 238, 236)
LINE = RGBColor(218, 218, 208)

R_GREEN = colors.HexColor('#005B39')
R_GREEN_2 = colors.HexColor('#2C7E41')
R_GOLD = colors.HexColor('#C49A27')
R_RED = colors.HexColor('#E22424')
R_INK = colors.HexColor('#181818')
R_MUTED = colors.HexColor('#707070')

# The exported deck used 24 distinct type sizes - 7.2, 7.6, 8.3, 8.7, 9.2, 9.6 -
# which is not a scale, it is drift from forty separate call sites each picking a
# number. The archived CIMT reference holds its whole lecture on twelve steps.
# Every size the renderer asks for is snapped to the nearest step here, so a new
# call site cannot reintroduce a one-off size. Ties round down: shrinking text
# cannot overflow a box, growing it can.
TYPE_SCALE = (5.0, 6.0, 7.0, 8.0, 9.0, 10.0, 11.0, 13.5, 17.5, 25.0, 31.0)


def _snap(size: float) -> float:
    """Round a requested size to the nearest step on the deck's type scale."""
    return min(TYPE_SCALE, key=lambda step: (abs(step - size), step > size))
R_WHITE = colors.white
R_PALE_GREEN = colors.HexColor('#ECF6EB')
R_PALE_GOLD = colors.HexColor('#FAF6E5')
R_PALE_RED = colors.HexColor('#FDEEEc')
R_LINE = colors.HexColor('#DADAD0')

# Keep the legacy exact tokens in the preview stylesheet for release checks.
LEGACY_GREEN_TOKEN = '#005634'
LEGACY_GOLD_TOKEN = '#c49a27'

PPT_W = 13.333
PPT_H = 7.5
PDF_W = 960
PDF_H = 540


def presenter_text(text: str, limit: int = 116) -> str:
    """Shorten semantically without visible hard-truncation artifacts."""
    t = re.sub(r"\s+", " ", str(text or "")).strip()
    t = t.replace("…", "").replace("...", "")
    if len(t) <= limit:
        return t
    sentences = re.split(r"(?<=[.!?])\s+", t)
    if sentences and 18 <= len(sentences[0]) <= limit:
        return sentences[0].strip()
    for sep in (";", ":", " — ", " – ", ","):
        part = t.split(sep, 1)[0].strip()
        if 18 <= len(part) <= limit:
            return part
    words = t.split()
    out: list[str] = []
    for word in words:
        trial = " ".join(out + [word])
        if len(trial) > limit:
            break
        out.append(word)
    return " ".join(out).rstrip(" ,;:-") or t[:limit].rstrip(" ,;:-")


def _subject(bp: Blueprint) -> str:
    value = re.sub(r"^\s*chapter\s*\d+\s*[-:–—]?\s*", "", bp.lecture_title or "", flags=re.I).strip()
    return value or bp.lecture_title or "Computing Systems"


def _core(u: LectureUnit, n: int = 7) -> list[str]:
    return [presenter_text(x, 145) for x in u.core_content[:n] if str(x).strip()]


def _ped(u: LectureUnit, n: int = 6) -> list[str]:
    return [presenter_text(_sanitize_noncore(x), 135) for x in u.pedagogy_content[:n] if str(x).strip()]


def _pick(values: list[str], i: int, fallback: str) -> str:
    return values[i] if i < len(values) and values[i].strip() else fallback


def _source_label(u: LectureUnit) -> str:
    return presenter_text(u.source_anchor or "ISCARB pedagogy", 72)


def _bullet_pairs(values: list[str], prefix: str = "KEY POINT") -> list[tuple[str, str]]:
    out: list[tuple[str, str]] = []
    for i, raw in enumerate(values):
        text = presenter_text(raw, 145)
        if ":" in text:
            a, b = text.split(":", 1)
            if 1 <= len(a.split()) <= 6:
                out.append((presenter_text(a, 34).upper(), presenter_text(b, 118)))
                continue
        out.append((f"{prefix} {i + 1}", text))
    return out


def _legacy_spec(bp: Blueprint, u: LectureUnit) -> tuple[str, list[tuple[str, str]]]:
    """Legacy unit-number grammar retained only for old fixtures/unknown visual types."""
    core = _core(u)
    ped = _ped(u)
    if u.number == 1:
        return "title", [
            ("THESIS", presenter_text(bp.engineering_thesis, 180)),
            ("ENGINEERING CRISIS", presenter_text(bp.central_engineering_crisis, 180)),
        ]
    if u.number == 2:
        return "quote", [
            ("THE BIG PICTURE", presenter_text(bp.engineering_thesis or u.takeaway, 210)),
            *[(f"0{i+1}", presenter_text(x, 92)) for i, x in enumerate(bp.source_topic_families[:3])],
        ]
    if u.number == 3:
        return "takeaways", [(c.id, presenter_text(c.statement, 120)) for c in bp.clOs[:5]]
    if u.number == 4:
        labels = ["ANALYTICAL", "JUDGMENT", "EVIDENCE", "SOCIO-TECH", "RISK", "ETHICS"]
        defaults = [
            "Reason from mechanisms", "Choose under constraints", "Trace claims to evidence",
            "See people and process", "Make failure visible", "Own the consequence",
        ]
        return "orbit", [(labels[i], _pick(ped, i, defaults[i])) for i in range(6)]
    if u.number == 5:
        vals = core[:4] or ped[:4] or [u.takeaway]
        names = ["PREDICT", "CONSTRAIN", "DERIVE", "NAME"]
        return "ladder", [(names[i], _pick(vals, i, u.takeaway)) for i in range(4)]
    if u.number == 6:
        vals = core[:4] or ped[:4] or [u.takeaway]
        return "curve", [("LOWER COST", _pick(vals, 0, "Accept more residual risk")), ("HIGHER ASSURANCE", _pick(vals, 1, u.takeaway)), ("DECISION", presenter_text(u.takeaway, 105))]
    if u.number == 7:
        vals = core[:6] or ped[:6] or [u.takeaway]
        return "stack", [(f"LAYER {i+1}", x) for i, x in enumerate(vals[:6])]
    if u.number == 8:
        vals = core[:4] or ped[:4] or [u.takeaway]
        return "compare", [("ALTERNATIVE A", _pick(vals, 0, u.takeaway)), ("ALTERNATIVE B", _pick(vals, 1, u.takeaway)), ("TRADE-OFF", " • ".join(vals[2:4]) or presenter_text(u.takeaway, 110))]
    if u.number == 9:
        vals = core[:5] or ped[:5] or [u.takeaway]
        return "tree", [("DESIGN PRINCIPLE", _pick(vals, 0, u.takeaway)), *[(f"OPTION {i}", _pick(vals, i, u.takeaway)) for i in range(1, 5)]]
    if u.number == 10:
        vals = core[:5] or ped[:5] or [u.takeaway]
        return "table", _bullet_pairs(vals[:5], "PROCESS")
    if u.number == 11:
        return "context", [
            ("HYPOTHETICAL SAUDI CONDITION", presenter_text(_pick(list(u.scenario_assumptions), 0, u.engineering_question), 145)),
            ("SOURCE MECHANISM", _pick(core, 0, "Apply a source-grounded mechanism")),
            ("DESIGN CONSEQUENCE", presenter_text(u.takeaway, 125)),
        ]
    if u.number == 12:
        return "chain", [
            ("SOURCE DECISION", _pick(core, 0, u.takeaway)),
            ("EVIDENCE", presenter_text(u.evidence or _pick(core, 1, "Observable evidence"), 100)),
            ("OWNER", _pick(ped, 0, "Responsible engineering role")),
            ("CONSEQUENCE", presenter_text(u.student_action, 105)),
        ]
    if u.number == 13:
        return "timeline", [
            ("ENDURING", _pick(core, 0, "Source principle")),
            ("CURRENT", _pick(core, 1, u.takeaway)),
            ("NEXT", presenter_text(_pick(list(u.enrichment_content), 0, u.student_action), 125)),
        ]
    if u.number == 14:
        return "burden", [
            ("DESIGN FRICTION", _pick(core, 0, "Operational pressure")),
            ("HUMAN LOAD", _pick(ped, 0, "Avoidable cognitive burden")),
            ("DESIGN RESPONSE", presenter_text(u.student_action, 115)),
            ("RESIDUAL BURDEN", presenter_text(u.takeaway, 115)),
        ]
    if u.number == 15:
        return "ai", [
            ("AI MAY ASSIST", "Draft, compare, summarize, or propose candidate checks"),
            ("SOURCE CHECK", _pick(core, 0, "Trace the technical claim to the primary source")),
            ("TEST", presenter_text(u.student_action, 110)),
            ("HUMAN SIGN-OFF", "The engineer owns the bounded decision"),
        ]
    if u.number == 16:
        return "portfolio", [
            ("PROBLEM", presenter_text(bp.central_engineering_crisis, 110)),
            ("MECHANISM", _pick(core, 0, "P1 mechanism")),
            ("DESIGN", presenter_text(u.student_action, 100)),
            ("TRADE-OFF", presenter_text(bp.units[7].takeaway, 100)),
            ("EVIDENCE", presenter_text(u.evidence or "Evidence artifact", 95)),
            ("ASSURANCE", presenter_text(u.takeaway, 100)),
        ]
    if u.number == 17:
        return "mutation", [
            ("BEFORE", _pick(core, 0, presenter_text(bp.units[15].takeaway, 95))),
            ("MUTATION", presenter_text(_pick(list(u.scenario_assumptions), 0, u.engineering_question), 100)),
            ("REDESIGN", presenter_text(u.student_action, 100)),
            ("CRITIQUE", _pick(ped, 0, "Peer challenges the revised decision")),
        ]
    if u.number == 18:
        vals = ped or []
        return "argument", [
            ("CLAIM", _pick(vals, 0, u.takeaway)),
            ("EVIDENCE", presenter_text(u.evidence or _pick(vals, 1, "Observed evidence"), 100)),
            ("WARRANT", _pick(vals, 2, "Explain why evidence supports the claim")),
            ("COUNTER-EVIDENCE", _pick(vals, 3, "State what would weaken the claim")),
            ("UNCERTAINTY", _pick(vals, 4, "Keep the residual bound visible")),
        ]
    if u.number == 19:
        return "rubric", [(presenter_text(c.criterion, 45), presenter_text(c.ready, 90)) for c in bp.rubric_criteria[:6]]
    if u.number == 20:
        # Unit 20's pedagogy is the four verdict options followed by the
        # residual-uncertainty instruction. Taking ped[0] put the APPROVE line
        # under a RESIDUAL UNCERTAINTY heading, so the slide contradicted its
        # own labels. Each side now carries the content its label promises.
        options = [x for x in ped if _names_a_verdict(x)]
        residual = next((x for x in ped if not _names_a_verdict(x)), "")
        return "verdict", [
            ("TOP CLAIM", presenter_text(u.takeaway, 165)),
            ("THE OPTIONS", presenter_text(" · ".join(options) if options
                                           else (u.evidence or "Weigh the decision against the evidence"), 300)),
            ("RESIDUAL UNCERTAINTY", presenter_text(residual or "State what remains unknown", 220)),
            ("VERDICT", "APPROVE | CONDITIONAL | REDESIGN | REJECT"),
        ]
    vals = core[:5] or ped[:5] or [u.takeaway]
    return "takeaways", _bullet_pairs(vals)


def _sanitize_noncore(text: str) -> str:
    """Keep learner-facing synthetic activities qualitative unless a source owns the number."""
    t = re.sub(r"\b\d+(?:\.\d+)?\s*%", "a bounded amount", str(text or ""), flags=re.I)
    t = re.sub(r"\b\d+(?:\.\d+)?\s*percent\b", "a bounded amount", t, flags=re.I)
    t = re.sub(r"\s+", " ", t).strip()
    return t


def _dedupe_repeated_phrase(title: str) -> str:
    """Collapse a heading that repeats its own opening phrase.

    Slide headers and deck titles often both survive extraction, producing
    "Chapter 2 - Software Processes Chapter 2 Software Processes 1". Cutting at
    the first recurrence of the opening phrase keeps the heading readable
    whatever source produced the repetition.
    """
    words = str(title or "").split()
    if len(words) < 6:
        return str(title or "").strip()
    norm = [re.sub(r"[^a-z0-9]", "", w.lower()) for w in words]
    for size in range(min(6, len(words) // 2), 1, -1):
        head = norm[:size]
        for start in range(size, len(norm) - size + 1):
            if norm[start:start + size] == head:
                return " ".join(words[:start]).rstrip(" -–—,;:/")
    return str(title or "").strip()


def _trim_on_word_boundary(text: str, limit: int) -> str:
    """Never end a classroom heading mid-word."""
    text = str(text or "").strip()
    if len(text) <= limit:
        return text
    head = text[:limit]
    cut = head.rfind(" ")
    return (head[:cut] if cut >= limit * 0.6 else head).rstrip(" -–—,;:/")


def _display_title(bp: Blueprint, u: LectureUnit) -> str:
    """Project internal ISCARB scaffolds into a source-first classroom headline."""
    core = " ".join(u.core_content).lower()
    if u.number == 13 and "process improvement" in core:
        return "Process Improvement: Measure, Analyze, Change"
    if u.number == 14 and "incremental development" in core:
        return "Incremental Development: Visibility, Structure, and Refactoring"
    if u.number == 15 and ("maturity" in core or "capability maturity" in core):
        return "Process Maturity: From Initial to Optimising"
    if u.number == 19:
        return "What You Should Now Be Able to Do"
    if u.number == 20:
        return "Take-home Decision"
    title = str(u.title or "").strip()
    substitutions = (
        (r"^Saudi Context:\s*", ""),
        (r"^Accountability:\s*", ""),
        (r"^Trend\s*&\s*Future:\s*", ""),
        (r"^Practitioner Well(?:being|-being):\s*", ""),
        (r"^Critical AI Literacy\s*&\s*", ""),
        (r"^Portfolio Challenge:\s*", ""),
        (r"^Constraint Mutation\s*&\s*Peer Critique:\s*", ""),
        (r"^Evidence Policy:\s*", "Defending a Process Decision: "),
        (r"^Mechanism Deep Dive:\s*", ""),
        (r"^Authentic Implementation:\s*", ""),
        (r"^First-Principles Prediction Gate:\s*", "Prediction: "),
        (r"^Measurement\s*&\s*Falsification:\s*", ""),
    )
    for pattern, repl in substitutions:
        title = re.sub(pattern, repl, title, flags=re.I)
    if u.number == 8 and "trade" in title.lower():
        title = "Choosing a Process Model: Alternatives and Trade-offs"
    if u.number == 10:
        title = "Senior Design Review: Known, Unknown, What We Monitor"
    if u.number == 16:
        title = re.sub(r"^Complete\s+", "", title, flags=re.I)
    if u.number == 17 and title.lower().startswith("adapting"):
        title = "Adapting the Process When Constraints Change"
    title = _dedupe_repeated_phrase(title or str(u.title or ""))
    return _trim_on_word_boundary(title, 76)


def _display_question(u: LectureUnit) -> str:
    core = " ".join(u.core_content).lower()
    if u.number == 13 and "process improvement" in core:
        return "How do measurement, analysis, and change work together to improve a software process?"
    if u.number == 14 and "incremental development" in core:
        return "What does incremental development trade for adaptability, and when does refactoring become necessary?"
    if u.number == 15 and "maturity" in core:
        return "What do the SEI maturity levels reveal about process control and continuous improvement?"
    if u.number == 19:
        return "Which engineering capabilities should your final artifact demonstrate?"
    if u.number == 20:
        return "What evidence supports approval, redesign, or rejection of the process decision?"
    return presenter_text(_sanitize_noncore(u.engineering_question), 132)


def _compact(text: str, limit: int = 178) -> str:
    """Word-boundary compaction that preserves comma-separated technical lists."""
    t = re.sub(r"\s+", " ", str(text or "")).replace("…", "").replace("...", "").strip()
    if len(t) <= limit:
        return t
    words, out = t.split(), []
    for word in words:
        trial = " ".join(out + [word])
        if len(trial) > limit:
            break
        out.append(word)
    return " ".join(out).rstrip(" ,;:-")


# A checkpoint label and its scaffolding are separated by either punctuation.
_ITEM_LABEL_SPLIT = re.compile(r"\s*[:\u2014\u2013]\s+")

# The archived CIMT lectures use these glyphs as a deliberate two-level bullet
# hierarchy. Finding one mid-sentence therefore does not mean the text is dirty:
# it means a bulleted list was flattened into one line during extraction. The
# structure is recoverable, so split it back out instead of deleting the evidence.
_FLATTENED_BULLET = re.compile(r"[\u25a0\u25aa\u2751\u2752\u274f\u2022\u00b7\u25c6\u25cf\u25b6\uf0b2\uf0a7]+\s*")


# The four closing options all open by naming themselves, which is what
# separates them from the residual-uncertainty instruction beside them.
_VERDICT_WORDS = ("APPROVE", "CONDITIONALLY APPROVE", "REDESIGN", "REJECT")


def _names_a_verdict(line: str) -> bool:
    return str(line or "").strip().upper().startswith(_VERDICT_WORDS)


def _unflatten_bullets(raw: str) -> list[str]:
    """Recover the separate points from a bullet list flattened into one line."""
    parts = [x.strip(" ;,·-") for x in _FLATTENED_BULLET.split(str(raw or ""))]
    parts = [x for x in parts if len(x.split()) >= 3]
    return parts if len(parts) > 1 else [str(raw or "").strip()]


def _split_label(raw: str, fallback: str) -> tuple[str, str]:
    """Use a short leading phrase as the item's heading when the line has one."""
    parts = _ITEM_LABEL_SPLIT.split(raw, maxsplit=1)
    if len(parts) == 2 and parts[1].strip() and 1 <= len(parts[0].split()) <= 7:
        return presenter_text(parts[0], 36).upper(), parts[1].strip()
    return fallback, raw


def _source_first_items(u: LectureUnit, limit: int = 6) -> list[tuple[str, str]]:
    """Source content first, then ISCARB scaffolding, up to the slide's capacity.

    Pedagogy used to appear only when a unit had no source content at all, so a
    unit carrying one source line and four scaffolding steps rendered as a single
    large box and the scaffolding never reached the learner. Source stays first
    and dominant; scaffolding fills the remaining slots rather than being dropped.

    Provenance stays visible: a scaffolding item is headed ISCARB STEP unless the
    line names its own step, so a learner never reads a teaching move as a claim
    the source made.
    """
    focal = [presenter_text(x, 32).upper() for x in (u.visual_plan.focal_elements or []) if str(x).strip()]
    items: list[tuple[str, str]] = []

    source_lines: list[str] = []
    for raw in [str(x).strip() for x in u.core_content if str(x).strip()]:
        source_lines.extend(_unflatten_bullets(raw))
    for i, raw in enumerate(source_lines[:limit]):
        fallback = focal[i] if i < len(focal) else f"KEY POINT {i + 1}"
        label, body = _split_label(raw, fallback)
        items.append((label, _compact(body, 178)))

    if len(items) < limit:
        scaffold = [_sanitize_noncore(x) for x in u.pedagogy_content if str(x).strip()]
        for j, raw in enumerate(scaffold[: limit - len(items)]):
            label, body = _split_label(raw, f"ISCARB STEP {j + 1}")
            items.append((label, _compact(body, 178)))

    return items or [("TAKE-HOME", _compact(u.takeaway, 160))]


def _spec(bp: Blueprint, u: LectureUnit) -> tuple[str, list[tuple[str, str]]]:
    """CIMT grammar driven by the source visual/content job, never by Unit number alone."""
    # Stable pedagogical bookends.
    if u.number in {1, 3, 4, 5, 18, 20}:
        return _legacy_spec(bp, u)
    if u.number == 10:
        known = _pick(_core(u), 0, u.takeaway)
        ped = _ped(u, 6)
        labels = ["KNOWN", "UNKNOWN", "DECISION-SENSITIVE", "WHAT WE MONITOR"]
        bodies = [
            known,
            next((x.split(":",1)[1].strip() for x in ped if x.upper().startswith("UNKNOWN:")), "State the bounded unknown"),
            next((x.split(":",1)[1].strip() for x in ped if x.upper().startswith("DECISION-SENSITIVE UNKNOWN:")), "Identify the unknown that could change the decision"),
            next((x.split(":",1)[1].strip() for x in ped if x.upper().startswith("WHAT WE MONITOR:")), presenter_text(u.student_action, 120)),
        ]
        return "chain", list(zip(labels, map(lambda x: presenter_text(_sanitize_noncore(x), 112), bodies)))
    if u.number == 15 and "maturity" in " ".join(u.core_content).lower():
        return "stack", [
            ("INITIAL", "Work is largely uncontrolled and depends on individual practice."),
            ("REPEATABLE", "Basic procedures are defined so successful practice can be repeated."),
            ("DEFINED", "An organization-wide process strategy is defined and used."),
            ("MANAGED", "Quality and process performance are measured and managed."),
            ("OPTIMISING", "Continuous process improvement is an explicit organizational practice."),
        ]
    if u.number == 19:
        items = [(presenter_text(r.criterion, 46), presenter_text(r.ready, 105)) for r in bp.rubric_criteria[:6]]
        return "takeaways", items

    raw_type = str(getattr(u.visual_plan, "visual_type", "") or "").strip().lower()
    if raw_type in {"process", "workflow", "flow", "sequence"}:
        items = _source_first_items(u, 4)
        # Process slides should visibly read left-to-right like the archived CIMT mechanism diagrams.
        return "chain", items[:4]
    if raw_type in {"trade-off", "tradeoff", "comparison", "compare"}:
        vals = _core(u, 3)
        if len(vals) < 2:
            vals += [_sanitize_noncore(x) for x in _ped(u, 3) if x not in vals]
        a = _pick(vals, 0, u.takeaway)
        b = _pick(vals, 1, u.takeaway)
        focal = [x for x in (u.visual_plan.focal_elements or []) if str(x).strip()]
        if u.number == 14 and "incremental development" in " ".join(u.core_content).lower():
            la, lb = "VISIBILITY & CONTROL", "MEASUREMENT & REFACTORING"
            trade = "Adaptability improves; visibility and structure still require deliberate management and refactoring."
        else:
            la = presenter_text(focal[0] if focal else "ALTERNATIVE A", 28).upper()
            lb = presenter_text(focal[1] if len(focal) > 1 else "ALTERNATIVE B", 28).upper()
            trade = _sanitize_noncore(u.takeaway)
        return "compare", [(la, _compact(a, 165)), (lb, _compact(b, 165)), ("ENGINEERING TRADE-OFF", _compact(trade, 135))]
    if raw_type in {"concept-map", "concept map", "concept", "map"}:
        return "takeaways", _source_first_items(u, 6)

    # Forms below were previously unreachable from a real lecture: only four
    # visual types were mapped, so architecture, table, timeline and measurement
    # content all fell through to a unit-number layout. That is why 41% of every
    # teaching span rendered as the same chain. Each form here already draws on
    # all three surfaces; they simply had no route from the content.
    if raw_type in {"architecture", "layered", "stack", "layers", "system-stack"}:
        # The layered stack is the archived CIMT signature: each layer sits on the
        # one below and hides its detail from the one above.
        return "stack", _source_first_items(u, 6)
    if raw_type in {"table", "matrix", "grid"}:
        return "table", _source_first_items(u, 6)
    if raw_type in {"timeline", "evolution", "lifecycle", "phases"}:
        return "timeline", _source_first_items(u, 3)
    if raw_type in {"measurement", "metric", "equation", "curve", "cost", "trend"}:
        # Cost/benefit and measurement content reads as a curve with annotations,
        # the way the archived cost/dependability slide does.
        return "curve", _source_first_items(u, 4)
    if raw_type in {"failure", "fault", "causal", "chain", "propagation"}:
        return "chain", _source_first_items(u, 4)

    # Archived/synthetic fixtures use historical visual-type tokens; preserve their
    # layout diversity without leaking that unit-number mapping into real lectures.
    return _legacy_spec(bp, u)



def _readiness_for_unit(bp: Blueprint, u: LectureUnit) -> str:
    refs: list[str] = []
    for row in bp.readiness_alignment:
        if u.number in row.evidence_units:
            refs.extend(row.slo_refs[:2])
    refs = list(dict.fromkeys(refs))
    return "READINESS EVIDENCE · " + " / ".join(refs) if refs else ""


def _try_text(u: LectureUnit) -> str:
    return presenter_text(_sanitize_noncore(u.student_action or u.engineering_question), 96)


def _check_text(u: LectureUnit) -> str:
    return presenter_text(_sanitize_noncore(u.takeaway), 86)


# ---------------------------------------------------------------------------
# Source-visual cropping
# ---------------------------------------------------------------------------

def _source_body_path(plan) -> Path | None:
    """Crop source-slide furniture so the teaching figure, not a duplicated deck, is reused."""
    if not getattr(plan, "asset", None):
        return None
    path = local_asset(plan.asset)
    if not path or not path.exists():
        return None
    try:
        out = path.with_name(path.stem + "__cimt_body.png")
        if out.exists() and out.stat().st_mtime_ns >= path.stat().st_mtime_ns and out.stat().st_size > 1000:
            return out
        with Image.open(path) as raw:
            im = raw.convert("RGB")
            w, h = im.size
            # Archived CIMT pages place title furniture in the upper ~17% and a
            # thin course/footer strip in the lower ~7%.  The crop also removes
            # narrow side furniture while preserving diagrams and tables.
            # Preserve edge-to-edge source labels; crop title/footer only vertically.
            box = (0, int(h * 0.17), w, int(h * 0.925))
            crop = im.crop(box)
            crop.save(out, "PNG", optimize=True)
        return out
    except Exception:
        return path


def _data_uri(path: Path | None) -> str:
    if not path or not path.exists():
        return ""
    mime = mimetypes.guess_type(path.name)[0] or "image/png"
    return f"data:{mime};base64,{base64.b64encode(path.read_bytes()).decode('ascii')}"


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def _ppt_text(slide, x, y, w, h, text, size=16, color=INK, bold=False, font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = str(text or ""); p.alignment = align
    p.font.name = font; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    return shape


def _ppt_rect(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(1)
    return shape


def _ppt_connector(slide, x1, y1, x2, y2, color=GREEN_2, width=2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color; line.line.width = Pt(width)
    return line


def _ppt_frame(slide, bp: Blueprint, u: LectureUnit, *, show_decision: bool = False):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
    # CIMT corner rule: an L, not a full-width decorative dashboard bar.
    _ppt_rect(slide, .44, .22, 11.92, .018, GOLD)
    _ppt_rect(slide, .44, .22, .018, .35, GOLD)
    _ppt_text(slide, .54, .34, 10.55, .66, _display_title(bp, u), 30, GREEN, False, "Georgia")
    _ppt_text(slide, .62, 1.02, 10.90, .46, _display_question(u), 12.8, INK, False, "Aptos")

    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.65), Inches(.31), Inches(.23), Inches(.23))
    ring.fill.background(); ring.line.color.rgb = GOLD; ring.line.width = Pt(2)
    _ppt_text(slide, 11.92, .27, .88, .20, "CIMT", 8.2, GREEN, True)
    _ppt_text(slide, 11.92, .46, .88, .15, "ISCARB", 5.8, MUTED, True)

    _ppt_rect(slide, .54, 7.06, 11.72, .014, GOLD)
    _ppt_text(slide, .55, 7.12, 5.60, .18, f"{_subject(bp)} · source-grounded presenter", 6.4, MUTED)
    _ppt_text(slide, 10.85, 7.11, 1.40, .18, f"{u.number:02d}/20", 6.4, MUTED, True, align=PP_ALIGN.RIGHT)
    if show_decision:
        _ppt_rect(slide, .62, 6.43, 11.58, .02, GOLD)
        _ppt_text(slide, .65, 6.52, .62, .22, "TRY", 8.8, RED, True)
        _ppt_text(slide, 1.18, 6.48, 5.12, .36, _try_text(u), 9.5, INK, True)
        _ppt_text(slide, 6.48, 6.52, .72, .22, "CHECK", 8.3, GREEN, True)
        _ppt_text(slide, 7.18, 6.48, 4.20, .36, _check_text(u), 9.1, INK, False)
        ready = _readiness_for_unit(bp, u)
        if ready:
            _ppt_text(slide, 9.25, 6.88, 2.95, .18, ready, 6.2, GREEN, True, align=PP_ALIGN.RIGHT)


def _ppt_title_slide(slide, bp: Blueprint, u: LectureUnit, items):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
    _ppt_rect(slide, .62, .72, 11.90, .018, GOLD)
    _ppt_rect(slide, .62, .72, .018, .55, GOLD)
    _ppt_text(slide, 1.02, 1.28, 11.10, 1.45, presenter_text(bp.lecture_title, 92), 36, GREEN, False, "Georgia", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _ppt_text(slide, 1.70, 2.78, 9.75, .56, "ISCARB Faculty Studio · CIMT-native lecture", 14, INK, False, "Aptos", PP_ALIGN.CENTER)
    _ppt_rect(slide, 3.90, 3.60, 5.55, .018, GOLD)
    _ppt_text(slide, 1.35, 4.02, 10.65, .85, presenter_text(items[0][1], 175), 19, RED, False, "Georgia", PP_ALIGN.CENTER)
    _ppt_text(slide, 2.20, 5.15, 8.90, .48, "20 teaching units · 90 live minutes · source fidelity retained", 10.5, MUTED, True, align=PP_ALIGN.CENTER)
    _ppt_rect(slide, .62, 6.90, 11.90, .018, GOLD)
    _ppt_text(slide, .68, 7.00, 6.5, .22, _source_label(u), 6.3, MUTED)
    _ppt_text(slide, 10.9, 7.00, 1.5, .22, "01/20", 6.3, MUTED, True, align=PP_ALIGN.RIGHT)


# The PPTX equivalent of BAND_TOP/BAND_BOTTOM. The question line clears at
# about 1.55in and the TRY rule is drawn at 6.43in.
PPT_BAND_TOP = 1.55
PPT_BAND_BOTTOM = 6.30


def _ppt_bullets(slide, items, *, x=.78, y=1.55, w=11.5, max_items=6, body_size=19.0):
    items = items[:max_items]
    if not items:
        return
    row_h = min(.82, 4.85 / max(1, len(items)))
    for i, (title, body) in enumerate(items):
        yy = y + i * row_h
        # Was a single .08in square per row and nothing else, which is why the
        # bullet units were the emptiest slides in the deck.
        if i % 2 == 0:
            _ppt_rect(slide, x - .10, yy - .06, w + .20, row_h - .06, PALE_GREEN)
        _ppt_rect(slide, x, yy - .06, .05, row_h - .06, GOLD)
        _ppt_text(slide, x + .12, yy + .02, .34, .22, f"{i + 1:02d}", 8.0, MUTED, True)
        _ppt_text(slide, x + .50, yy, 1.85, .34, presenter_text(title, 32), 11.0, RED if i == 0 else GREEN, True)
        _ppt_text(slide, x + 2.18, yy - .02, w - 2.2, .55, presenter_text(body, 138), body_size, INK, False, "Aptos")


def _ppt_quote(slide, items):
    quote = items[0][1]
    _ppt_text(slide, 1.15, 1.75, 11.0, 2.25, f'“{quote}”', 27, RED, False, "Georgia", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _ppt_rect(slide, 4.55, 4.22, 4.20, .018, GOLD)
    for i, (_, body) in enumerate(items[1:4]):
        x = 1.0 + i * 4.08
        _ppt_text(slide, x, 4.58, .55, .32, f"0{i+1}", 12, GOLD, True, "Georgia")
        _ppt_text(slide, x + .50, 4.48, 3.20, .88, body, 13.5, INK, True)


def _ppt_orbit(slide, items):
    cx, cy = 6.55, 3.72
    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - .78), Inches(cy - .55), Inches(1.56), Inches(1.10))
    hub.fill.solid(); hub.fill.fore_color.rgb = GREEN_2; hub.line.fill.background()
    _ppt_text(slide, cx - .66, cy - .22, 1.32, .42, "DECISION", 12, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    positions = [(1.05,1.70),(1.05,3.20),(1.05,4.70),(9.20,1.70),(9.20,3.20),(9.20,4.70)]
    for i, (title, body) in enumerate(items[:6]):
        x, y = positions[i]
        side_right = x > cx
        x2 = x if side_right else x + 3.0
        _ppt_connector(slide, cx + (.8 if side_right else -.8), cy, x2, y + .35, GOLD, 1.4)
        _ppt_text(slide, x, y, 3.05, .28, title, 10.2, RED if i in {1,4} else GREEN, True)
        _ppt_text(slide, x, y + .30, 3.05, .62, presenter_text(body, 92), 12.3, INK)


def _ppt_ladder(slide, items):
    base_x, base_y = 1.05, 5.55
    step_w, step_h = 2.55, .72
    for i, (title, body) in enumerate(items[:4]):
        x = base_x + i * 2.67
        y = base_y - i * .86
        _ppt_rect(slide, x, y, step_w, step_h, PALE_GOLD if i % 2 == 0 else PALE_GREEN, GOLD)
        _ppt_text(slide, x + .10, y + .08, step_w - .20, .20, title, 8.5, RED if i == 0 else GREEN, True)
        _ppt_text(slide, x + .10, y + .29, step_w - .20, .35, presenter_text(body, 78), 10.8, INK)
        if i < 3:
            _ppt_connector(slide, x + step_w, y + .36, x + 2.67, y - .50, GOLD, 1.4)
    _ppt_text(slide, 1.10, 1.50, 3.0, .35, "FROM FIRST PRINCIPLES", 10.5, GREEN, True)
    _ppt_text(slide, 1.10, 1.90, 4.2, .55, "Predict first. Name the mechanism only after the constraint becomes visible.", 16, INK)


def _ppt_curve(slide, items):
    x0, y0, x1, y1 = 1.25, 5.62, 7.45, 1.72
    _ppt_connector(slide, x0, y0, x1, y0, INK, 1.2)
    _ppt_connector(slide, x0, y0, x0, y1, INK, 1.2)
    pts = [(1.55,5.42),(2.35,5.28),(3.20,5.05),(4.05,4.70),(4.85,4.15),(5.55,3.35),(6.10,2.45),(6.50,1.82)]
    for a, b in zip(pts, pts[1:]):
        _ppt_connector(slide, a[0], a[1], b[0], b[1], GREEN_2, 2.2)
    _ppt_text(slide, 3.15, 5.78, 2.6, .32, "assurance / dependability →", 8.5, MUTED, True, align=PP_ALIGN.CENTER)
    _ppt_text(slide, .70, 3.15, .42, 1.3, "COST", 8.5, MUTED, True, align=PP_ALIGN.CENTER)
    _ppt_text(slide, 8.05, 1.78, 4.10, .32, items[1][0], 10.3, RED, True)
    _ppt_text(slide, 8.05, 2.16, 4.10, 1.18, presenter_text(items[1][1], 115), 14.2, INK)
    _ppt_text(slide, 8.05, 3.62, 4.10, .32, items[2][0], 10.3, GREEN, True)
    _ppt_text(slide, 8.05, 4.00, 4.10, 1.18, presenter_text(items[2][1], 115), 14.2, INK)


def _ppt_stack(slide, items):
    shown = items[:6]
    if not shown:
        return
    x, w = 2.05, 9.25
    # Matches the PDF stack: rows divide the whole band instead of stepping a
    # fixed .74in, so a five-item stack no longer stops a quarter of the way
    # up the slide, and the body type grows with the row.
    pitch = (PPT_BAND_BOTTOM - PPT_BAND_TOP) / len(shown)
    h = pitch - .10
    body_size = 12.5 if pitch < .78 else (15.0 if pitch < .95 else 17.5)
    for i, (title, body) in enumerate(shown):
        yy = PPT_BAND_TOP + i * pitch
        _ppt_rect(slide, x, yy, w, h, PALE_GREEN if i % 2 else WHITE, GREEN_2)
        _ppt_text(slide, x + .12, yy + h / 2 - .12, 1.35, .24, title, 8.8, RED if i in {0,5} else GREEN, True)
        _ppt_text(slide, x + 1.55, yy + .06, w - 1.70, h - .12, presenter_text(body, 112), body_size, INK, False, "Aptos", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        _ppt_text(slide, x - .38, yy + h / 2 - .12, .30, .24, f"L{len(shown) - i}", 8.0, MUTED, True, align=PP_ALIGN.RIGHT)
        if i:
            _ppt_connector(slide, x + w + .18, yy - pitch + h, x + w + .18, yy, GOLD, 1.4)
    mid = PPT_BAND_TOP + (PPT_BAND_BOTTOM - PPT_BAND_TOP) / 2
    _ppt_text(slide, x + w + .30, mid - .22, .95, .44, "CHANGE\nRIPPLES", 8.0, GOLD, True, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)


def _ppt_compare(slide, items):
    _ppt_text(slide, .95, 1.55, 5.0, .32, items[0][0], 11, GREEN, True)
    _ppt_text(slide, .95, 2.02, 5.05, 2.55, presenter_text(items[0][1], 150), 19, INK, False, "Georgia", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _ppt_rect(slide, 6.48, 1.52, .018, 4.40, GOLD)
    _ppt_text(slide, 7.05, 1.55, 5.0, .32, items[1][0], 11, RED, True)
    _ppt_text(slide, 7.05, 2.02, 5.05, 2.55, presenter_text(items[1][1], 150), 19, INK, False, "Georgia", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _ppt_rect(slide, 3.40, 5.22, 6.45, .56, PALE_GOLD, GOLD)
    _ppt_text(slide, 3.58, 5.35, 1.15, .22, "TRADE-OFF", 8.6, RED, True)
    _ppt_text(slide, 4.82, 5.28, 4.82, .30, presenter_text(items[2][1], 105), 10.5, INK, True, align=PP_ALIGN.CENTER)


def _ppt_tree(slide, items):
    _ppt_rect(slide, 4.72, 1.55, 3.90, .72, PALE_GOLD, GOLD)
    _ppt_text(slide, 4.92, 1.70, 3.50, .24, items[0][0], 9, RED, True, align=PP_ALIGN.CENTER)
    _ppt_text(slide, 4.92, 1.96, 3.50, .22, presenter_text(items[0][1], 60), 10.2, INK, True, align=PP_ALIGN.CENTER)
    positions = [(1.0,3.40),(4.0,3.40),(7.0,3.40),(10.0,3.40)]
    for i, (title, body) in enumerate(items[1:5]):
        x, y = positions[i]
        _ppt_connector(slide, 6.67, 2.27, x + 1.10, y, GREEN_2, 1.3)
        _ppt_text(slide, x, y + .10, 2.2, .25, title, 8.5, GREEN, True, align=PP_ALIGN.CENTER)
        _ppt_text(slide, x, y + .48, 2.2, 1.32, presenter_text(body, 88), 12.2, INK, False, "Aptos", PP_ALIGN.CENTER)
        _ppt_rect(slide, x + .18, y + 1.94, 1.84, .018, GOLD)


def _ppt_table(slide, items):
    x, y, w = .95, 1.50, 11.45
    _ppt_rect(slide, x, y, w, .50, GOLD)
    _ppt_text(slide, x + .16, y + .12, 2.2, .24, "CHARACTERISTIC", 8.5, WHITE, True)
    _ppt_text(slide, x + 2.55, y + .12, 8.45, .24, "WHAT IT MEANS IN THE ENGINEERING DECISION", 8.5, WHITE, True)
    row_h = .78
    for i, (title, body) in enumerate(items[:5]):
        yy = y + .50 + i * row_h
        _ppt_rect(slide, x, yy, w, row_h, PALE_GOLD if i % 2 == 0 else WHITE, LINE)
        _ppt_text(slide, x + .16, yy + .16, 2.15, .32, presenter_text(title, 30), 10.8, GREEN, True)
        _ppt_text(slide, x + 2.55, yy + .10, 8.45, .50, presenter_text(body, 130), 12.2, INK)


def _ppt_context(slide, items):
    _ppt_rect(slide, .95, 1.55, 4.15, 3.90, PALE_RED, RED)
    _ppt_text(slide, 1.20, 1.82, 3.65, .55, items[0][0], 11.5, RED, True, align=PP_ALIGN.CENTER)
    _ppt_text(slide, 1.20, 2.58, 3.65, 2.20, items[0][1], 18.5, INK, False, "Georgia", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _ppt_connector(slide, 5.32, 3.48, 6.18, 3.48, GOLD, 2)
    _ppt_text(slide, 6.38, 1.75, 5.35, .32, items[1][0], 10.5, GREEN, True)
    _ppt_text(slide, 6.38, 2.18, 5.35, 1.14, items[1][1], 15.5, INK)
    _ppt_text(slide, 6.38, 3.72, 5.35, .32, items[2][0], 10.5, RED, True)
    _ppt_text(slide, 6.38, 4.15, 5.35, 1.14, items[2][1], 15.5, INK)


def _ppt_chain(slide, items):
    n = len(items); left = .82; gap = .18; total = 11.65; w = (total - gap * (n - 1)) / n
    for i, (title, body) in enumerate(items):
        x = left + i * (w + gap)
        _ppt_text(slide, x, 2.10, w, .30, title, 9.2, RED if i in {1,3} else GREEN, True, align=PP_ALIGN.CENTER)
        _ppt_rect(slide, x, 2.62, w, 2.55, WHITE if i % 2 else PALE_GREEN, GREEN_2)
        _ppt_text(slide, x + .14, 2.86, w - .28, 1.95, body, 14.2, INK, False, "Aptos", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            _ppt_connector(slide, x + w, 3.88, x + w + gap, 3.88, GOLD, 2)


def _ppt_timeline(slide, items):
    _ppt_connector(slide, 1.35, 3.35, 11.95, 3.35, GOLD, 2)
    xs = [2.0, 6.65, 11.25]
    for i, (title, body) in enumerate(items[:3]):
        x = xs[i]
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - .12), Inches(3.22), Inches(.24), Inches(.24))
        dot.fill.solid(); dot.fill.fore_color.rgb = RED if i == 2 else GREEN_2; dot.line.fill.background()
        _ppt_text(slide, x - 1.40, 2.25, 2.80, .28, title, 10.2, RED if i == 2 else GREEN, True, align=PP_ALIGN.CENTER)
        _ppt_text(slide, x - 1.55, 3.82, 3.10, 1.22, body, 14.0, INK, False, "Aptos", PP_ALIGN.CENTER)


def _ppt_burden(slide, items):
    _ppt_text(slide, .98, 1.60, 4.85, .32, items[0][0], 10.5, GREEN, True)
    _ppt_text(slide, .98, 2.03, 4.85, 1.28, items[0][1], 17, INK, False, "Georgia", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _ppt_text(slide, .98, 3.68, 4.85, .32, items[1][0], 10.5, RED, True)
    _ppt_text(slide, .98, 4.10, 4.85, 1.28, items[1][1], 16, INK, False, "Georgia", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _ppt_connector(slide, 6.15, 3.38, 7.05, 3.38, GOLD, 2.2)
    _ppt_text(slide, 7.35, 1.68, 4.45, .32, items[2][0], 10.5, GREEN, True)
    _ppt_text(slide, 7.35, 2.10, 4.45, 1.30, items[2][1], 16.5, INK)
    _ppt_text(slide, 7.35, 3.78, 4.45, .32, items[3][0], 10.5, RED, True)
    _ppt_text(slide, 7.35, 4.20, 4.45, 1.25, items[3][1], 16.5, INK)


def _ppt_ai(slide, items):
    cols = [(.85, PALE_GREEN, GREEN_2), (4.47, PALE_GOLD, GOLD), (8.09, PALE_RED, RED)]
    data = [items[0], items[1], items[3]]
    for i, ((x, fill, line), (title, body)) in enumerate(zip(cols, data)):
        _ppt_rect(slide, x, 1.85, 3.35, 3.68, fill, line)
        _ppt_text(slide, x + .18, 2.12, 2.99, .36, title, 10.5, line, True, align=PP_ALIGN.CENTER)
        _ppt_text(slide, x + .25, 2.78, 2.85, 1.75, body, 16, INK, False, "Georgia", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _ppt_text(slide, 4.72, 5.72, 3.0, .24, "TEST", 8.5, RED, True, align=PP_ALIGN.CENTER)
    _ppt_text(slide, 3.18, 6.02, 6.1, .32, presenter_text(items[2][1], 105), 10.2, INK, True, align=PP_ALIGN.CENTER)


def _ppt_portfolio(slide, items):
    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.45), Inches(2.75), Inches(2.35), Inches(1.55))
    hub.fill.solid(); hub.fill.fore_color.rgb = PALE_GOLD; hub.line.color.rgb = GOLD; hub.line.width = Pt(1.2)
    _ppt_text(slide, 5.65, 3.06, 1.95, .52, "ENGINEERING\nMISSION", 12.5, GREEN, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    pos = [(1.0,1.55),(4.0,1.35),(8.95,1.55),(1.0,4.55),(4.05,5.0),(8.95,4.55)]
    for i, (title, body) in enumerate(items[:6]):
        x, y = pos[i]
        _ppt_connector(slide, 6.62, 3.48, x + 1.3, y + .55, GOLD, 1.2)
        _ppt_text(slide, x, y, 2.65, .25, title, 8.8, RED if i in {0,3} else GREEN, True, align=PP_ALIGN.CENTER)
        _ppt_text(slide, x, y + .35, 2.65, .86, body, 11.8, INK, False, "Aptos", PP_ALIGN.CENTER)


def _ppt_mutation(slide, items):
    _ppt_chain(slide, items)
    _ppt_text(slide, 4.15, 5.62, 5.05, .32, "Change one constraint. Re-run the decision.", 12.5, RED, True, "Georgia", PP_ALIGN.CENTER)


def _ppt_argument(slide, items):
    y = 1.50
    widths = [10.4, 9.5, 8.6, 7.7, 6.8]
    for i, (title, body) in enumerate(items[:5]):
        x = (13.333 - widths[i]) / 2
        _ppt_text(slide, x, y, 1.75, .25, title, 8.7, RED if i in {0,3,4} else GREEN, True)
        _ppt_rect(slide, x + 1.70, y - .04, widths[i] - 1.70, .56, PALE_GOLD if i in {0,3} else (PALE_GREEN if i in {1,2} else PALE_RED), GOLD if i in {0,3} else GREEN_2)
        _ppt_text(slide, x + 1.88, y + .08, widths[i] - 2.05, .30, presenter_text(body, 105), 11.8, INK, True, align=PP_ALIGN.CENTER)
        y += .92


def _ppt_rubric(slide, items):
    cols = 3; x0 = .90; y0 = 1.52; gapx = .22; gapy = .25; w = 3.68; h = 1.82
    for i, (title, body) in enumerate(items[:6]):
        r, c = divmod(i, cols); x = x0 + c * (w + gapx); y = y0 + r * (h + gapy)
        _ppt_text(slide, x, y, w, .30, title, 9.4, GREEN, True, align=PP_ALIGN.CENTER)
        _ppt_rect(slide, x, y + .45, w, 1.15, WHITE if i % 2 else PALE_GREEN, GREEN_2)
        _ppt_text(slide, x + .18, y + .68, w - .36, .72, body, 12.8, INK, False, "Aptos", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def _ppt_verdict(slide, items):
    _ppt_text(slide, 1.08, 1.58, 11.15, .28, items[0][0], 10.5, GREEN, True, align=PP_ALIGN.CENTER)
    _ppt_text(slide, 1.08, 2.02, 11.15, 1.02, items[0][1], 21.5, INK, False, "Georgia", PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    _ppt_rect(slide, 1.12, 3.55, 5.25, 1.30, PALE_GREEN, GREEN_2)
    _ppt_text(slide, 1.35, 3.76, 1.20, .25, items[1][0], 8.8, GREEN, True)
    _ppt_text(slide, 2.48, 3.68, 3.55, .68, items[1][1], 12.3, INK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _ppt_rect(slide, 6.95, 3.55, 5.25, 1.30, PALE_RED, RED)
    _ppt_text(slide, 7.18, 3.76, 1.75, .25, items[2][0], 8.8, RED, True)
    _ppt_text(slide, 8.72, 3.68, 3.10, .68, items[2][1], 12.3, INK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _ppt_rect(slide, 2.65, 5.40, 8.05, .62, PALE_GOLD, GOLD)
    _ppt_text(slide, 2.82, 5.57, 7.70, .26, items[3][1], 11, GREEN, True, align=PP_ALIGN.CENTER)


def _ppt_source(slide, bp: Blueprint, u: LectureUnit, plan) -> bool:
    path = _source_body_path(plan)
    if not path or not path.exists():
        return False
    _ppt_frame(slide, bp, u, show_decision=u.number in set(range(5, 19)))
    try:
        with Image.open(path) as im:
            iw, ih = im.size
        box_x, box_y, box_w, box_h = .75, 1.33, 11.80, 4.95
        scale = min(box_w / iw, box_h / ih)
        w, h = iw * scale, ih * scale
        slide.shapes.add_picture(str(path), Inches(box_x + (box_w - w) / 2), Inches(box_y + (box_h - h) / 2), width=Inches(w), height=Inches(h))
        _ppt_text(slide, 9.30, 6.28, 3.0, .20, f"ADAPTED VISUAL · P1 {plan.source_slide}", 6.4, GREEN, True, align=PP_ALIGN.RIGHT)
        return True
    except Exception:
        return False


def _ppt_redraw(slide, bp: Blueprint, u: LectureUnit):
    kind, items = _spec(bp, u)
    if kind == "title":
        _ppt_title_slide(slide, bp, u, items); return
    _ppt_frame(slide, bp, u, show_decision=u.number in set(range(5, 19)))
    if kind == "quote": _ppt_quote(slide, items)
    elif kind == "takeaways": _ppt_bullets(slide, items)
    elif kind == "orbit": _ppt_orbit(slide, items)
    elif kind == "ladder": _ppt_ladder(slide, items)
    elif kind == "curve": _ppt_curve(slide, items)
    elif kind == "stack": _ppt_stack(slide, items)
    elif kind == "compare": _ppt_compare(slide, items)
    elif kind == "tree": _ppt_tree(slide, items)
    elif kind == "table": _ppt_table(slide, items)
    elif kind == "context": _ppt_context(slide, items)
    elif kind == "chain": _ppt_chain(slide, items)
    elif kind == "timeline": _ppt_timeline(slide, items)
    elif kind == "burden": _ppt_burden(slide, items)
    elif kind == "ai": _ppt_ai(slide, items)
    elif kind == "portfolio": _ppt_portfolio(slide, items)
    elif kind == "mutation": _ppt_mutation(slide, items)
    elif kind == "argument": _ppt_argument(slide, items)
    elif kind == "rubric": _ppt_rubric(slide, items)
    elif kind == "verdict": _ppt_verdict(slide, items)
    else: _ppt_bullets(slide, items)


def export_cimt_presenter_pptx_v43(bp: Blueprint, out: Path) -> Path:
    out = Path(out)
    prs = Presentation(); prs.slide_width = Inches(PPT_W); prs.slide_height = Inches(PPT_H)
    plans = plans_for_blueprint_v42(bp)
    for u, plan in zip(bp.units, plans):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Preserve a recognisable title page instead of pasting a source visual on unit 1.
        if u.number != 1 and plan.reuse_mode == "USE" and _ppt_source(slide, bp, u, plan):
            continue
        _ppt_redraw(slide, bp, u)
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _r_wrap(c, text, x, y, width, size=12, color=R_INK, bold=False, max_lines=5, align="left", font=None):
    size = _snap(size)
    font_name = font or ("Helvetica-Bold" if bold else "Helvetica")
    words = str(text or "").split(); lines: list[str] = []; line = ""
    for word in words:
        trial = (line + " " + word).strip()
        if c.stringWidth(trial, font_name, size) <= width:
            line = trial
        else:
            if line: lines.append(line)
            line = word
            if len(lines) >= max_lines - 1: break
    if line and len(lines) < max_lines: lines.append(line)
    c.setFont(font_name, size); c.setFillColor(color)
    for i, ln in enumerate(lines):
        yy = y - i * size * 1.28
        if align == "center": c.drawCentredString(x + width / 2, yy, ln)
        elif align == "right": c.drawRightString(x + width, yy, ln)
        else: c.drawString(x, yy, ln)


def _r_frame(c, bp: Blueprint, u: LectureUnit, show_decision=False):
    c.setFillColor(R_WHITE); c.rect(0, 0, PDF_W, PDF_H, fill=1, stroke=0)
    c.setFillColor(R_GOLD); c.rect(32, 516, 865, 1.3, fill=1, stroke=0); c.rect(32, 493, 1.3, 24, fill=1, stroke=0)
    c.setFillColor(R_GREEN); c.setFont("Times-Roman", _snap(25)); c.drawString(40, 482, _display_title(bp, u))
    _r_wrap(c, _display_question(u), 45, 452, 780, 9.6, R_INK, False, 2)
    c.setStrokeColor(R_GOLD); c.setLineWidth(1.6); c.circle(850, 493, 7.5, fill=0, stroke=1)
    c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(7)); c.drawString(865, 495, "CIMT")
    c.setFillColor(R_MUTED); c.setFont("Helvetica-Bold", _snap(4.8)); c.drawString(865, 487, "ISCARB")
    c.setFillColor(R_GOLD); c.rect(40, 31, 850, 1, fill=1, stroke=0)
    c.setFillColor(R_MUTED); c.setFont("Helvetica", _snap(5.8)); c.drawString(40, 20, f"{_subject(bp)} · source-grounded presenter")
    c.drawRightString(890, 20, f"{u.number:02d}/20")
    if show_decision:
        c.setStrokeColor(R_GOLD); c.setLineWidth(.8); c.line(48, 64, 875, 64)
        c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(7)); c.drawString(48, 50, "TRY")
        _r_wrap(c, _try_text(u), 78, 51, 360, 7.5, R_INK, True, 2)
        c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(6.8)); c.drawString(470, 50, "CHECK")
        _r_wrap(c, _check_text(u), 515, 51, 300, 7.2, R_INK, False, 2)
        ready = _readiness_for_unit(bp, u)
        if ready:
            c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(4.8)); c.drawRightString(875, 36, ready)


def _r_title(c, bp: Blueprint, u: LectureUnit, items):
    c.setFillColor(R_WHITE); c.rect(0, 0, PDF_W, PDF_H, fill=1, stroke=0)
    c.setFillColor(R_GOLD); c.rect(42, 480, 850, 1.3, fill=1, stroke=0); c.rect(42, 446, 1.3, 35, fill=1, stroke=0)
    _r_wrap(c, bp.lecture_title, 90, 392, 780, 31, R_GREEN, False, 2, "center", "Times-Roman")
    c.setFillColor(R_INK); c.setFont("Helvetica", _snap(11)); c.drawCentredString(480, 312, "ISCARB Faculty Studio · CIMT-native lecture")
    c.setFillColor(R_GOLD); c.rect(330, 280, 300, 1.2, fill=1, stroke=0)
    _r_wrap(c, items[0][1], 115, 225, 730, 17.5, R_RED, False, 3, "center", "Times-Roman")
    c.setFillColor(R_MUTED); c.setFont("Helvetica-Bold", _snap(8)); c.drawCentredString(480, 110, "20 teaching units · 90 live minutes · source fidelity retained")
    c.setFillColor(R_GOLD); c.rect(42, 50, 850, 1.2, fill=1, stroke=0)
    c.setFillColor(R_MUTED); c.setFont("Helvetica", _snap(5.8)); c.drawString(45, 36, _source_label(u)); c.drawRightString(892, 36, "01/20")


def _r_bullets(c, items, x=65, y=405, width=820, body_size=13.5):
    """Numbered rows separated by rules, not a column of floating text.

    This drew one 6pt square per item and nothing else, so the four units that
    use it were the emptiest pages in the deck. The row band, the index and the
    separator cost nothing and give the eye somewhere to land.
    """
    shown = items[:6]
    row = min(56, 290 / max(1, len(shown)))
    for i, (title, body) in enumerate(shown):
        yy = y - i * row
        top, height = yy - row + 16, row - 8
        if i % 2 == 0:
            c.setFillColor(R_PALE_GREEN); c.rect(x - 8, top, width + 16, height, fill=1, stroke=0)
        c.setFillColor(R_GOLD); c.rect(x, top, 3.5, height, fill=1, stroke=0)
        c.setFillColor(R_MUTED); c.setFont("Helvetica-Bold", _snap(6))
        c.drawString(x + 10, yy + 5, f"{i + 1:02d}")
        c.setFillColor(R_RED if i == 0 else R_GREEN); c.setFont("Helvetica-Bold", _snap(8.3)); c.drawString(x + 15, yy - 7, presenter_text(title, 32))
        _r_wrap(c, body, x + 145, yy - 5, width - 155, body_size, R_INK, False, 2)
        if i < len(shown) - 1:
            c.setStrokeColor(R_LINE); c.setLineWidth(0.6)
            c.line(x - 8, top - 4, x + width + 8, top - 4)


def _r_source(c, bp: Blueprint, u: LectureUnit, plan) -> bool:
    path = _source_body_path(plan)
    if not path or not path.exists():
        return False
    _r_frame(c, bp, u, show_decision=u.number in set(range(5, 19)))
    try:
        img = ImageReader(str(path)); iw, ih = img.getSize(); box = (55, 80, 850, 350)
        scale = min(box[2] / iw, box[3] / ih); dw, dh = iw * scale, ih * scale
        c.drawImage(img, box[0] + (box[2] - dw) / 2, box[1] + (box[3] - dh) / 2, width=dw, height=dh, preserveAspectRatio=True, mask='auto')
        c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(5.7)); c.drawRightString(900, 62, f"ADAPTED VISUAL · P1 {plan.source_slide}")
        return True
    except Exception:
        return False


# The drawable band between the question line and the TRY/CHECK footer. Forms
# used to hard-code their own vertical extents and stop short of it, which is
# what left a third of the slide blank.
# The question line ends near y=445 and the TRY rule is drawn at y=64, so this
# is the real drawable height. Forms previously stopped at y=150 and left 72pt
# - a quarter of the band - blank above the footer.
BAND_TOP = 428.0
BAND_BOTTOM = 82.0


def _arrow(c, x1, y1, x2, y2, color=None, width=1.4, head=6.0):
    """A line that ends in a head.

    Every "diagram" in this deck was a row of rectangles: a chain with no
    arrows between its links, a stack with no direction, a curve with no axis.
    The connective marks are what make a form readable as a diagram rather than
    as boxes of text, so they are drawn, not implied.
    """
    import math

    color = color or R_GREEN_2
    c.saveState()
    c.setStrokeColor(color); c.setLineWidth(width)
    angle = math.atan2(y2 - y1, x2 - x1)
    # Stop the shaft short so the head tip lands exactly on the target point.
    c.line(x1, y1, x2 - head * math.cos(angle) * 0.9, y2 - head * math.sin(angle) * 0.9)
    c.setFillColor(color)
    path = c.beginPath()
    path.moveTo(x2, y2)
    for sign in (1, -1):
        path.lineTo(x2 - head * math.cos(angle - sign * 0.42),
                    y2 - head * math.sin(angle - sign * 0.42))
    path.close()
    c.drawPath(path, fill=1, stroke=0)
    c.restoreState()


def _tick_label(c, x, y, text, size=6.0, color=None, centred=True):
    """A small axis or index label."""
    c.setFillColor(color or R_MUTED)
    c.setFont("Helvetica-Bold", _snap(size))
    (c.drawCentredString if centred else c.drawString)(x, y, text)


def _r_redraw(c, bp: Blueprint, u: LectureUnit):
    kind, items = _spec(bp, u)
    if kind == "title":
        _r_title(c, bp, u, items); return
    _r_frame(c, bp, u, show_decision=u.number in set(range(5, 19)))
    if kind == "quote":
        _r_wrap(c, f'“{items[0][1]}”', 110, 365, 740, 22, R_RED, False, 4, "center", "Times-Roman")
        c.setFillColor(R_GOLD); c.rect(330, 220, 300, 1.2, fill=1, stroke=0)
        for i, (_, body) in enumerate(items[1:4]):
            x = 75 + i * 285
            c.setFillColor(R_GOLD); c.setFont("Helvetica-Bold", _snap(12)); c.drawString(x, 185, f"0{i+1}")
            _r_wrap(c, body, x + 32, 186, 235, 10.5, R_INK, True, 4)
        return
    if kind in {"takeaways", "table"}:
        if kind == "table":
            c.setFillColor(R_GOLD); c.rect(65, 400, 830, 28, fill=1, stroke=0)
            c.setFillColor(R_WHITE); c.setFont("Helvetica-Bold", _snap(7.5)); c.drawString(76, 410, "CHARACTERISTIC"); c.drawString(270, 410, "WHAT IT MEANS IN THE ENGINEERING DECISION")
            rows = items[:5]
            row_h = (400 - BAND_BOTTOM) / max(1, len(rows))
            for i, (title, body) in enumerate(rows):
                yy = 400 - (i + 1) * row_h
                c.setFillColor(R_PALE_GOLD if i % 2 == 0 else R_WHITE); c.setStrokeColor(R_LINE); c.rect(65, yy, 830, row_h, fill=1, stroke=1)
                c.setStrokeColor(R_LINE); c.setLineWidth(0.6); c.line(258, yy, 258, yy + row_h)
                c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(8)); c.drawString(76, yy + row_h/2 - 4, presenter_text(title, 28))
                _r_wrap(c, body, 270, yy + row_h - 14, 605, 11.0, R_INK, False, 4)
        else:
            _r_bullets(c, items)
        return
    if kind == "orbit":
        c.setFillColor(R_GREEN_2); c.circle(480, 270, 48, fill=1, stroke=0); c.setFillColor(R_WHITE); c.setFont("Helvetica-Bold", _snap(9)); c.drawCentredString(480, 267, "DECISION")
        pos = [(75,380),(75,270),(75,160),(650,380),(650,270),(650,160)]
        for i, (title, body) in enumerate(items[:6]):
            x,y = pos[i]
            # Spokes: without them the hub is decoration and the six items are
            # unrelated captions parked either side of it.
            _arrow(c, x + (245 if i < 3 else -10), y + 4, 480 + (-52 if i < 3 else 52), 270 + (30 - (i % 3) * 30),
                   R_LINE if i % 2 else R_GREEN_2, 0.9, 4.5)
            c.setFillColor(R_RED if i in {1,4} else R_GREEN); c.setFont("Helvetica-Bold", _snap(8)); c.drawString(x,y,title)
            _r_wrap(c, body, x, y-17, 235, 9.2, R_INK, False, 4)
        return
    if kind == "ladder":
        for i, (title, body) in enumerate(items[:4]):
            x = 80 + i * 205; y = 135 + i * 54
            c.setFillColor(R_PALE_GOLD if i % 2 == 0 else R_PALE_GREEN); c.setStrokeColor(R_GOLD); c.rect(x,y,185,48,fill=1,stroke=1)
            c.setFillColor(R_RED if i == 0 else R_GREEN); c.setFont("Helvetica-Bold", _snap(7)); c.drawString(x+8,y+31,title)
            _r_wrap(c, body, x+8, y+20, 168, 7.7, R_INK, False, 2)
        _r_wrap(c, "Predict first. Name the mechanism only after the constraint becomes visible.", 80, 400, 360, 13, R_INK, False, 3)
        return
    if kind == "curve":
        # Gridlines and ticks first, so the curve reads against a scale rather
        # than floating on an unlabelled pair of axes.
        c.setStrokeColor(R_LINE); c.setLineWidth(0.5)
        for g in range(1, 5):
            gy = 130 + g * 68
            c.line(95, gy, 560, gy)
        for g in range(1, 6):
            gx = 95 + g * 78
            c.line(gx, 130, gx, 410)
        c.setStrokeColor(R_INK); c.setLineWidth(1); c.line(95,130,560,130); c.line(95,130,95,410)
        for g in range(0, 6):
            gx = 95 + g * 78
            c.setStrokeColor(R_INK); c.line(gx, 126, gx, 130)
        for g in range(0, 5):
            gy = 130 + g * 68
            c.setStrokeColor(R_INK); c.line(91, gy, 95, gy)
        _tick_label(c, 327, 112, "DEPENDABILITY ACHIEVED  \u2192", 6.0)
        c.saveState(); c.translate(78, 270); c.rotate(90)
        _tick_label(c, 0, 0, "COST  \u2192", 6.0); c.restoreState()
        pts=[(110,145),(175,150),(240,165),(305,190),(365,225),(415,275),(455,340),(485,405)]
        c.setStrokeColor(R_GREEN_2); c.setLineWidth(2)
        for a,b in zip(pts,pts[1:]): c.line(a[0],a[1],b[0],b[1])
        # Mark where the curve turns: that knee is the engineering decision.
        c.setFillColor(R_RED); c.circle(415, 275, 3.6, fill=1, stroke=0)
        _arrow(c, 470, 210, 424, 266, R_RED, 1.1, 5.0)
        _tick_label(c, 505, 200, "COST TURNS", 5.8, R_RED)
        c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(8)); c.drawString(610,370,items[1][0]); _r_wrap(c,items[1][1],610,350,280,11,R_INK,False,5)
        c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(8)); c.drawString(610,245,items[2][0]); _r_wrap(c,items[2][1],610,225,280,11,R_INK,False,5)
        return
    if kind == "stack":
        shown = items[:6]
        # The row pitch used to be a fixed 47pt, so a five-item stack left the
        # bottom third of the slide empty while its text stayed at 8.5pt. Rows
        # now divide the whole content band, and the type grows with them.
        top, bottom = BAND_TOP, BAND_BOTTOM
        pitch = (top - bottom) / max(1, len(shown))
        box_h = pitch - 9
        body_size = 9.0 if pitch < 50 else (11.0 if pitch < 62 else 13.5)
        for i,(title,body) in enumerate(shown):
            yy = top - pitch * (i + 1) + (pitch - box_h)
            c.setFillColor(R_PALE_GREEN if i%2 else R_WHITE); c.setStrokeColor(R_GREEN_2); c.rect(175,yy,610,box_h,fill=1,stroke=1)
            c.setFillColor(R_RED if i in {0,5} else R_GREEN); c.setFont("Helvetica-Bold", _snap(7.5)); c.drawString(190,yy+box_h/2-4,title)
            _r_wrap(c,body,285,yy+box_h-10,480,body_size,R_INK,False,3,"center")
            _tick_label(c, 160, yy+box_h/2-4, f"L{len(shown)-i}", 6.0, R_MUTED)
            # A layer that hides detail from the one above it only reads as a
            # layer when the deck draws where the effect travels.
            if i:
                _arrow(c, 800, yy+pitch+2, 800, yy+box_h-2, R_GOLD, 1.2, 5.0)
        if len(shown) > 1:
            mid = (top + bottom) / 2
            _tick_label(c, 828, mid + 5, "CHANGE", 5.8, R_GOLD)
            _tick_label(c, 828, mid - 5, "RIPPLES", 5.8, R_GOLD)
        return
    if kind == "compare":
        c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(8.5)); c.drawCentredString(255,400,items[0][0]); _r_wrap(c,items[0][1],75,350,360,15,R_INK,False,6,"center","Times-Roman")
        c.setFillColor(R_GOLD); c.rect(478,135,1.1,270,fill=1,stroke=0)
        c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(8.5)); c.drawCentredString(705,400,items[1][0]); _r_wrap(c,items[1][1],525,350,360,15,R_INK,False,6,"center","Times-Roman")
        c.setFillColor(R_PALE_GOLD); c.setStrokeColor(R_GOLD); c.rect(250,95,460,38,fill=1,stroke=1); c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(7)); c.drawString(265,110,"TRADE-OFF"); _r_wrap(c,items[2][1],345,113,345,7.8,R_INK,True,2,"center")
        return
    if kind == "tree":
        c.setFillColor(R_PALE_GOLD); c.setStrokeColor(R_GOLD); c.rect(350,360,260,50,fill=1,stroke=1); c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(7.5)); c.drawCentredString(480,391,items[0][0]); _r_wrap(c,items[0][1],365,378,230,7.6,R_INK,True,2,"center")
        xs=[105,315,525,735]
        for i,(title,body) in enumerate(items[1:5]):
            x=xs[i]; c.setStrokeColor(R_GREEN_2); c.line(480,360,x+60,285); c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(7.2)); c.drawCentredString(x+60,260,title); _r_wrap(c,body,x,235,120,8.4,R_INK,False,5,"center")
        return
    if kind == "context":
        c.setFillColor(R_PALE_RED); c.setStrokeColor(R_RED); c.rect(70,145,300,250,fill=1,stroke=1); c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(8)); c.drawCentredString(220,365,items[0][0]); _r_wrap(c,items[0][1],95,320,250,14,R_INK,False,7,"center","Times-Roman")
        c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(8)); c.drawString(445,355,items[1][0]); _r_wrap(c,items[1][1],445,333,405,12,R_INK,False,5)
        c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(8)); c.drawString(445,245,items[2][0]); _r_wrap(c,items[2][1],445,223,405,12,R_INK,False,5)
        return
    if kind in {"chain","mutation"}:
        n=len(items); left=55; gap=26; total=850; w=(total-gap*(n-1))/n
        for i,(title,body) in enumerate(items):
            x=left+i*(w+gap); c.setFillColor(R_RED if i in {1,3} else R_GREEN); c.setFont("Helvetica-Bold", _snap(7.2)); c.drawCentredString(x+w/2,360,title)
            c.setFillColor(R_PALE_GREEN if i%2==0 else R_WHITE); c.setStrokeColor(R_GREEN_2)
            c.rect(x, BAND_BOTTOM, w, 348 - BAND_BOTTOM, fill=1, stroke=1)
            _r_wrap(c, body, x+10, 328, w-20, 12.0, R_INK, False, 12, "center")
            _tick_label(c, x+11, 173, f"{i+1:02d}", 6.0, R_MUTED, centred=False)
            # The link between two stages is the point of the form; without it
            # the slide is a row of boxes that happen to sit side by side.
            if i:
                _arrow(c, x-gap+3, (348+BAND_BOTTOM)/2, x-3, (348+BAND_BOTTOM)/2)
        if kind=="mutation": c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(11)); c.drawCentredString(480,120,"Change one constraint. Re-run the decision.")
        return
    if kind == "timeline":
        c.setStrokeColor(R_GOLD); c.setLineWidth(2); c.line(100,285,860,285)
        xs=[150,480,810]
        _arrow(c, 855, 285, 872, 285, R_GOLD, 2.0, 7.0)
        for i,(title,body) in enumerate(items[:3]):
            accent = R_RED if i==2 else R_GREEN_2
            # Stems tie each marker to its label; without them the captions
            # float free of the axis they are supposed to sit on.
            c.setStrokeColor(accent); c.setLineWidth(1)
            c.line(xs[i], 292, xs[i], 340); c.line(xs[i], 278, xs[i], 262)
            c.setFillColor(R_WHITE); c.setStrokeColor(accent); c.setLineWidth(1.6)
            c.circle(xs[i],285,7,fill=1,stroke=1)
            c.setFillColor(accent); c.circle(xs[i],285,3,fill=1,stroke=0)
            c.setFillColor(R_RED if i==2 else R_GREEN); c.setFont("Helvetica-Bold", _snap(8)); c.drawCentredString(xs[i],350,title); _r_wrap(c,body,xs[i]-115,245,230,11,R_INK,False,6,"center")
        return
    if kind == "burden":
        _r_wrap(c,items[0][1],75,365,320,14,R_INK,False,5,"center","Times-Roman"); c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(8)); c.drawCentredString(235,405,items[0][0])
        _r_wrap(c,items[1][1],75,235,320,13,R_INK,False,5,"center","Times-Roman"); c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(8)); c.drawCentredString(235,275,items[1][0])
        c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(8)); c.drawString(560,390,items[2][0]); _r_wrap(c,items[2][1],560,365,320,11.5,R_INK,False,5)
        c.setFillColor(R_RED); c.setFont("Helvetica-Bold", _snap(8)); c.drawString(560,260,items[3][0]); _r_wrap(c,items[3][1],560,235,320,11.5,R_INK,False,5)
        return
    if kind == "ai":
        cols=[(65,R_PALE_GREEN,R_GREEN_2),(345,R_PALE_GOLD,R_GOLD),(625,R_PALE_RED,R_RED)]; data=[items[0],items[1],items[3]]
        for (x,fill,stroke),(title,body) in zip(cols,data):
            c.setFillColor(fill); c.setStrokeColor(stroke); c.rect(x,145,250,245,fill=1,stroke=1); c.setFillColor(stroke); c.setFont("Helvetica-Bold", _snap(8)); c.drawCentredString(x+125,360,title); _r_wrap(c,body,x+20,310,210,13,R_INK,False,8,"center","Times-Roman")
        return
    if kind == "portfolio":
        c.setFillColor(R_PALE_GOLD); c.setStrokeColor(R_GOLD); c.circle(480,275,65,fill=1,stroke=1); c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(9)); c.drawCentredString(480,280,"ENGINEERING"); c.drawCentredString(480,267,"MISSION")
        pos=[(75,390),(350,420),(690,390),(75,145),(350,125),(690,145)]
        for i,(title,body) in enumerate(items[:6]):
            x,y=pos[i]
            c.setStrokeColor(R_LINE); c.setLineWidth(0.9); c.line(x+95, y-6, 480, 275)
            c.setFillColor(R_RED if i in {0,3} else R_GREEN); c.setFont("Helvetica-Bold", _snap(7)); c.drawCentredString(x+95,y,title); _r_wrap(c,body,x,y-20,190,8.7,R_INK,False,5,"center")
        # Redraw the hub so the spokes pass behind it, not across the label.
        c.setFillColor(R_PALE_GOLD); c.setStrokeColor(R_GOLD); c.circle(480,275,65,fill=1,stroke=1)
        c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(9)); c.drawCentredString(480,280,"ENGINEERING"); c.drawCentredString(480,267,"MISSION")
        return
    if kind == "argument":
        shown = items[:5]
        # CLAIM -> EVIDENCE -> WARRANT -> COUNTER -> RESIDUAL is a chain of
        # reasoning. It was drawn as five rows that stopped 75pt above the
        # footer with nothing joining them, so the sequence was invisible.
        pitch = (BAND_TOP - BAND_BOTTOM) / max(1, len(shown))
        box_h = pitch - 16
        widths=[760,690,620,550,480]
        for i,(title,body) in enumerate(shown):
            w=widths[i]; x=(960-w)/2
            y = BAND_TOP - pitch * i - box_h / 2
            c.setFillColor(R_RED if i in {0,3,4} else R_GREEN); c.setFont("Helvetica-Bold", _snap(7.5)); c.drawString(x,y,title)
            c.setFillColor(R_PALE_GOLD if i in {0,3} else (R_PALE_GREEN if i in {1,2} else R_PALE_RED))
            c.setStrokeColor(R_GOLD if i in {0,3} else R_GREEN_2)
            c.rect(x+110, y - box_h/2, w-110, box_h, fill=1, stroke=1)
            _r_wrap(c, body, x+122, y + box_h/2 - 12, w-134, 10.5, R_INK, True, 3, "center")
            if i:
                _arrow(c, 480, y + pitch - box_h/2, 480, y + box_h/2 + 2, R_GOLD, 1.2, 5.5)
        return
    if kind == "rubric":
        cols=3; w=260; h=105; x0=60; y0=295
        for i,(title,body) in enumerate(items[:6]):
            r,cc=divmod(i,cols); x=x0+cc*290; y=y0-r*145; c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(7.5)); c.drawCentredString(x+w/2,y+80,title); c.setFillColor(R_PALE_GREEN if i%2==0 else R_WHITE); c.setStrokeColor(R_GREEN_2); c.rect(x,y-10,w,70,fill=1,stroke=1); _r_wrap(c,body,x+12,y+35,w-24,9.5,R_INK,False,4,"center")
        return
    if kind == "verdict":
        # The closing decision slide left the lower half empty. The two options
        # now occupy the band, and the residual-uncertainty bar sits under them
        # rather than floating between them and the footer.
        c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(8)); c.drawCentredString(480,BAND_TOP-8,items[0][0])
        _r_wrap(c,items[0][1],100,BAND_TOP-26,760,17.5,R_INK,False,4,"center","Times-Roman")
        bar_h, gap = 44.0, 18.0
        top, bottom = BAND_TOP - 120, BAND_BOTTOM + bar_h + gap
        box_h = top - bottom
        for x, fill, stroke, (title, body) in (
            (90, R_PALE_GREEN, R_GREEN_2, items[1]),
            (505, R_PALE_RED, R_RED, items[2]),
        ):
            c.setFillColor(fill); c.setStrokeColor(stroke); c.rect(x, bottom, 365, box_h, fill=1, stroke=1)
            c.setFillColor(stroke); c.setFont("Helvetica-Bold", _snap(7.5)); c.drawString(x+15, bottom+box_h-20, title)
            _r_wrap(c, body, x+15, bottom+box_h-38, 330, 11.0, R_INK, True, 6, "center")
        c.setFillColor(R_PALE_GOLD); c.setStrokeColor(R_GOLD); c.rect(220, BAND_BOTTOM, 520, bar_h, fill=1, stroke=1)
        c.setFillColor(R_GREEN); c.setFont("Helvetica-Bold", _snap(9))
        c.drawCentredString(480, BAND_BOTTOM + bar_h/2 - 4, presenter_text(items[3][1], 92))
        return
    _r_bullets(c, items)


def export_cimt_presenter_pdf_v43(bp: Blueprint, out: Path) -> Path:
    out = Path(out)
    c = canvas.Canvas(str(out), pagesize=(PDF_W, PDF_H), pageCompression=1)
    c.setTitle(bp.lecture_title); c.setAuthor("ISCARB Faculty Studio")
    plans = plans_for_blueprint_v42(bp)
    for u, plan in zip(bp.units, plans):
        if u.number != 1 and plan.reuse_mode == "USE" and _r_source(c, bp, u, plan):
            pass
        else:
            _r_redraw(c, bp, u)
        c.showPage()
    c.save()
    return out


# ---------------------------------------------------------------------------
# Browser preview
# ---------------------------------------------------------------------------

def _h(value: str) -> str:
    return html.escape(str(value or ""))


def _html_items(items, cls="bulletList") -> str:
    return f'<div class="{cls}">' + ''.join(
        f'<article><i></i><b>{_h(t)}</b><span>{_h(b)}</span></article>' for t, b in items
    ) + '</div>'


def _html_redraw(bp: Blueprint, u: LectureUnit) -> str:
    kind, items = _spec(bp, u)
    if kind == "title":
        return f'<div class="titleSlide"><h1>{_h(bp.lecture_title)}</h1><p>ISCARB Faculty Studio · CIMT-native lecture</p><hr><blockquote>{_h(items[0][1])}</blockquote><small>20 teaching units · 90 live minutes · source fidelity retained</small></div>'
    if kind == "quote":
        tails=''.join(f'<article><b>0{i+1}</b><span>{_h(b)}</span></article>' for i,(_,b) in enumerate(items[1:4]))
        return f'<div class="quoteSlide"><blockquote>“{_h(items[0][1])}”</blockquote><hr><div>{tails}</div></div>'
    if kind in {"takeaways","table"}:
        return _html_items(items, "tableList" if kind=="table" else "bulletList")
    if kind == "orbit":
        return '<div class="orbit"><strong>DECISION</strong>' + ''.join(f'<article class="o{i}"><b>{_h(t)}</b><span>{_h(b)}</span></article>' for i,(t,b) in enumerate(items[:6])) + '</div>'
    if kind == "ladder":
        return '<div class="ladder">' + ''.join(f'<article style="--i:{i}"><b>{_h(t)}</b><span>{_h(b)}</span></article>' for i,(t,b) in enumerate(items[:4])) + '<p>Predict first. Name the mechanism only after the constraint becomes visible.</p></div>'
    if kind == "curve":
        return f'<div class="curve"><div class="chart"><span class="axisY">COST</span><span class="axisX">assurance →</span><svg viewBox="0 0 500 260" aria-label="assurance cost curve"><polyline points="20,230 90,225 160,212 230,190 300,150 355,105 400,55 435,20"/></svg></div><div class="curveNotes"><b>{_h(items[1][0])}</b><p>{_h(items[1][1])}</p><b>{_h(items[2][0])}</b><p>{_h(items[2][1])}</p></div></div>'
    if kind == "stack":
        return '<div class="stack">' + ''.join(f'<article><b>{_h(t)}</b><span>{_h(b)}</span></article>' for t,b in items[:6]) + '</div>'
    if kind == "compare":
        return f'<div class="compare"><section><b>{_h(items[0][0])}</b><p>{_h(items[0][1])}</p></section><i></i><section><b>{_h(items[1][0])}</b><p>{_h(items[1][1])}</p></section><footer><strong>TRADE-OFF</strong>{_h(items[2][1])}</footer></div>'
    if kind == "tree":
        return '<div class="tree"><header><b>{}</b><span>{}</span></header><div>{}</div></div>'.format(_h(items[0][0]),_h(items[0][1]),''.join(f'<article><b>{_h(t)}</b><span>{_h(b)}</span></article>' for t,b in items[1:5]))
    if kind == "context":
        return f'<div class="context"><section><b>{_h(items[0][0])}</b><p>{_h(items[0][1])}</p></section><div><b>{_h(items[1][0])}</b><p>{_h(items[1][1])}</p><b>{_h(items[2][0])}</b><p>{_h(items[2][1])}</p></div></div>'
    if kind in {"chain","mutation"}:
        extra='<footer>Change one constraint. Re-run the decision.</footer>' if kind=="mutation" else ''
        return '<div class="chain">' + ''.join(f'<article><b>{_h(t)}</b><span>{_h(b)}</span></article>' for t,b in items) + extra + '</div>'
    if kind == "timeline":
        return '<div class="timeline">' + ''.join(f'<article><b>{_h(t)}</b><span>{_h(b)}</span></article>' for t,b in items[:3]) + '</div>'
    if kind == "burden":
        return '<div class="burden">' + ''.join(f'<article><b>{_h(t)}</b><span>{_h(b)}</span></article>' for t,b in items) + '</div>'
    if kind == "ai":
        return '<div class="aiGate">' + ''.join(f'<article><b>{_h(t)}</b><span>{_h(b)}</span></article>' for t,b in (items[0],items[1],items[3])) + f'<footer><b>TEST</b>{_h(items[2][1])}</footer></div>'
    if kind == "portfolio":
        return '<div class="portfolio"><strong>ENGINEERING<br>MISSION</strong>' + ''.join(f'<article class="p{i}"><b>{_h(t)}</b><span>{_h(b)}</span></article>' for i,(t,b) in enumerate(items[:6])) + '</div>'
    if kind == "argument":
        return '<div class="argument">' + ''.join(f'<article><b>{_h(t)}</b><span>{_h(b)}</span></article>' for t,b in items[:5]) + '</div>'
    if kind == "rubric":
        return '<div class="rubric">' + ''.join(f'<article><b>{_h(t)}</b><span>{_h(b)}</span></article>' for t,b in items[:6]) + '</div>'
    if kind == "verdict":
        return f'<div class="verdict"><header><b>{_h(items[0][0])}</b><p>{_h(items[0][1])}</p></header><div><article><b>{_h(items[1][0])}</b><span>{_h(items[1][1])}</span></article><article><b>{_h(items[2][0])}</b><span>{_h(items[2][1])}</span></article></div><footer>{_h(items[3][1])}</footer></div>'
    return _html_items(items)


def render_cimt_presenter_preview_v43(bp: Blueprint, release_state: str = "BLOCKED") -> str:
    plans = plans_for_blueprint_v42(bp)
    slides: list[str] = []; thumbs: list[str] = []
    for i, (u, plan) in enumerate(zip(bp.units, plans)):
        visual = ""
        if u.number != 1 and plan.reuse_mode == "USE":
            uri = _data_uri(_source_body_path(plan))
            if uri:
                visual = f'<div class="sourceVisual"><img src="{uri}" alt="Adapted P1 teaching visual"><small>ADAPTED VISUAL · P1 {plan.source_slide}</small></div>'
        if not visual:
            visual = _html_redraw(bp, u)
        show_decision = u.number in set(range(5,19))
        if u.number == 1:
            content = f'<div class="visual first">{visual}</div>'
        else:
            ready = _readiness_for_unit(bp, u)
            lesson = (f'<div class="learningStrip"><section><b>TRY</b><span>{_h(_try_text(u))}</span></section>'
                      f'<section><b>CHECK</b><span>{_h(_check_text(u))}</span></section>'
                      + (f'<small>{_h(ready)}</small>' if ready else '') + '</div>') if show_decision else ''
            content = f'''<div class="corner"></div><header class="head"><div><h2>{_h(_display_title(bp, u))}</h2><p class="question">{_h(_display_question(u))}</p></div><div class="mark"><i></i><b>CIMT</b><small>ISCARB</small></div></header><div class="visual">{visual}</div>{lesson}<footer class="foot"><span>{_h(_subject(bp))} · source-grounded presenter</span><em>{u.number:02d}/20</em></footer>'''
        slides.append(f'<section class="slide{" show" if i==0 else ""}" data-i="{i}">{content}</section>')
        thumbs.append(f'<button class="thumb{" active" if i==0 else ""}" data-i="{i}"><b>{u.number:02d}</b><span>{_h(presenter_text(_display_title(bp, u),46))}</span></button>')

    css = f'''
    :root{{--green:#005b39;--green2:#2c7e41;--gold:{LEGACY_GOLD_TOKEN};--red:#e22424;--ink:#181818;--muted:#707070;--pale:#ecf6eb;--legacy:{LEGACY_GREEN_TOKEN}}}
    *{{box-sizing:border-box}}body{{margin:0;background:#eceee9;color:var(--ink);font-family:Arial,Helvetica,sans-serif}}.deck{{height:100vh;display:grid;grid-template-columns:220px 1fr}}.rail{{background:#f7f7f2;border-right:1px solid #d9ddd7;padding:18px;overflow:auto}}.brand{{font-family:Georgia,serif;color:var(--green);font-size:19px}}.state{{font-size:9px;color:var(--muted);margin:6px 0 16px}}.thumb{{display:grid;grid-template-columns:27px 1fr;width:100%;gap:7px;padding:8px 4px;border:0;border-bottom:1px solid #dde1da;background:transparent;text-align:left;cursor:pointer;color:#454b46}}.thumb b{{color:var(--gold)}}.thumb span{{font-size:10px}}.thumb.active{{background:#edf5eb;color:#111}}.stage{{display:grid;place-items:center;padding:22px}}.slide{{display:none;width:min(1180px,calc(100vw - 275px));aspect-ratio:16/9;background:#fff;box-shadow:0 22px 50px #1c2a2028;position:relative;overflow:hidden;padding:22px 34px 18px}}.slide.show{{display:block}}.corner{{position:absolute;left:34px;top:20px;width:calc(100% - 68px);height:34px;border-top:2px solid var(--gold);border-left:2px solid var(--gold)}}.head{{height:112px;display:flex;justify-content:space-between;align-items:start;padding:12px 8px 0}}.head h2{{margin:0;color:var(--green);font:400 clamp(27px,2.7vw,40px)/1.05 Georgia,serif;letter-spacing:-.02em}}.question{{margin:8px 0 0;max-width:900px;font-size:13px;line-height:1.32;color:#333}}.mark{{display:grid;grid-template-columns:20px auto;grid-template-rows:15px 12px;align-items:center;column-gap:6px;color:var(--green);font-size:8px;margin-top:2px}}.mark i{{grid-row:1/3;width:18px;height:18px;border:2px solid var(--gold);border-radius:50%}}.mark small{{font-size:6px;color:var(--muted);font-weight:800}}.visual{{height:calc(100% - 196px);display:grid;align-items:center;padding:4px 10px 8px}}.visual.first{{height:100%;padding:0}}.foot{{position:absolute;left:42px;right:42px;bottom:14px;border-top:2px solid var(--gold);padding-top:6px;display:flex;justify-content:space-between;color:var(--muted);font-size:7px}}.foot em{{font-style:normal;font-weight:800}}.decision{{position:absolute;left:48px;right:48px;bottom:39px;display:grid;grid-template-columns:64px 1fr;align-items:center;gap:8px;font-size:9px}}.decision b{{color:var(--red)}}.decision span{{font-weight:700}}.learningStrip{{position:absolute;left:48px;right:48px;bottom:39px;border-top:1px solid var(--gold);padding-top:7px;display:grid;grid-template-columns:1fr 1fr;gap:18px;font-size:9px}}.learningStrip section{{display:grid;grid-template-columns:42px 1fr;gap:7px;align-items:start}}.learningStrip section:first-child b{{color:var(--red)}}.learningStrip section:last-child b{{color:var(--green)}}.learningStrip span{{font-weight:650;line-height:1.25}}.learningStrip small{{position:absolute;right:0;bottom:-14px;color:var(--green);font-size:6px;font-weight:800}}.titleSlide{{height:100%;display:grid;grid-template-rows:1.2fr auto auto 1.4fr auto;align-items:center;text-align:center;padding:35px 70px;position:relative}}.titleSlide:before{{content:'';position:absolute;left:0;right:0;top:0;border-top:2px solid var(--gold)}}.titleSlide h1{{margin:0;color:var(--green);font:400 clamp(34px,4vw,58px)/1.08 Georgia,serif}}.titleSlide p{{font-size:14px}}.titleSlide hr{{width:42%;border:0;border-top:2px solid var(--gold)}}.titleSlide blockquote{{margin:0;color:var(--red);font:400 clamp(18px,2vw,29px)/1.25 Georgia,serif}}.titleSlide small{{color:var(--muted);font-weight:700}}.quoteSlide blockquote{{margin:12px auto;color:var(--red);font:400 clamp(25px,3.1vw,45px)/1.12 Georgia,serif;text-align:center;max-width:880px}}.quoteSlide hr{{width:38%;border:0;border-top:2px solid var(--gold)}}.quoteSlide>div{{display:grid;grid-template-columns:repeat(3,1fr);gap:26px;margin-top:24px}}.quoteSlide article{{display:grid;grid-template-columns:36px 1fr;gap:8px;align-items:start}}.quoteSlide b{{color:var(--gold);font:700 18px Georgia,serif}}.quoteSlide span{{font-size:14px;font-weight:700;line-height:1.35}}.bulletList,.tableList{{display:grid;gap:7px;padding:2px 6px}}.bulletList article{{display:grid;grid-template-columns:10px 150px 1fr;gap:12px;align-items:start;padding:7px 0}}.bulletList i{{width:7px;height:7px;background:var(--gold);margin-top:7px}}.bulletList b{{font-size:11px;color:var(--green)}}.bulletList article:first-child b{{color:var(--red)}}.bulletList span{{font-size:18px;line-height:1.28}}.tableList{{border-top:34px solid var(--gold);position:relative}}.tableList:before{{content:'CHARACTERISTIC                    WHAT IT MEANS IN THE ENGINEERING DECISION';position:absolute;top:-26px;left:12px;color:white;font-size:8px;font-weight:800;white-space:pre}}.tableList article{{display:grid;grid-template-columns:190px 1fr;border:1px solid #dadad0;border-top:0;min-height:50px;align-items:center;padding:6px 12px}}.tableList article:nth-child(odd){{background:#faf6e5}}.tableList i{{display:none}}.tableList b{{color:var(--green);font-size:11px}}.tableList span{{font-size:12px}}.orbit{{height:100%;position:relative}}.orbit>strong{{position:absolute;left:50%;top:47%;transform:translate(-50%,-50%);width:126px;height:88px;border-radius:50%;display:grid;place-items:center;background:var(--green2);color:#fff;font-size:12px}}.orbit article{{position:absolute;width:27%;font-size:12px}}.orbit article b{{display:block;color:var(--green);font-size:10px;margin-bottom:5px}}.orbit .o1 b,.orbit .o4 b{{color:var(--red)}}.orbit .o0{{left:2%;top:9%}}.orbit .o1{{left:2%;top:39%}}.orbit .o2{{left:2%;top:69%}}.orbit .o3{{right:2%;top:9%}}.orbit .o4{{right:2%;top:39%}}.orbit .o5{{right:2%;top:69%}}.ladder{{height:100%;position:relative;padding-top:10px}}.ladder article{{position:absolute;left:calc(3% + var(--i)*23%);bottom:calc(10% + var(--i)*13%);width:22%;height:73px;border:1px solid var(--gold);background:#faf6e5;padding:9px}}.ladder article:nth-child(even){{background:#ecf6eb}}.ladder b{{font-size:9px;color:var(--green)}}.ladder article:first-child b{{color:var(--red)}}.ladder span{{display:block;font-size:10px;margin-top:4px}}.ladder p{{position:absolute;left:4%;top:6%;width:32%;font:400 18px/1.25 Georgia,serif}}.curve{{display:grid;grid-template-columns:1.45fr .9fr;gap:32px;height:100%;align-items:center}}.chart{{height:80%;position:relative;border-left:2px solid #222;border-bottom:2px solid #222}}.chart svg{{position:absolute;inset:8%;width:88%;height:82%}}.chart polyline{{fill:none;stroke:var(--green2);stroke-width:4}}.axisY{{position:absolute;left:-35px;top:42%;font-size:8px;font-weight:800;transform:rotate(-90deg)}}.axisX{{position:absolute;bottom:-22px;left:34%;font-size:8px;color:var(--muted)}}.curveNotes b{{display:block;color:var(--red);font-size:10px;margin:14px 0 5px}}.curveNotes b:nth-of-type(2){{color:var(--green)}}.curveNotes p{{font-size:14px;line-height:1.35}}.stack{{display:grid;gap:8px;padding:4px 80px}}.stack article{{border:1px solid var(--green2);display:grid;grid-template-columns:90px 1fr;padding:8px 12px;align-items:center}}.stack article:nth-child(odd){{background:#ecf6eb}}.stack b{{color:var(--green);font-size:9px}}.stack article:first-child b,.stack article:last-child b{{color:var(--red)}}.stack span{{font-size:13px;text-align:center}}.compare{{height:100%;display:grid;grid-template-columns:1fr 2px 1fr;grid-template-rows:1fr auto;gap:18px 26px;align-items:center}}.compare>i{{height:78%;background:var(--gold)}}.compare section{{text-align:center;padding:8px 24px}}.compare section>b{{color:var(--green);font-size:11px}}.compare section:nth-of-type(2)>b{{color:var(--red)}}.compare section p{{font:400 22px/1.25 Georgia,serif}}.compare footer{{grid-column:1/4;justify-self:center;border:1px solid var(--gold);background:#faf6e5;padding:9px 18px;font-size:10px}}.compare footer strong{{color:var(--red);margin-right:16px}}.tree{{display:grid;grid-template-rows:auto 1fr;gap:24px;height:100%;padding:4px 30px}}.tree header{{justify-self:center;width:36%;border:1px solid var(--gold);background:#faf6e5;padding:9px;text-align:center}}.tree header b{{display:block;color:var(--red);font-size:9px}}.tree header span{{font-size:10px;font-weight:700}}.tree>div{{display:grid;grid-template-columns:repeat(4,1fr);gap:20px;align-items:start}}.tree article{{text-align:center;border-top:2px solid var(--gold);padding-top:12px}}.tree article b{{color:var(--green);font-size:9px}}.tree article span{{display:block;font-size:12px;margin-top:8px;line-height:1.3}}.context{{display:grid;grid-template-columns:.85fr 1.25fr;gap:42px;height:100%;align-items:center}}.context>section{{border:1px solid var(--red);background:#fdeeee;padding:25px;text-align:center;min-height:70%;display:grid;align-content:center}}.context>section b{{color:var(--red);font-size:10px}}.context>section p{{font:400 20px/1.3 Georgia,serif}}.context>div b{{display:block;color:var(--green);font-size:10px;margin-top:12px}}.context>div b:nth-of-type(2){{color:var(--red)}}.context>div p{{font-size:15px;line-height:1.35}}.chain{{display:grid;grid-template-columns:repeat(4,1fr);gap:14px;align-items:center;height:100%;position:relative}}.chain article{{border:1px solid var(--green2);min-height:180px;padding:15px;display:grid;align-content:start;text-align:center}}.chain article:nth-child(odd){{background:#ecf6eb}}.chain b{{color:var(--green);font-size:9px}}.chain article:nth-child(2) b,.chain article:nth-child(4) b{{color:var(--red)}}.chain span{{font-size:14px;line-height:1.35;margin-top:14px}}.chain footer{{position:absolute;bottom:6px;left:25%;right:25%;text-align:center;color:var(--red);font:700 13px Georgia,serif}}.timeline{{display:grid;grid-template-columns:repeat(3,1fr);gap:35px;height:100%;align-items:center;position:relative}}.timeline:before{{content:'';position:absolute;left:7%;right:7%;top:49%;height:2px;background:var(--gold)}}.timeline article{{position:relative;text-align:center;z-index:2}}.timeline article:before{{content:'';display:block;width:16px;height:16px;border-radius:50%;background:var(--green2);margin:0 auto 18px}}.timeline article:last-child:before{{background:var(--red)}}.timeline b{{display:block;color:var(--green);font-size:10px}}.timeline article:last-child b{{color:var(--red)}}.timeline span{{display:block;font-size:14px;line-height:1.35;margin-top:50px}}.burden{{display:grid;grid-template-columns:1fr 1fr;gap:60px;padding:10px 20px}}.burden article{{padding:13px 20px;border-bottom:2px solid var(--gold)}}.burden b{{display:block;color:var(--green);font-size:10px}}.burden article:nth-child(2) b,.burden article:nth-child(4) b{{color:var(--red)}}.burden span{{display:block;font:400 17px/1.3 Georgia,serif;margin-top:10px}}.aiGate{{display:grid;grid-template-columns:repeat(3,1fr);gap:22px;height:100%;align-items:center;position:relative}}.aiGate article{{min-height:220px;padding:20px;border:1px solid var(--green2);background:#ecf6eb;text-align:center;display:grid;align-content:center}}.aiGate article:nth-child(2){{background:#faf6e5;border-color:var(--gold)}}.aiGate article:nth-child(3){{background:#fdeeee;border-color:var(--red)}}.aiGate b{{color:var(--green);font-size:10px}}.aiGate article:nth-child(3) b{{color:var(--red)}}.aiGate span{{font:400 16px/1.3 Georgia,serif;margin-top:12px}}.aiGate footer{{position:absolute;bottom:0;left:28%;right:28%;text-align:center;font-size:10px}}.aiGate footer b{{color:var(--red);margin-right:10px}}.portfolio{{height:100%;position:relative}}.portfolio>strong{{position:absolute;left:50%;top:48%;transform:translate(-50%,-50%);width:150px;height:105px;border:1px solid var(--gold);background:#faf6e5;border-radius:50%;display:grid;place-items:center;text-align:center;color:var(--green);font-size:12px}}.portfolio article{{position:absolute;width:27%;text-align:center}}.portfolio article b{{display:block;color:var(--green);font-size:9px}}.portfolio .p0 b,.portfolio .p3 b{{color:var(--red)}}.portfolio article span{{display:block;font-size:11px;margin-top:6px}}.portfolio .p0{{left:1%;top:4%}}.portfolio .p1{{left:36%;top:0}}.portfolio .p2{{right:1%;top:4%}}.portfolio .p3{{left:1%;bottom:4%}}.portfolio .p4{{left:36%;bottom:0}}.portfolio .p5{{right:1%;bottom:4%}}.argument{{display:grid;gap:8px;padding:5px 65px}}.argument article{{display:grid;grid-template-columns:130px 1fr;align-items:center;margin:0 auto;border:1px solid var(--green2);background:#ecf6eb;padding:7px 12px}}.argument article:nth-child(1),.argument article:nth-child(4){{background:#faf6e5;border-color:var(--gold)}}.argument article:nth-child(5){{background:#fdeeee;border-color:var(--red)}}.argument article:nth-child(1){{width:100%}}.argument article:nth-child(2){{width:92%}}.argument article:nth-child(3){{width:84%}}.argument article:nth-child(4){{width:76%}}.argument article:nth-child(5){{width:68%}}.argument b{{font-size:9px;color:var(--green)}}.argument article:nth-child(1) b,.argument article:nth-child(4) b,.argument article:nth-child(5) b{{color:var(--red)}}.argument span{{font-size:11px;text-align:center}}.rubric{{display:grid;grid-template-columns:repeat(3,1fr);gap:18px 22px;padding:10px 22px}}.rubric article{{text-align:center}}.rubric b{{display:block;color:var(--green);font-size:10px;margin-bottom:7px}}.rubric span{{display:grid;place-items:center;border:1px solid var(--green2);background:#ecf6eb;min-height:80px;padding:10px;font-size:12px}}.verdict header{{text-align:center;padding:8px 40px}}.verdict header b{{color:var(--green);font-size:10px}}.verdict header p{{font:400 21px/1.25 Georgia,serif}}.verdict>div{{display:grid;grid-template-columns:1fr 1fr;gap:28px;margin:15px 50px}}.verdict article{{border:1px solid var(--green2);background:#ecf6eb;padding:14px}}.verdict article:nth-child(2){{border-color:var(--red);background:#fdeeee}}.verdict article b{{color:var(--green);font-size:9px}}.verdict article:nth-child(2) b{{color:var(--red)}}.verdict article span{{display:block;font-size:12px;margin-top:8px}}.verdict footer{{margin:18px auto 0;width:66%;border:1px solid var(--gold);background:#faf6e5;padding:10px;text-align:center;color:var(--green);font-weight:800;font-size:11px}}.sourceVisual{{height:100%;display:grid;place-items:center;position:relative;overflow:hidden}}.sourceVisual img{{max-width:100%;max-height:100%;object-fit:contain}}.sourceVisual small{{position:absolute;right:2px;bottom:2px;background:#fff;padding:3px 6px;color:var(--green);font-size:6px;font-weight:800}}
    @media(max-width:850px){{html,body,.deck,.stage{{max-width:100%;overflow-x:hidden}}.deck{{grid-template-columns:minmax(0,1fr)}}.rail{{display:none}}.stage{{padding:8px;min-width:0}}.slide{{width:calc(100vw - 16px);max-width:calc(100vw - 16px);min-width:0;padding-left:18px;padding-right:18px}}.head,.visual,.slide>*{{min-width:0}}.corner{{left:18px;width:calc(100% - 36px)}}.foot{{left:22px;right:22px}}.decision{{left:24px;right:24px}}.visual{{height:calc(100% - 168px)}}.learningStrip{{left:24px;right:24px;gap:8px;grid-template-columns:1fr}}.learningStrip section{{grid-template-columns:36px 1fr}}.bulletList article{{grid-template-columns:8px minmax(72px,28%) minmax(0,1fr)}}.tableList article{{grid-template-columns:minmax(88px,34%) minmax(0,1fr)}}.stack{{padding-left:12px;padding-right:12px}}.argument{{padding-left:8px;padding-right:8px}}.argument article{{grid-template-columns:minmax(68px,30%) minmax(0,1fr)}}.compare section{{padding-left:8px;padding-right:8px}}.tree{{padding-left:8px;padding-right:8px}}.context{{gap:14px}}.chain{{gap:6px}}.timeline{{gap:10px}}.burden{{gap:12px;padding-left:4px;padding-right:4px}}.verdict>div{{margin-left:8px;margin-right:8px;gap:10px}}}}
    '''
    js = '''<script>(function(){const s=[...document.querySelectorAll('.slide')],b=[...document.querySelectorAll('.thumb')];function go(i){s.forEach((x,j)=>x.classList.toggle('show',j===i));b.forEach((x,j)=>x.classList.toggle('active',j===i));}b.forEach((x,i)=>x.onclick=()=>go(i));document.onkeydown=e=>{let i=s.findIndex(x=>x.classList.contains('show'));if(e.key==='ArrowRight'||e.key==='PageDown')go(Math.min(s.length-1,i+1));if(e.key==='ArrowLeft'||e.key==='PageUp')go(Math.max(0,i-1));};})();</script>'''
    return f'<!doctype html><html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>{_h(bp.lecture_title)} · ISCARB Presenter</title><style>{css}</style></head><body><div class="deck"><aside class="rail"><div class="brand">ISCARB · CIMT-native Presenter</div><div class="state">{_h(release_state)} · 20 units · 90 minutes</div>{"".join(thumbs)}</aside><main class="stage">{"".join(slides)}</main></div>{js}</body></html>'
