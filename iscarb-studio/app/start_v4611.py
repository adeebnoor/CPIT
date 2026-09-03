from __future__ import annotations

"""ISCARB Faculty Studio v4.6.11 quality hotfix.

This patch makes the 20-unit grammar a presentation-quality contract, not just
metadata. It auto-replaces stored hollow drafts with the deterministic,
source-complete review draft before preview/export, raises the density floor,
and preserves any canvas overflow in the unit evidence / PPTX speaker notes.
"""
from pathlib import Path

from . import start_v440 as base
from . import gate_v14
from . import main as engine_main
from . import deterministic_blueprint_fallback as draft_builder
from . import presenter_v44
from .free_workflow import load_bundle
from .storage import UPLOADS

app = base.app
engine = base.engine
PUBLIC_VERSION = "4.6.11"
PIPELINE_ID = "faculty-studio-v4.6.11-strict-20-unit-output-quality"

# ---------------------------------------------------------------------------
# 1) A slide with a few labels is not a teaching unit.
# ---------------------------------------------------------------------------
gate_v14.MIN_TEACHING_WORDS_WITH_SOURCE_VISUAL = 24
gate_v14.MIN_TEACHING_WORDS_WITHOUT_VISUAL = 35


def _strict_presenter_density(bp):
    for u in bp.units:
        words = gate_v14._teaching_payload_words(u)
        floor = 20 if u.number == 19 else 24
        if words < floor:
            return False
        visible_items = len([x for x in (*u.core_content, *u.pedagogy_content) if str(x).strip()])
        if u.number != 19 and visible_items < 3:
            return False
        if not str(u.student_action or "").strip() or not str(u.takeaway or "").strip():
            return False
    return True


gate_v14._presenter_density_ok = _strict_presenter_density

# ---------------------------------------------------------------------------
# 2) Never silently throw source detail away while fitting the canvas.
#    The visible slide stays readable; removed detail is retained in evidence
#    and therefore can be carried into editable-PPTX speaker notes.
# ---------------------------------------------------------------------------
_original_drop = draft_builder._drop_last_statement


def _drop_last_statement_preserving_source(unit, teaching: bool):
    before = [str(x).strip() for x in unit.core_content if str(x).strip()]
    removed = before[-1] if len(before) >= 2 else ""
    changed = _original_drop(unit, teaching)
    if changed and removed:
        marker = "PRESERVED SOURCE DETAIL (speaker notes): "
        if marker + removed not in unit.evidence:
            unit.evidence = (str(unit.evidence or "").strip() + " " + marker + removed).strip()
    return changed


draft_builder._drop_last_statement = _drop_last_statement_preserving_source

# ---------------------------------------------------------------------------
# 3) Treat hollow/illegible presentation checks as critical. A semantic draft
#    does not get to stay merely because its metadata is technically complete.
# ---------------------------------------------------------------------------
_original_critical = engine_main._critical_presenter_failures


def _critical_presenter_failures(checks):
    failures = list(_original_critical(checks))
    for name in (
        "v14_no_unit_is_a_near_empty_slide",
        "v14_technical_units_have_teaching_density",
        "v15_complete_20_unit_grammar",
        "v15_presenter_fits_readable_canvas",
        "v15_technical_units_retain_source_detail",
    ):
        if checks.get(name) is False and name not in failures:
            failures.append(name)
    return failures


engine_main._critical_presenter_failures = _critical_presenter_failures

# ---------------------------------------------------------------------------
# 4) Self-heal old jobs at the moment the user opens/downloads them.
# ---------------------------------------------------------------------------
_original_presenter_job = base._presenter_job


def _presenter_job_v4611(job_id: str):
    job = _original_presenter_job(job_id)
    failures = _critical_presenter_failures(job.deterministic_checks or {})
    if failures and job.source_profile is not None:
        try:
            bundle = load_bundle(job, UPLOADS / job_id)
            job.blueprint = engine_main._source_preserving_draft(job.source_profile, bundle)
            checks = engine.deterministic_gate(job.blueprint, job.source_profile, bundle.combined_local_text())
            checks.update(engine_main.session_scope_gate(job.blueprint, job.source_profile, bundle))
            job.deterministic_checks = checks
            engine.save_job(job)
        except Exception:
            # Never destroy a saved draft if its historical source bundle has
            # expired; the normal endpoint will still show the saved artifact.
            pass
    return job


base._presenter_job = _presenter_job_v4611

# ---------------------------------------------------------------------------
# 5) Editable PPTX notes retain source overflow/evidence for faculty use.
# ---------------------------------------------------------------------------
_original_pptx = base.export_presenter_pptx


def _pptx_with_preserved_notes(bp, out: Path, source_root=None, release_state="REVIEW"):
    path = _original_pptx(bp, out, source_root=source_root, release_state=release_state)
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        for slide, unit in zip(prs.slides, bp.units):
            evidence = str(unit.evidence or "").strip()
            if "PRESERVED SOURCE DETAIL (speaker notes):" in evidence:
                tf = slide.notes_slide.notes_text_frame
                tf.text = tf.text + "\n\n[Preserved source detail]\n" + evidence
        prs.save(str(path))
    except Exception:
        pass
    return path


base.export_presenter_pptx = _pptx_with_preserved_notes

# Update all user-facing health/version strings without replacing routes.
base.PUBLIC_VERSION = PUBLIC_VERSION
base.PIPELINE_ID = PIPELINE_ID
base.engine.health = base._health_v440
