from __future__ import annotations

"""v7.2.3 — source-faithful Kickoff / Orientation mode.

The 20-unit grammar remains fixed, but a course-launch P1 must not be forced
through chapter-mechanism templates that turn grading, roadmap and logistics
into fake technical mechanisms. This patch also makes PPTX table cells visible
to deterministic profiling so schedules and milestone tables cannot disappear.
"""

import re
from pathlib import Path
from fastapi.responses import HTMLResponse

from . import main as engine
from . import start_v440 as base
from . import source_bundle as source_bundle_mod
from . import source_profile_fallback as profile_fallback
from . import source_text as source_text_mod
from .deterministic_blueprint_fallback import fit_presenter_text
from .models import CoverageEvidence, TopicFamily, VisualPlan

_PATCHED = False


def _shape_lines(shape):
    out=[]
    try:
        if getattr(shape, "has_text_frame", False) and getattr(shape, "text", "").strip():
            out.extend(x.strip() for x in str(shape.text).splitlines() if x.strip())
    except Exception:
        pass
    try:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells=[" ".join(str(cell.text or "").split()) for cell in row.cells]
                cells=[x for x in cells if x]
                if cells:
                    out.append(" | ".join(cells))
    except Exception:
        pass
    try:
        # GroupShape is iterable; ordinary shapes are not.
        for child in shape.shapes:
            out.extend(_shape_lines(child))
    except Exception:
        pass
    return out


def _pptx_text(path: Path, limit: int = 600_000) -> str:
    from pptx import Presentation
    prs=Presentation(str(path)); chunks=[]
    for slide in prs.slides:
        for shape in slide.shapes:
            chunks.extend(_shape_lines(shape))
    return "\n".join(chunks)[:limit]


def _pptx_chunks(path: Path):
    from pptx import Presentation
    prs=Presentation(str(path)); out=[]
    for i,slide in enumerate(prs.slides,1):
        raw="\n".join(line for sh in slide.shapes for line in _shape_lines(sh))
        lines=profile_fallback._meaningful_lines(raw)
        if not lines:
            continue
        label=profile_fallback._choose_label(lines,i)
        orientation_contact=any(k in label.lower() for k in ("logistics","contact","next steps","course information","course details"))
        body=[]
        for x in lines:
            furniture=profile_fallback._is_furniture_line(x)
            if not furniture or (orientation_contact and profile_fallback._CONTACT_RE.search(x)):
                body.append(x)
        out.append((i,label," · ".join(body)))
    return out


def _is_kickoff(profile, bundle) -> bool:
    try:
        text=bundle.combined_local_text().lower()
    except Exception:
        text=" ".join([getattr(profile,"weekly_focus","") or "", getattr(profile,"lecture_title","") or ""]).lower()
    course_signal=any(k in text for k in ("welcome to", "semester", "syllabus", "course launch", "kick-off", "kickoff"))
    assessment_signal=any(k in text for k in ("midterm", "final exam", "earning your grade", "grading", "assessment"))
    schedule_signal=any(k in text for k in ("roadmap", "milestones", "logistics", "theory lectures", "lab sessions"))
    return course_signal and assessment_signal and schedule_signal


def _find(profile, *needles):
    for item in list(getattr(profile,"coverage_items",[]) or []):
        blob=(str(getattr(item,"label","") or "")+" "+str(getattr(item,"why_important","") or "")).lower()
        if any(n.lower() in blob for n in needles):
            return item
    return None


def _statements(row, max_items=20):
    if row is None:
        return []
    raw=str(getattr(row,"why_important","") or getattr(row,"label","") or "")
    label=" ".join(str(getattr(row,"label","") or "").split()).strip()
    parts=re.split(r"\s*[·•▪■◆❑❒❏]\s*|\n+",raw)
    out=[]
    for part in parts:
        text=" ".join(str(part).split()).strip(" -:;,\t")
        if not text or text==label or text in out:
            continue
        out.append(text)
    return out[:max_items]


def _course_title(bundle, profile):
    try:
        text=bundle.combined_local_text()
        for line in text.splitlines():
            clean=" ".join(line.split()).strip()
            if re.search(r"\b[A-Z]{2,6}\s*[- ]?\s*\d{3,4}\b",clean) and 2 <= len(clean.split()) <= 12:
                if "PRIMARY:" not in clean.upper():
                    return clean
    except Exception:
        pass
    focus=str(getattr(bundle,"lecture_focus","") or "").strip()
    if focus:
        return focus.split(":",1)[0].strip(" -—")
    return str(getattr(profile,"lecture_title","") or "Course Kickoff").strip()


def _ensure_kickoff_profile(profile,bundle):
    if not _is_kickoff(profile,bundle):
        return profile
    flow=_find(profile,"iscarb engineering flow")
    case=_find(profile,"60-second micro-case","micro-case")
    ai=_find(profile,"professional standards & ai literacy","ai literacy")
    grade=_find(profile,"earning your grade","midterm exam")
    roadmap=_find(profile,"roadmap","milestones")
    logistics=_find(profile,"logistics & next steps","theory lectures","lab sessions")
    rows=[x for x in (flow,case,ai,grade,roadmap,logistics) if x is not None]
    families=[]
    for row in rows:
        name=str(row.label).strip()
        if name and name not in [x.name for x in families]:
            families.append(TopicFamily(name=name,source_anchor=row.source_anchor,why_important=row.why_important))
    if families:
        profile.topic_families=families
        profile.in_scope_families=[x.name for x in families]
        profile.deferred_topics=[]
    profile.lecture_title=_course_title(bundle,profile)
    profile.weekly_focus=str(getattr(bundle,"lecture_focus","") or profile.weekly_focus or "Course kickoff and engineering readiness")
    profile.source_warnings=[x for x in profile.source_warnings if "source too thin" not in str(x).lower()]
    return profile


def _compact(lines, start=0, stop=None):
    rows=list(lines[start:stop])
    return [x for x in rows if str(x).strip()]


def _combine(lines, max_groups=4):
    """Group short source rows without deleting them; useful for roadmap tables."""
    rows=[" ".join(str(x).split()) for x in lines if str(x).strip()]
    if len(rows)<=max_groups:
        return rows
    groups=[[] for _ in range(max_groups)]
    for i,row in enumerate(rows):
        groups[min(max_groups-1, i*max_groups//len(rows))].append(row)
    return ["; ".join(g) for g in groups if g]


def _visual(kind,purpose,anchor="",source_slide=""):
    use=bool(source_slide)
    return VisualPlan(
        visual_type=kind,teaching_purpose=purpose,
        source_visual_available=use,source_page_or_slide=source_slide,
        reuse_mode="USE" if use else "NEW",citation=anchor or "ISCARB source-derived visual",
        focal_elements=[],annotation_plan=[],
        visual_evidence_role="Use the P1 visual when it carries information; otherwise use an ISCARB native redraw or text-first composition.",
    )


def _specialize_kickoff(bp,profile,bundle):
    if not _is_kickoff(profile,bundle):
        return _strip_retired_readiness(bp)

    intro=_find(profile,"welcome to","stop memorizing","design is how it works")
    flow=_find(profile,"iscarb engineering flow")
    case=_find(profile,"60-second micro-case","micro-case")
    ai=_find(profile,"professional standards & ai literacy","ai literacy")
    grade=_find(profile,"earning your grade","midterm exam")
    roadmap=_find(profile,"roadmap","milestones")
    logistics=_find(profile,"logistics & next steps","theory lectures","lab sessions")

    intro_s=_statements(intro); flow_s=_statements(flow); case_s=_statements(case)
    ai_s=_statements(ai); grade_s=_statements(grade); road_s=_statements(roadmap); log_s=_statements(logistics)
    title=_course_title(bundle,profile)
    crisis=next((x for x in case_s if x.lower().startswith("crisis")),None) or next((x for x in ai_s if "failure of accountability" in x.lower()),None) or "What evidence is enough to defend an engineering decision?"

    bp.lecture_title=title
    bp.central_engineering_crisis=crisis
    bp.engineering_thesis="CPIT-455 asks learners to stop memorizing definitions and start making critical engineering decisions through an explicit crisis-to-verdict evidence flow."
    bp.named_ethical_purpose="Use AI to accelerate work without delegating accountability: independent validation and explicit human sign-off remain required."

    clos=list(bp.clOs)
    clos[0].statement="Explain the ISCARB flow from CRISIS to VERDICT and the Predict → Constraint → Derive → Name reasoning sequence."
    clos[0].evidence_expected="A source-anchored map of the ISCARB flow."
    clos[1].statement="Apply the 60-second micro-case to distinguish technical test success from organizational assurance evidence."
    clos[1].evidence_expected="A bounded case decision with MAP, EVIDENCE and VERDICT."
    clos[2].statement="Use AI as an engineering aid while retaining independent validation, accountability and human sign-off."
    clos[2].evidence_expected="An AI-assisted artifact plus an independent verification note."
    clos[3].statement="Explain how assessment rewards defended decisions, Decision Cards and live constraint mutation rather than rote recall."
    clos[3].evidence_expected="A personal evidence plan aligned to the assessment structure."
    clos[4].statement="Use the course roadmap, logistics and action items to prepare the next milestone."
    clos[4].evidence_expected="A concrete next-step checklist tied to the published schedule."
    bp.clOs=clos

    source_rows=[x for x in (flow,case,ai,grade,roadmap,logistics) if x is not None]
    bp.source_topic_families=[str(x.label) for x in source_rows]

    def assign(n,title_,question,core,ped,action,takeaway,rows=(),kind="concept-map",source_slide=""):
        u=bp.units[n-1]
        u.title=title_; u.engineering_question=question
        u.core_content=[str(x).strip() for x in core if str(x).strip()]
        u.pedagogy_content=[str(x).strip() for x in ped if str(x).strip()]
        u.enrichment_content=[]; u.enrichment_basis=[]; u.scenario_assumptions=[]
        u.student_action=action; u.takeaway=takeaway; u.visual_suggestion=kind
        anchors=[]; evid=[]
        for row in rows:
            if row is None: continue
            if row.source_anchor not in anchors: anchors.append(row.source_anchor)
            # Evidence excerpt must be learner-visible and substantive.
            row_blob=(str(row.label)+" "+str(row.why_important)).lower()
            excerpt=next((x for x in u.core_content if len(x.split())>=4 and any(tok in row_blob for tok in x.lower().split()[:3])),None)
            if excerpt and len(excerpt)>=20:
                evid.append(CoverageEvidence(coverage_id=row.id,source_anchor=row.source_anchor,visible_excerpt=excerpt))
        u.source_anchor="; ".join(anchors) if anchors else "N/A — ISCARB PEDAGOGY"
        u.coverage_evidence=evid
        u.evidence=u.source_anchor
        u.verify_before_release=False
        u.visual_plan=_visual(kind,question,u.source_anchor,source_slide)
        return u

    # U01-U05: launch the learner into the course decision culture.
    welcome=[x for x in intro_s if "welcome" in x.lower() or "stop memorizing" in x.lower() or "critical engineering" in x.lower()]
    if not welcome: welcome=intro_s[:2]
    assign(1,"Welcome to engineering decisions","What changes when this course rewards defended decisions rather than memorized definitions?",welcome,
           ["DECISION — commit to one engineering judgment you expect to defend this semester.","EVIDENCE — name what would have to be inspectable before that judgment deserves confidence."],
           "CHECKPOINT — follow the course contract; respond only if the lecturer calls for a quick check.","The course is organized around decisions that can be defended with evidence.",(intro,),"title")

    spine=["ISCARB Engineering Flow","60-Second Engineering Micro-Case","Professional Standards & AI Literacy","Assessment & Engineering Defense","Roadmap & Milestones","Logistics & Next Steps"]
    assign(2,"Week 1 map","How do the six parts of this kickoff fit into one readiness story?",spine,
           ["MAP — connect method, case, professional rules, assessment, milestones and next action; this is a curated map, not a source-heading dump."],
           "CHECKPOINT — sketch the six-node map.","Week 1 establishes how you will think, prove, be assessed and prepare.",source_rows,"concept-map")

    assign(3,"Five outcomes for the kickoff","What will count as visible evidence that you understood the course contract?",[],[f"{c.id}: {c.statement}" for c in bp.clOs],
           "CHECKPOINT — choose the outcome that will be hardest to prove.","Every outcome requires an observable learner artifact, not recognition alone.",(),"table")

    quote=[x for x in intro_s if "design" in x.lower()]
    defense=[x for x in grade_s if "engineering defense" in x.lower() or "mutate your constraints" in x.lower()]
    assign(4,"The engineering judgment contract","Why is 'design is how it works' a stronger course rule than 'know the definition'?",quote+defense,
           ["REASON — connect function, evidence and defendability.","ACCOUNTABILITY — a decision is not complete until someone can inspect why it was made."],
           "CHECKPOINT — name one difference between recall and defendable engineering judgment.","Good design is judged by how it works under evidence and changing constraints.",(intro,grade),"concept-map")

    case_crisis=[x for x in case_s if x.lower().startswith("crisis")]
    assign(5,"Prediction gate: what failed?","PREDICT before seeing the full diagnosis: if the software passes its tests, what could still block deployment?",case_crisis,
           ["PREDICT: Commit before the MAP is revealed.","CONSTRAINT: Technical tests are not the only evidence in the case.","DERIVE: State what other assurance layer could fail.","NAME: Name the layer only after defending the reasoning."],
           "THINK–PAIR–SHARE · 1 MIN — make one prediction, compare, then name the evidence that would change it.","Prediction comes before explanation.",(case,),"process")

    # U06-U10: teach the source's method, AI rule, assessment and roadmap.
    flow_core=[]
    steps=[x for x in flow_s if re.match(r"^[1-5]\.\s",x)]
    if steps: flow_core=[" → ".join(re.sub(r"^\d+\.\s*","",x) for x in steps)]
    flow_core += [x for x in flow_s if "predict" in x.lower() or "solving failures" in x.lower()]
    assign(6,"ISCARB Engineering Flow","How does a crisis become a bounded engineering verdict?",flow_core,
           ["TRACE — CRISIS → MAP → TRADE-OFF → EVIDENCE → VERDICT.","REASON — Predict → Constraint → Derive → Name before you label the principle.","DECISION BOX — What evidence changes the verdict?"],
           "CHECKPOINT — trace the five stages in order.","The flow turns a failure into an inspectable decision.",(flow,),"process","SLIDE 3")

    assign(7,"ISCARB in action: 60-second micro-case","Can you trace the hospital case from observed crisis to verdict without skipping evidence?",_compact(case_s,0,8),
           ["TRACE — label CRISIS, MAP, TRADE-OFF, EVIDENCE and VERDICT.","DECISION BOX — reject or accept only after the evidence step."],
           "CHECKPOINT — run the source micro-case in 60 seconds.","Technical test success does not substitute for organizational assurance evidence.",(case,),"decision-chain")

    assign(8,"Professional Standards & AI Literacy","When does AI accelerate engineering, and when does reliance on it become an accountability failure?",_compact(ai_s,0,8),
           ["ALTERNATIVE A — use AI to accelerate preparation.","ALTERNATIVE B — accept no AI output until independently validated.","TRADE-OFF — speed is useful only while assurance remains owned by a human.","DECISION BOX — What evidence would let you sign?"],
           "THINK–PAIR–SHARE · 1 MIN — choose the boundary between assistance and delegated accountability.","AI may prepare evidence; it does not own assurance.",(ai,),"assurance-chain")

    assign(9,"Earning your grade: defend your decisions","What behavior is each assessment component actually rewarding?",_compact(grade_s,0,10),
           ["MEASURE — connect each percentage to the learner performance it evaluates.","EVIDENCE — Decision Cards and Engineering Defense make reasoning inspectable, not merely submitted.","DECISION BOX — What would a strong artifact have to prove?"],
           "CHECKPOINT — identify which assessment demands the strongest live defense.","The grading model rewards bounded reasoning and defended decisions, not rote homework.",(grade,),"assessment-map")

    road_rows=[x for x in road_s if re.search(r"\b(?:sept|oct|nov|dec|jan|feb|mar|apr|may|jun|jul|aug)\b",x,re.I) or "|" in x]
    first_half=road_rows[:max(1,(len(road_rows)+1)//2)] if road_rows else road_s[:4]
    assign(10,"Roadmap & milestones: first half","Which milestone changes what you must prepare next?",_combine(first_half,3),
           ["KNOWN — use the published sequence, not an invented calendar.","DECISION — choose the next milestone that changes your preparation plan.","DECISION BOX — What evidence must exist before that milestone?"],
           "THINK–PAIR–SHARE · 1 MIN — identify the next decision-sensitive milestone.","A roadmap is useful only when it changes preparation before the deadline.",(roadmap,),"timeline","SLIDE 7")

    # U11-U15: deepen the same source story without inventing a later chapter.
    assign(11,"Core case: technical success, assurance failure","The code passed its tests; what exactly still prevents a responsible deployment?",_compact(case_s,0,8),
           ["MICRO-EXAMPLE — the source already supplies the bounded hospital case; use it before generalizing.","TRANSFER — observed condition → affected assurance layer → evidence owner → decision → reversal condition.","DECISION BOX — What evidence would reverse REJECT?"],
           "CORE IN-CLASS CASE · 5–7 MIN — analyse the source case as a team and defend one verdict.","A complete system decision includes organizational evidence, not only technical test results.",(case,),"case-map")

    accountability=[x for x in case_s if x.lower().startswith(("trade-off","evidence","verdict"))]
    accountability += [x for x in ai_s if "sign-off" in x.lower() or "independent validation" in x.lower() or "failure of accountability" in x.lower()]
    assign(12,"Core case: speed versus accountable assurance","How much deployment speed are you willing to trade for evidence that another engineer can independently inspect?",_combine(accountability,4),
           ["OWNER — name who owns the decision.","EVIDENCE — name what must be independently inspectable.","SIGN-OFF — state who accepts the residual risk.","DECISION BOX — What evidence changes the speed/assurance trade-off?"],
           "CORE IN-CLASS CASE · 5–7 MIN — defend one side, then state the evidence that would make you switch.","Owner, evidence and sign-off are the control against blind acceleration.",(case,ai),"assurance-chain")

    second_half=road_rows[len(first_half):] if road_rows else road_s[4:]
    if not second_half: second_half=road_rows[-3:] if road_rows else road_s[-3:]
    assign(13,"Roadmap & milestones: pressure points","Where do project review and final defense change the evidence you must have ready?",_combine(second_half,3),
           ["TREND — read the semester as a sequence of increasing evidence obligations.","FAIL-FIRST — identify the milestone at which weak evidence would become visible.","REDESIGN — move preparation earlier rather than compressing it into the deadline."],
           "CHECKPOINT — mark the first milestone your current preparation would fail.","Milestones are evidence gates, not dates to memorize.",(roadmap,),"timeline","SLIDE 7")

    ai_mis=[x for x in ai_s if "prepare evidence" in x.lower() or "not an engineer" in x.lower() or "independent validation" in x.lower()]
    assign(14,"Misconception: AI output is not assurance","What is wrong with treating a fluent AI answer as engineering evidence?",ai_mis,
           ["PLAUSIBLE-BUT-WRONG — 'the model produced it, so the work is done.'","SOURCE CHECK — AI can prepare evidence; the learner still owns validation and justification.","OPERATING CONSEQUENCE — unsupported AI output fails accountability."],
           "THINK–PAIR–SHARE · 1 MIN — give one AI use you would allow and one you would reject.","Fluency is not evidence; accountability remains human.",(ai,),"misconception")

    audit=[x for x in ai_s if "human sign-off" in x.lower() or "independent validation" in x.lower()]
    audit += [x for x in grade_s if "engineering defense" in x.lower() or "mutate" in x.lower()]
    assign(15,"Assurance gate: defend before you sign","What must remain true when AI assists, constraints change, or a claim cannot be explained directly?",_combine(audit,3),
           ["AI ASSURANCE LENS — model output and system assurance are different claims.","BLACK-BOX TEST — require observable evidence that an independent reviewer can inspect.","HUMAN SIGN-OFF — defend the result under a changed constraint before accepting it."],
           "THINK–PAIR–SHARE · 1 MIN — name the one piece of evidence you would require before signing.","The engineer signs because the evidence survives challenge, not because the model sounds confident.",(ai,grade),"audit")

    # U16-U20: convert the kickoff into concrete learner artifacts and next actions.
    decision_card=[x for x in grade_s if "decision card" in x.lower() or "assignments" in x.lower()]
    assign(16,"Post-class artifact: Decision Card #0","What would your first inspectable decision record contain?",decision_card,
           ["CLAIM — one bounded decision.","P1 EVIDENCE — the source statement that supports it.","TRADE-OFF — the cost you accept.","REVERSAL CONDITION — what new evidence would change the decision."],
           "POST-CLASS ARTIFACT — draft one Decision Card individually after class.","Your first artifact should make a decision challengeable by another engineer.",(grade,),"artifact")

    mutation=[x for x in grade_s if "mutate" in x.lower() or "group project" in x.lower() or "engineering defense" in x.lower()]
    assign(17,"Post-class artifact: constraint mutation","What happens to your decision when one live constraint changes?",mutation,
           ["CHANGE ONE CONSTRAINT — do not rewrite the whole problem.","RERUN — apply the same reasoning sequence under the new condition.","COMPARE — state what survived and what changed."],
           "POST-CLASS ARTIFACT — mutate one constraint in your Decision Card and revise the verdict.","A defensible design survives re-analysis when constraints move.",(grade,),"constraint-card")

    assign(18,"Post-class artifact: Engineering Defense rehearsal","Can you defend your architecture when the examiner changes the constraint live?",mutation,
           ["DEFEND — state the decision and its evidence.","CHALLENGE — accept one changed constraint.","RESPOND — revise only what the new constraint invalidates.","RESIDUAL UNCERTAINTY — state what still needs verification."],
           "POST-CLASS ARTIFACT — rehearse a 60-second Engineering Defense.","The course assesses reasoning that remains coherent under challenge.",(grade,),"defense")

    assign(19,"Post-class artifact: two-question peer review","Can another student inspect your evidence and falsify your claim?",decision_card,
           ["PEER-REVIEW CARD — (1) Can another person independently inspect the evidence? (2) What variable or edge case would invalidate the claim?","Use the full 6×4 rubric after class; do not run all rubric cells live in the room."],
           "POST-CLASS ARTIFACT — peer-review one Decision Card with the two-question matrix.","Fast peer review tests inspectability and falsifiability before full grading.",(grade,),"peer-review")

    logistics_core=_compact(log_s,0,12)
    assign(20,"Logistics & next steps","What must you do before the next class, and where do you go if you need help?",logistics_core,
           ["ACTION TODAY — review the syllabus and configure the lab environment.","TAKE-HOME CHECKPOINT — verify the next class/lab commitment from the published logistics, then arrive ready to engineer."],
           "CHECKPOINT — confirm your next action before leaving.","Leave Week 1 with the schedule, contact path and next action explicit.",(logistics,),"checklist")

    # Update first-taught evidence so the ledger follows the actual learner story.
    mapping={}
    for unit in bp.units:
        for ev in unit.coverage_evidence:
            mapping.setdefault(ev.coverage_id,unit.number)
    for entry in bp.coverage_ledger:
        if entry.coverage_id in mapping:
            entry.first_taught_unit=mapping[entry.coverage_id]
            entry.reinforced_units=[n for n in (16,18,19) if n!=entry.first_taught_unit]
            entry.representation=f"Unit {entry.first_taught_unit}: source-faithful kickoff teaching move"
    family_first={
        "iscarb engineering flow":6,"micro-case":5,"professional standards":8,"ai literacy":8,
        "earning your grade":9,"roadmap":10,"milestones":10,"logistics":20,"next steps":20,
    }
    for tc in bp.topic_coverage:
        low=tc.topic_family.lower()
        tc.first_taught_unit=next((n for k,n in family_first.items() if k in low),min(tc.first_taught_unit,15))
        tc.reinforced_units=[n for n in (11,12,15,16,18,20) if n!=tc.first_taught_unit][:4]

    bp.release_notes=[x for x in bp.release_notes if "gulf.edu.sa" not in str(x).lower()]
    return _strip_retired_readiness(fit_presenter_text(bp))


def _strip_retired_readiness(bp):
    for u in list(getattr(bp,"units",[]) or []):
        cleaned=[]
        for x in list(getattr(u,"pedagogy_content",[]) or []):
            low=str(x).lower()
            if "gulf.edu.sa/standardized-exams-readiness" in low:
                continue
            if low.startswith("etec readiness target — unverified"):
                cleaned.append("READINESS CHECK — assert a course/program outcome mapping only when an approved authoritative mapping is supplied.")
                continue
            cleaned.append(x)
        u.pedagogy_content=cleaned
        u.enrichment_content=[x for x in (u.enrichment_content or []) if "gulf.edu.sa/standardized-exams-readiness" not in str(x).lower()]
        u.enrichment_basis=[x for x in (u.enrichment_basis or []) if "gulf.edu.sa/standardized-exams-readiness" not in str(x).lower()]
    return bp


def apply_v723_kickoff_patch(app):
    global _PATCHED
    if _PATCHED:
        return
    _PATCHED=True

    # Exact PPTX tables must be visible to both source profiling and hard gates.
    old_extract=source_text_mod.extract_source_text
    def extract_v723(path,limit=600_000):
        p=Path(path)
        if p.suffix.lower()==".pptx":
            try:return _pptx_text(p,limit)
            except Exception:return old_extract(p,limit)
        return old_extract(p,limit)
    source_text_mod.extract_source_text=extract_v723
    source_bundle_mod.extract_source_text=extract_v723
    profile_fallback._pptx_chunks=_pptx_chunks

    # Enrich deterministic and reconciled profiles before blueprint construction.
    prev_profile=engine.build_deterministic_source_profile
    def build_profile(bundle,*args,**kwargs):
        return _ensure_kickoff_profile(prev_profile(bundle,*args,**kwargs),bundle)
    engine.build_deterministic_source_profile=build_profile
    base.engine.build_deterministic_source_profile=build_profile

    prev_reconcile=engine.reconcile_source_profile
    def reconcile(profile,bundle):
        return _ensure_kickoff_profile(prev_reconcile(profile,bundle),bundle)
    engine.reconcile_source_profile=reconcile
    base.engine.reconcile_source_profile=reconcile

    # Apply orientation storytelling to both source-only fallback and model output.
    prev_draft=engine._source_preserving_draft
    def kickoff_draft(profile,bundle):
        return _specialize_kickoff(prev_draft(profile,bundle),profile,bundle)
    engine._source_preserving_draft=kickoff_draft
    base.engine._source_preserving_draft=kickoff_draft

    prev_timebox=engine.apply_90_minute_timebox
    def kickoff_timebox(bp,profile,bundle):
        return _specialize_kickoff(prev_timebox(bp,profile,bundle),profile,bundle)
    engine.apply_90_minute_timebox=kickoff_timebox
    base.engine.apply_90_minute_timebox=kickoff_timebox

    # Final public release identity.
    root_route=next((r for r in app.router.routes if getattr(r,"path",None)=="/"),None)
    old_root=getattr(root_route,"endpoint",None)
    if root_route is not None: app.router.routes.remove(root_route)
    if old_root is not None:
        @app.get("/")
        def home_v723():
            response=old_root(); body=response.body.decode("utf-8") if hasattr(response,"body") else str(response)
            body=body.replace("7.2.2","7.2.3")
            headers=dict(getattr(response,"headers",{}) or {})
            headers.update({"X-ISCARB-Version":"7.2.3","X-ISCARB-UI":"7.2.3","Cache-Control":"no-store, no-cache, must-revalidate, max-age=0"})
            return HTMLResponse(body,headers=headers)

    health_route=next((r for r in app.router.routes if getattr(r,"path",None)=="/api/health"),None)
    old_health=getattr(health_route,"endpoint",None)
    if health_route is not None: app.router.routes.remove(health_route)
    if old_health is not None:
        @app.get("/api/health")
        def health_v723():
            data=dict(old_health())
            data.update({
                "version":"7.2.3","release_ui":"7.2.3",
                "pipeline":"iscarb-v7.2.3-final-clean-it-wide-kickoff-aware",
                "pptx_table_extraction":True,
                "kickoff_orientation_mode":True,
                "kickoff_preserves_assessment_roadmap_logistics":True,
                "retired_readiness_link_removed":True,
                "domain_spine_kickoff_curated":True,
            })
            return data
