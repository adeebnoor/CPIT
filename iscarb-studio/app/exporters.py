from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet

from .models import Blueprint


# -----------------------------------------------------------------------------
# Shared helpers
# -----------------------------------------------------------------------------

def _readiness_for_unit(bp: Blueprint, unit_no: int) -> list[str]:
    refs: list[str] = []
    for r in bp.readiness_alignment:
        if unit_no in r.evidence_units or unit_no in {3, 16, 19, 20}:
            label = f"{r.sku}: {', '.join(r.slo_refs)} → {', '.join(r.klo_refs)}"
            if label not in refs:
                refs.append(label)
    return refs


def _short(text: str, limit: int = 130) -> str:
    text = " ".join((text or "").split())
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 1)].rstrip() + "…"


# -----------------------------------------------------------------------------
# Detailed instructor DOCX
# -----------------------------------------------------------------------------

def export_docx(bp: Blueprint, out: Path) -> Path:
    doc = Document()
    doc.add_heading(bp.lecture_title, 0)
    doc.add_paragraph(bp.engineering_thesis)
    doc.add_paragraph(f"Live session: {bp.session_minutes} minutes | 20 Units")
    doc.add_paragraph(f"Named ethical purpose: {bp.named_ethical_purpose}")

    doc.add_heading("Lecture Source Bundle", level=1)
    for s in bp.source_manifest:
        doc.add_paragraph(s, style="List Bullet")

    doc.add_heading("Central Engineering Crisis", level=1)
    doc.add_paragraph(bp.central_engineering_crisis)

    doc.add_heading("Weekly CLOs", level=1)
    for c in bp.clOs:
        doc.add_paragraph(f"{c.id}: {c.statement}", style="List Bullet")
        doc.add_paragraph(f"Evidence: {c.evidence_expected}")

    doc.add_heading("ETEC IT Readiness Alignment", level=1)
    for r in bp.readiness_alignment:
        doc.add_paragraph(
            f"{r.gku} | {r.sku} | {', '.join(r.slo_refs)} | {', '.join(r.klo_refs)} | {r.strength}",
            style="List Bullet",
        )
        doc.add_paragraph(f"Rationale: {r.rationale}")
        doc.add_paragraph(f"Atomicity evidence: {r.atomicity_evidence}")
        doc.add_paragraph(
            f"Evidence units: {', '.join(str(x) for x in r.evidence_units)} | "
            f"ETEC pages: {', '.join(str(x) for x in r.standard_source_pages)}"
        )

    for u in bp.units:
        doc.add_page_break()
        doc.add_heading(f"UNIT {u.number:02d} — {u.phase} — {u.title} — {u.planned_minutes} min", level=1)
        doc.add_paragraph(f"Engineering question: {u.engineering_question}")
        if u.core_content:
            doc.add_heading("Lecture-bundle technical content", level=2)
            for bullet in u.core_content:
                doc.add_paragraph(bullet, style="List Bullet")
        if u.pedagogy_content:
            doc.add_heading("ISCARB Pedagogy / Decision Work", level=2)
            for bullet in u.pedagogy_content:
                doc.add_paragraph(bullet, style="List Bullet")
        if u.enrichment_content:
            doc.add_heading("ISCARB Contextual Enrichment", level=2)
            for bullet in u.enrichment_content:
                doc.add_paragraph(bullet, style="List Bullet")
            for basis in u.enrichment_basis:
                doc.add_paragraph(f"Basis: {basis}")
        if u.scenario_assumptions:
            doc.add_heading("Scenario assumptions", level=2)
            for assumption in u.scenario_assumptions:
                doc.add_paragraph(assumption, style="List Bullet")
        refs = _readiness_for_unit(bp, u.number)
        if refs:
            doc.add_heading("Readiness trace", level=2)
            for ref in refs:
                doc.add_paragraph(ref, style="List Bullet")
        doc.add_paragraph(f"Visual suggestion: {u.visual_suggestion}")
        doc.add_paragraph(f"Student action: {u.student_action}")
        doc.add_paragraph(f"Takeaway: {u.takeaway}")
        doc.add_paragraph(f"CIMT: {', '.join(u.cimtlens)} | CLO: {', '.join(u.clo_ids)}")
        doc.add_paragraph(f"Source anchor: {u.source_anchor or 'N/A — ISCARB pedagogy'}")
        doc.add_paragraph(f"Evidence: {u.evidence}")

    doc.add_page_break()
    doc.add_heading("Four-Level Engineering Rubric", level=1)
    for r in bp.rubric_criteria:
        doc.add_heading(r.criterion, level=2)
        doc.add_paragraph(f"4 — Distinguished: {r.distinguished}")
        doc.add_paragraph(f"3 — Ready: {r.ready}")
        doc.add_paragraph(f"2 — Developing: {r.developing}")
        doc.add_paragraph(f"1 — Not Yet Ready: {r.not_yet_ready}")
        if r.readiness_refs:
            doc.add_paragraph(f"Readiness: {', '.join(r.readiness_refs)}")
    doc.save(out)
    return out


# -----------------------------------------------------------------------------
# Visual presenter PowerPoint — 20 Units, 20 distinct teaching stops
# -----------------------------------------------------------------------------

BG = RGBColor(247, 248, 250)
INK = RGBColor(20, 29, 38)
MUTED = RGBColor(91, 103, 115)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(220, 225, 231)
DARK = RGBColor(14, 25, 37)
PHASE = {
    "IFHAM": RGBColor(37, 99, 235),
    "MARIS": RGBColor(13, 143, 104),
    "ATQAN": RGBColor(171, 109, 16),
    "MAYYIZ": RGBColor(116, 70, 232),
}
SOFT = {
    "IFHAM": RGBColor(235, 242, 255),
    "MARIS": RGBColor(232, 248, 241),
    "ATQAN": RGBColor(252, 244, 227),
    "MAYYIZ": RGBColor(242, 237, 255),
}
GREEN = RGBColor(15, 118, 85)
RED = RGBColor(181, 55, 55)
AMBER = RGBColor(181, 116, 24)
BLUE = RGBColor(37, 99, 235)


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _line(shape, color: RGBColor = LINE, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def _box(slide, x, y, w, h, *, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(shape, fill)
    _line(shape, line)
    return shape


def _text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = str(text or "")
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tb


def _label(slide, x, y, w, text, color, *, fill=None):
    fill = fill or color
    shape = _box(slide, x, y, w, 0.36, fill=fill, line=fill)
    _text(slide, x, y + 0.01, w, 0.31, text.upper(), size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return shape


def _slide_base(prs: Presentation, bp: Blueprint, u):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    accent = PHASE[u.phase]

    # top rule + phase/unit marker
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.09))
    _fill(rule, accent); rule.line.fill.background()
    _label(slide, 0.62, 0.35, 1.15, u.phase, accent)
    _text(slide, 1.93, 0.34, 2.15, 0.36, f"UNIT {u.number:02d} · {u.planned_minutes} MIN", size=9, bold=True, color=MUTED, valign=MSO_ANCHOR.MIDDLE)
    _text(slide, 0.62, 0.82, 12.05, 0.64, u.title, size=25, bold=True, color=INK)
    _text(slide, 0.62, 1.48, 12.0, 0.74, u.engineering_question, size=15.5, bold=True, color=accent)
    return slide, accent


def _footer(slide, bp: Blueprint, u, accent):
    y = 7.10
    _text(slide, 0.62, y, 8.1, 0.22, _short(u.source_anchor or "ISCARB pedagogy", 115), size=7.3, color=MUTED)
    refs = _readiness_for_unit(bp, u.number)
    if refs:
        _text(slide, 8.55, y, 3.55, 0.22, _short(refs[0], 62), size=7.2, bold=True, color=accent, align=PP_ALIGN.RIGHT)
    _text(slide, 12.22, y, 0.48, 0.22, f"{u.number}/20", size=7.5, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)


def _action_bar(slide, u, accent):
    _box(slide, 0.62, 6.37, 12.05, 0.58, fill=SOFT[u.phase], line=SOFT[u.phase])
    _text(slide, 0.83, 6.47, 1.2, 0.22, "YOU TRY", size=8.5, bold=True, color=accent)
    _text(slide, 1.75, 6.43, 10.6, 0.31, _short(u.student_action, 150), size=10.5, bold=True, color=INK)


def _card(slide, x, y, w, h, title, body, *, accent=BLUE, fill=WHITE, number=None):
    _box(slide, x, y, w, h, fill=fill, line=LINE)
    if number is not None:
        _box(slide, x + 0.18, y + 0.18, 0.42, 0.42, fill=accent, line=accent)
        _text(slide, x + 0.18, y + 0.195, 0.42, 0.34, str(number), size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        tx = x + 0.72
        tw = w - 0.9
    else:
        tx = x + 0.2
        tw = w - 0.4
    _text(slide, tx, y + 0.18, tw, 0.32, title, size=10.5, bold=True, color=accent)
    _text(slide, x + 0.2, y + 0.62, w - 0.4, h - 0.78, _short(body, 160), size=9.2, color=INK)


def _layout_crisis(slide, bp, u, accent):
    _box(slide, 0.62, 2.25, 7.2, 3.75, fill=DARK, line=DARK)
    _text(slide, 0.95, 2.55, 6.55, 0.32, "INCIDENT / INCOMPLETE EVIDENCE", size=9, bold=True, color=RGBColor(157, 177, 196))
    _text(slide, 0.95, 3.02, 6.45, 1.1, _short(bp.central_engineering_crisis, 260), size=19, bold=True, color=WHITE)
    _text(slide, 0.95, 4.65, 6.35, 0.75, "Professional purpose", size=9, bold=True, color=RGBColor(157, 177, 196))
    _text(slide, 0.95, 5.02, 6.4, 0.65, _short(bp.named_ethical_purpose, 150), size=12.5, color=WHITE)
    items = (u.core_content + u.scenario_assumptions + u.pedagogy_content)[:3]
    for i, item in enumerate(items):
        _card(slide, 8.12, 2.25 + i * 1.18, 4.55, 1.0, f"Signal {i+1}", item, accent=accent, number=i+1)


def _layout_spine(slide, bp, u, accent):
    fams = bp.source_topic_families[:6]
    for i, fam in enumerate(fams):
        col, row = i % 3, i // 3
        x = 0.72 + col * 4.05
        y = 2.42 + row * 1.58
        _card(slide, x, y, 3.65, 1.25, f"{i+1:02d}", fam, accent=accent, fill=WHITE)
    _text(slide, 0.78, 5.72, 11.7, 0.42, _short(u.takeaway, 175), size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)


def _layout_clos(slide, bp, u, accent):
    positions = [(0.68, 2.35, 3.85), (4.73, 2.35, 3.85), (8.78, 2.35, 3.85), (2.7, 4.3, 3.85), (6.75, 4.3, 3.85)]
    for i, c in enumerate(bp.clOs[:5]):
        x, y, w = positions[i]
        _card(slide, x, y, w, 1.55, c.id, c.statement, accent=accent, fill=SOFT[u.phase], number=i+1)


def _layout_hstack(slide, bp, u, accent):
    competencies = [
        ("Analytical Reasoning", "Decompose the system and isolate failure mechanisms."),
        ("Engineering Judgment", "Choose under competing constraints and imperfect evidence."),
        ("Evidence-Based Reasoning", "Defend claims with inspectable evidence."),
        ("Socio-Technical Thinking", "Account for people, policy, context, and operations."),
        ("Risk-Aware Design", "Design for failure, uncertainty, and residual risk."),
        ("Ethical Responsibility", "Own the consequences of engineering decisions."),
    ]
    for i, (title, body) in enumerate(competencies):
        col, row = i % 3, i // 3
        _card(slide, 0.7 + col * 4.05, 2.35 + row * 1.68, 3.65, 1.34, title, body, accent=accent, fill=WHITE, number=i+1)


def _layout_predict(slide, bp, u, accent):
    labels = ["PREDICT", "CONSTRAINT", "DERIVE", "NAME"]
    bodies = [
        _short(u.student_action, 105),
        _short((u.scenario_assumptions or u.core_content)[0] if (u.scenario_assumptions or u.core_content) else "What limits the choice?", 105),
        _short(u.pedagogy_content[0] if u.pedagogy_content else "Derive the mechanism from the constraint.", 105),
        _short(u.takeaway, 105),
    ]
    for i, label in enumerate(labels):
        x = 0.68 + i * 3.05
        _card(slide, x, 2.65, 2.7, 2.45, label, bodies[i], accent=accent, fill=WHITE, number=i+1)
        if i < 3:
            _text(slide, x + 2.72, 3.58, 0.34, 0.35, "→", size=21, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _text(slide, 1.05, 5.45, 11.2, 0.44, "Prediction must happen before the model name is revealed.", size=11.5, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _layout_mechanism(slide, bp, u, accent):
    labels = ["INPUT", "MECHANISM", "OUTPUT", "ASSUMPTION", "FAILURE MODE"]
    source = u.core_content + u.pedagogy_content
    for i, label in enumerate(labels):
        x = 0.62 + i * 2.45
        body = source[i] if i < len(source) else u.takeaway
        _card(slide, x, 2.72, 2.15, 2.2, label, body, accent=accent, fill=WHITE, number=i+1)
        if i < 4:
            _text(slide, x + 2.15, 3.55, 0.3, 0.3, "→", size=18, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _text(slide, 0.78, 5.35, 11.9, 0.48, _short(u.takeaway, 180), size=11.2, bold=True, color=INK, align=PP_ALIGN.CENTER)


def _layout_layers(slide, bp, u, accent):
    layers = ["PLATFORM / OUTER", "APPLICATION / MIDDLE", "RECORD / CORE"]
    widths = [10.8, 8.5, 6.1]
    xs = [1.25, 2.4, 3.6]
    ys = [2.4, 3.25, 4.1]
    fills = [SOFT[u.phase], WHITE, SOFT[u.phase]]
    for i in range(3):
        _box(slide, xs[i], ys[i], widths[i], 1.05, fill=fills[i], line=accent)
        body = u.core_content[i] if i < len(u.core_content) else u.takeaway
        _text(slide, xs[i] + 0.25, ys[i] + 0.18, widths[i] - 0.5, 0.25, layers[i], size=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
        _text(slide, xs[i] + 0.35, ys[i] + 0.5, widths[i] - 0.7, 0.32, _short(body, 100), size=9, color=INK, align=PP_ALIGN.CENTER)
    _text(slide, 0.85, 5.52, 11.7, 0.45, "Defense-in-depth: compromise of one layer must not imply compromise of the asset.", size=11.2, bold=True, color=INK, align=PP_ALIGN.CENTER)


def _layout_tradeoff(slide, bp, u, accent):
    left = "\n".join(_short(x, 78) for x in (u.core_content[:2] or ["Alternative A"])); right = "\n".join(_short(x, 78) for x in (u.core_content[2:4] or ["Alternative B"]))
    _card(slide, 0.72, 2.45, 5.35, 2.65, "ALTERNATIVE A", left, accent=accent, fill=WHITE)
    _card(slide, 7.26, 2.45, 5.35, 2.65, "ALTERNATIVE B", right, accent=accent, fill=WHITE)
    _box(slide, 5.93, 3.0, 1.45, 1.55, fill=accent, line=accent)
    _text(slide, 5.93, 3.22, 1.45, 0.42, "VS", size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _text(slide, 5.93, 3.72, 1.45, 0.32, "TRADE-OFF", size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    matrix = "COST  ·  PERFORMANCE  ·  RESILIENCE  ·  COMPLEXITY"
    _text(slide, 1.1, 5.45, 11.1, 0.35, matrix, size=10.5, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _layout_falsification(slide, bp, u, accent):
    cards = [
        ("CLAIM", u.takeaway),
        ("MEASURE", u.core_content[0] if u.core_content else u.evidence),
        ("BREAK IT", next((x for x in u.pedagogy_content if "fals" in x.lower() or "invalid" in x.lower()), u.student_action)),
    ]
    for i, (title, body) in enumerate(cards):
        _card(slide, 0.72 + i * 4.12, 2.55, 3.7, 2.55, title, body, accent=accent, fill=WHITE, number=i+1)
    _text(slide, 1.0, 5.4, 11.3, 0.45, "Passing a test is not the same as supporting an engineering claim.", size=12, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _layout_uncertainty(slide, bp, u, accent):
    labels = ["KNOWN", "UNKNOWN", "DECISION-SENSITIVE UNKNOWN", "WHAT WE MONITOR"]
    for i, label in enumerate(labels):
        col, row = i % 2, i // 2
        x, y = 0.72 + col * 6.02, 2.35 + row * 1.75
        body = next((b for b in u.pedagogy_content if label.lower().split()[0] in b.lower()), u.takeaway)
        _card(slide, x, y, 5.55, 1.43, label, body, accent=accent, fill=WHITE, number=i+1)


def _layout_context(slide, bp, u, accent):
    _box(slide, 0.72, 2.4, 4.1, 3.15, fill=SOFT[u.phase], line=accent)
    _text(slide, 1.03, 2.73, 3.45, 0.34, "SAUDI CONTEXT", size=11, bold=True, color=accent, align=PP_ALIGN.CENTER)
    scenario = (u.scenario_assumptions + u.enrichment_content + u.pedagogy_content)
    _text(slide, 1.03, 3.28, 3.45, 1.55, _short(scenario[0] if scenario else "Context changes the engineering decision.", 190), size=13, bold=True, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _text(slide, 5.25, 2.48, 1.35, 2.8, "→", size=34, bold=True, color=accent, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _box(slide, 6.65, 2.4, 6.0, 3.15, fill=WHITE, line=LINE)
    _text(slide, 6.98, 2.72, 5.3, 0.34, "HOW DOES THE CONTEXT CHANGE THE DESIGN?", size=10, bold=True, color=accent)
    for i, item in enumerate((u.core_content + u.pedagogy_content)[-3:]):
        _text(slide, 7.0, 3.3 + i * 0.62, 5.05, 0.5, f"{i+1}. {_short(item, 92)}", size=10.2, color=INK)


def _layout_accountability(slide, bp, u, accent):
    labels = ["SYSTEM EVENT", "AUDIT EVIDENCE", "RESPONSIBILITY", "AMANAH / DUTY"]
    source = u.core_content + u.pedagogy_content
    for i, label in enumerate(labels):
        x = 0.72 + i * 3.03
        body = source[i] if i < len(source) else u.takeaway
        _card(slide, x, 2.7, 2.68, 2.2, label, body, accent=accent, fill=WHITE, number=i+1)
        if i < 3:
            _text(slide, x + 2.68, 3.55, 0.35, 0.3, "→", size=18, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _text(slide, 1.0, 5.38, 11.3, 0.45, _short(u.takeaway, 170), size=11.3, bold=True, color=INK, align=PP_ALIGN.CENTER)


def _layout_trend(slide, bp, u, accent):
    cards = [
        ("ENDURING PRINCIPLE", u.core_content[0] if u.core_content else u.takeaway),
        ("CURRENT PRACTICE", u.enrichment_content[0] if u.enrichment_content else "How is the principle implemented today?"),
        ("WATCH NEXT", u.student_action),
    ]
    for i, (title, body) in enumerate(cards):
        _card(slide, 0.72 + i * 4.12, 2.5, 3.7, 2.7, title, body, accent=accent, fill=WHITE, number=i+1)
        if i < 2:
            _text(slide, 4.34 + i * 4.12, 3.55, 0.45, 0.3, "→", size=19, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _layout_wellbeing(slide, bp, u, accent):
    cards = [
        ("DESIGN LOAD", "What creates cognitive/operational burden?"),
        ("AUTOMATE / CLARIFY", "Which mechanism removes avoidable toil or ambiguity?"),
        ("RESILIENT PRACTICE", "How does the team respond safely under pressure?"),
    ]
    bodies = u.pedagogy_content + u.core_content + [u.takeaway]
    for i, (title, default) in enumerate(cards):
        _card(slide, 0.72 + i * 4.12, 2.55, 3.7, 2.6, title, bodies[i] if i < len(bodies) else default, accent=accent, fill=WHITE, number=i+1)
    _text(slide, 1.0, 5.45, 11.25, 0.4, "Reliable engineering should reduce alert fatigue, ambiguity, and recovery chaos—not merely add controls.", size=11.1, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _layout_ai(slide, bp, u, accent):
    _card(slide, 0.72, 2.45, 5.55, 2.25, "AI MAY ASSIST", next((x for x in u.pedagogy_content if "ai may assist" in x.lower()), "Generate alternatives, draft checks, and brainstorm test cases."), accent=GREEN, fill=RGBColor(238, 249, 244))
    _card(slide, 6.72, 2.45, 5.55, 2.25, "AI MUST NOT BE TRUSTED AUTONOMOUSLY", next((x for x in u.pedagogy_content if "ai must not" in x.lower()), "Do not delegate critical architecture approval or unverifiable security claims."), accent=RED, fill=RGBColor(252, 240, 240))
    pipeline = ["CLAIM", "ASSUMPTION", "SOURCE", "TEST", "FAILURE SEARCH", "HUMAN SIGN-OFF"]
    for i, label in enumerate(pipeline):
        x = 0.76 + i * 2.0
        _box(slide, x, 5.05, 1.62, 0.72, fill=SOFT[u.phase], line=accent)
        _text(slide, x + 0.08, 5.22, 1.46, 0.33, label, size=7.8, bold=True, color=accent, align=PP_ALIGN.CENTER)
        if i < len(pipeline)-1:
            _text(slide, x + 1.62, 5.23, 0.34, 0.25, "→", size=13, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _layout_portfolio(slide, bp, u, accent):
    items = [
        "Problem framing", "Threat / risk ledger", "Architecture", "Trade-offs",
        "Deployment / operations", "Evidence + assurance"
    ]
    for i, item in enumerate(items):
        col, row = i % 3, i // 3
        _card(slide, 0.72 + col * 4.05, 2.35 + row * 1.55, 3.65, 1.25, item, "Professional artifact evidence", accent=accent, fill=WHITE, number=i+1)
    refs = _readiness_for_unit(bp, u.number)
    if refs:
        _text(slide, 1.0, 5.65, 11.3, 0.32, "Readiness trace: " + _short(" | ".join(refs), 150), size=9.8, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _layout_mutation(slide, bp, u, accent):
    before = _short(u.takeaway, 100)
    mutation = _short((u.scenario_assumptions or u.pedagogy_content)[0] if (u.scenario_assumptions or u.pedagogy_content) else "Constraint changes", 115)
    after = _short(u.student_action, 115)
    cards = [("BEFORE", before), ("MUTATION", mutation), ("ADAPT + CRITIQUE", after)]
    for i, (title, body) in enumerate(cards):
        _card(slide, 0.72 + i * 4.12, 2.55, 3.7, 2.65, title, body, accent=accent, fill=WHITE, number=i+1)
        if i < 2:
            _text(slide, 4.34 + i * 4.12, 3.55, 0.45, 0.3, "→", size=19, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _layout_evidence(slide, bp, u, accent):
    labels = ["CLAIM", "EVIDENCE", "WARRANT", "COUNTER-EVIDENCE", "RESIDUAL UNCERTAINTY"]
    source = u.pedagogy_content
    for i, label in enumerate(labels):
        x = 0.62 + i * 2.45
        body = next((b for b in source if label.lower().split("-")[0] in b.lower()), u.takeaway)
        _card(slide, x, 2.65, 2.15, 2.55, label, body, accent=accent, fill=WHITE, number=i+1)
        if i < 4:
            _text(slide, x + 2.15, 3.65, 0.3, 0.3, "→", size=18, bold=True, color=accent, align=PP_ALIGN.CENTER)


def _descriptor_keyword(text: str) -> str:
    part = (text or "").split(";")[0].split(".")[0]
    return _short(part, 54)


def _layout_rubric(slide, bp, u, accent):
    criteria = bp.rubric_criteria[:6]
    x0, y0 = 0.62, 2.32
    widths = [3.15, 2.2, 2.2, 2.2, 2.2]
    headers = ["CRITERION", "4 · DISTINGUISHED", "3 · READY", "2 · DEVELOPING", "1 · NOT YET"]
    x = x0
    for i, h in enumerate(headers):
        _box(slide, x, y0, widths[i], 0.55, fill=DARK if i == 0 else accent, line=DARK if i == 0 else accent, radius=False)
        _text(slide, x + 0.06, y0 + 0.12, widths[i]-0.12, 0.28, h, size=7.2, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += widths[i]
    for r, crit in enumerate(criteria):
        y = y0 + 0.55 + r * 0.56
        x = x0
        vals = [crit.criterion, _descriptor_keyword(crit.distinguished), _descriptor_keyword(crit.ready), _descriptor_keyword(crit.developing), _descriptor_keyword(crit.not_yet_ready)]
        for c, val in enumerate(vals):
            fill = WHITE if r % 2 == 0 else RGBColor(250, 251, 252)
            _box(slide, x, y, widths[c], 0.56, fill=fill, line=LINE, radius=False)
            _text(slide, x + 0.06, y + 0.08, widths[c]-0.12, 0.4, val, size=6.7 if c else 7.2, bold=(c==0), color=INK, align=PP_ALIGN.LEFT if c==0 else PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            x += widths[c]


def _layout_assurance(slide, bp, u, accent):
    _box(slide, 1.1, 2.25, 11.1, 0.92, fill=DARK, line=DARK)
    top_claim = next((x for x in u.pedagogy_content if "claim" in x.lower()), u.takeaway)
    _text(slide, 1.35, 2.47, 10.6, 0.44, _short(top_claim, 150), size=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, c in enumerate(bp.clOs[:5]):
        x = 0.7 + i * 2.52
        _card(slide, x, 3.55, 2.25, 1.65, c.id, c.statement, accent=accent, fill=WHITE, number=i+1)
    decisions = ["APPROVE", "CONDITIONAL", "REDESIGN", "REJECT"]
    for i, d in enumerate(decisions):
        color = [GREEN, AMBER, accent, RED][i]
        _box(slide, 2.05 + i * 2.35, 5.48, 2.02, 0.62, fill=color, line=color)
        _text(slide, 2.05 + i * 2.35, 5.65, 2.02, 0.26, d, size=8.6, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def _layout_generic(slide, bp, u, accent):
    items = u.core_content + u.pedagogy_content + u.enrichment_content
    for i, item in enumerate(items[:6]):
        col, row = i % 2, i // 2
        _card(slide, 0.72 + col * 6.0, 2.35 + row * 1.25, 5.55, 1.0, f"POINT {i+1}", item, accent=accent, fill=WHITE, number=i+1)


def export_pptx(bp: Blueprint, out: Path) -> Path:
    """Export a presenter-first visual deck.

    The content model remains source-locked; this renderer changes only visual grammar.
    Exactly one slide is created per ISCARB Unit (20 slides total).
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layouts = {
        1: _layout_crisis,
        2: _layout_spine,
        3: _layout_clos,
        4: _layout_hstack,
        5: _layout_predict,
        6: _layout_mechanism,
        7: _layout_layers,
        8: _layout_tradeoff,
        9: _layout_falsification,
        10: _layout_uncertainty,
        11: _layout_context,
        12: _layout_accountability,
        13: _layout_trend,
        14: _layout_wellbeing,
        15: _layout_ai,
        16: _layout_portfolio,
        17: _layout_mutation,
        18: _layout_evidence,
        19: _layout_rubric,
        20: _layout_assurance,
    }

    for u in bp.units:
        slide, accent = _slide_base(prs, bp, u)
        renderer = layouts.get(u.number, _layout_generic)
        renderer(slide, bp, u, accent)
        _action_bar(slide, u, accent)
        _footer(slide, bp, u, accent)

    prs.save(out)
    return out


# -----------------------------------------------------------------------------
# Detailed PDF — source-grounded reference/handout
# -----------------------------------------------------------------------------

def export_pdf(bp: Blueprint, out: Path) -> Path:
    styles = getSampleStyleSheet()
    story = [
        Paragraph(bp.lecture_title, styles["Title"]),
        Paragraph(bp.engineering_thesis, styles["BodyText"]),
        Paragraph(f"<b>Live session:</b> {bp.session_minutes} minutes | 20 Units", styles["BodyText"]),
        Paragraph(f"<b>Named ethical purpose:</b> {bp.named_ethical_purpose}", styles["BodyText"]),
        Spacer(1, 8),
        Paragraph("<b>LECTURE SOURCE BUNDLE</b>", styles["BodyText"]),
    ]
    for source in bp.source_manifest:
        story.append(Paragraph("• " + source, styles["BodyText"]))
    story.append(PageBreak())

    for u in bp.units:
        story.append(Paragraph(f"UNIT {u.number:02d} — {u.phase} — {u.title} — {u.planned_minutes} min", styles["Heading1"]))
        story.append(Paragraph(f"<b>Engineering question:</b> {u.engineering_question}", styles["BodyText"]))
        if u.core_content:
            story.append(Paragraph("<b>LECTURE-BUNDLE TECHNICAL CONTENT</b>", styles["BodyText"]))
            for bullet in u.core_content:
                story.append(Paragraph("• " + bullet, styles["BodyText"]))
        if u.pedagogy_content:
            story.append(Paragraph("<b>ISCARB PEDAGOGY / DECISION WORK</b>", styles["BodyText"]))
            for bullet in u.pedagogy_content:
                story.append(Paragraph("• " + bullet, styles["BodyText"]))
        if u.enrichment_content:
            story.append(Paragraph("<b>ISCARB CONTEXTUAL ENRICHMENT</b>", styles["BodyText"]))
            for bullet in u.enrichment_content:
                story.append(Paragraph("• " + bullet, styles["BodyText"]))
            story.append(Paragraph(f"<b>Basis:</b> {' | '.join(u.enrichment_basis)}", styles["BodyText"]))
        if u.scenario_assumptions:
            story.append(Paragraph(f"<b>Scenario assumptions:</b> {' | '.join(u.scenario_assumptions)}", styles["BodyText"]))
        refs = _readiness_for_unit(bp, u.number)
        if refs:
            story.append(Paragraph(f"<b>ETEC readiness:</b> {' | '.join(refs)}", styles["BodyText"]))
        if u.number == 19:
            story.append(Paragraph("<b>FOUR-LEVEL RUBRIC MATRIX</b>", styles["BodyText"]))
            for r in bp.rubric_criteria:
                story.append(Paragraph(f"<b>{r.criterion}</b> — 4: {r.distinguished} | 3: {r.ready} | 2: {r.developing} | 1: {r.not_yet_ready}", styles["BodyText"]))
        story.append(Paragraph(f"<b>Student action:</b> {u.student_action}", styles["BodyText"]))
        story.append(Paragraph(f"<b>Takeaway:</b> {u.takeaway}", styles["BodyText"]))
        story.append(Paragraph(f"<b>Source:</b> {u.source_anchor or 'N/A — ISCARB pedagogy'}", styles["BodyText"]))
        story.append(PageBreak())
    SimpleDocTemplate(str(out), pagesize=A4).build(story)
    return out
