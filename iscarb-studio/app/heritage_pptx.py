from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from .models import Blueprint
from .visual_engine import export_presenter_pptx

WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(8, 117, 84)
DARK_GREEN = RGBColor(18, 62, 49)
GOLD = RGBColor(201, 163, 61)
MUTED = RGBColor(101, 118, 111)
LIGHT = RGBColor(235, 241, 238)


def _set_text_color(shape, color: RGBColor) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color


def export_cimt_heritage_pptx(blueprint: Blueprint, path: Path) -> Path:
    """Render the normal ISCARB visual grammar, then apply CIMT heritage DNA.

    The post-process is deliberately conservative: it does not change content or
    geometry. It restores the white academic canvas, green editorial titles,
    thin gold rules and a lighter footer while preserving phase cues and each
    Unit's purpose-built visual structure.
    """
    path = Path(path)
    base = path.with_name(path.stem + "__visual_base.pptx")
    export_presenter_pptx(blueprint, base)
    prs = Presentation(str(base))

    for slide in prs.slides:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = WHITE

        # Existing visual shapes stay intact; only editorial chrome is refined.
        for shape in slide.shapes:
            top = shape.top / Inches(1)
            left = shape.left / Inches(1)
            height = shape.height / Inches(1)

            # Main editorial title region.
            if getattr(shape, "has_text_frame", False) and 0.70 <= top <= 1.35 and height <= 1.15:
                _set_text_color(shape, GREEN)

            # Footer bar created by Visual Grammar v1.
            if top >= 7.0 and not getattr(shape, "has_text_frame", False):
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = WHITE
                    shape.line.fill.background()
                except Exception:
                    pass
            if getattr(shape, "has_text_frame", False) and top >= 7.0:
                _set_text_color(shape, MUTED if left < 9.3 else GREEN)

        # CIMT-style thin rules: visual identity without competing with content.
        for y in (0.16, 6.98):
            rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.56), Inches(y), Inches(12.18), Inches(0.025))
            rule.fill.solid(); rule.fill.fore_color.rgb = GOLD; rule.line.fill.background()

        # Quiet lineage marker, intentionally smaller than the teaching content.
        marker = slide.shapes.add_textbox(Inches(10.45), Inches(0.19), Inches(2.25), Inches(0.25))
        tf = marker.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = "ISCARB · CIMT HERITAGE"
        p.font.name = "Aptos"; p.font.size = Pt(7.5); p.font.bold = True; p.font.color.rgb = DARK_GREEN

    prs.save(str(path))
    try:
        base.unlink(missing_ok=True)
    except Exception:
        pass
    return path
