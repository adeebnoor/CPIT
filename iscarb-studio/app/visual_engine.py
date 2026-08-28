from __future__ import annotations

import html
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .models import Blueprint, LectureUnit

# ISCARB Visual Grammar v1 — deterministic/local, no model calls.
INK = RGBColor(21, 31, 42)
MUTED = RGBColor(92, 107, 120)
PAPER = RGBColor(247, 249, 251)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(221, 227, 232)
BLUE = RGBColor(53, 104, 232)
GREEN = RGBColor(22, 133, 107)
AMBER = RGBColor(183, 123, 29)
VIOLET = RGBColor(121, 85, 220)
RED = RGBColor(182, 64, 64)
SOFT_BLUE = RGBColor(235, 241, 255)
SOFT_GREEN = RGBColor(234, 248, 242)
SOFT_AMBER = RGBColor(255, 245, 226)
SOFT_VIOLET = RGBColor(243, 239, 255)
SOFT_RED = RGBColor(255, 239, 239)

PHASE_COLOR = {"IFHAM": BLUE, "MARIS": GREEN, "ATQAN": AMBER, "MAYYIZ": VIOLET}
PHASE_SOFT = {"IFHAM": SOFT_BLUE, "MARIS": SOFT_GREEN, "ATQAN": SOFT_AMBER, "MAYYIZ": SOFT_VIOLET}


def _short(text: str, n: int = 145) -> str:
    text = re.sub(r"\s+", " ", str(text or "")).strip()
    return text if len(text) <= n else text[: n - 1].rstrip() + "…"


def _strip_prefix(text: str) -> str:
    text = re.sub(r"^\s*[•\-–—]\s*", "", text or "")
    text = re.sub(r"^\s*\d+[\.:)]\s*", "", text)
    return text.strip()


def _find_line(unit: LectureUnit, keyword: str, fallback: str = "") -> str:
    k = keyword.lower()
    for item in [*unit.pedagogy_content, *unit.core_content, *unit.scenario_assumptions]:
        if k in item.lower():
            value = re.sub(r"^.*?[:\-]\s*", "", item, count=1)
            return _short(value or item)
    return fallback


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_bg(slide, color=PAPER):
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = color


def _text(slide, x, y, w, h, text, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = str(text or "")
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return shape


def _box(slide, x, y, w, h, title, body="", fill=WHITE, line=LINE, title_color=INK, body_color=MUTED, title_size=13, body_size=10.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(1)
    tf = sh.text_frame; tf.clear(); tf.word_wrap = True; tf.margin_left = Inches(.14); tf.margin_right = Inches(.14); tf.margin_top = Inches(.10); tf.margin_bottom = Inches(.08)
    p = tf.paragraphs[0]; p.text = _short(title, 70); p.font.name = "Aptos"; p.font.size = Pt(title_size); p.font.bold = True; p.font.color.rgb = title_color
    if body:
        p = tf.add_paragraph(); p.text = _short(body, 180); p.font.name = "Aptos"; p.font.size = Pt(body_size); p.font.color.rgb = body_color; p.space_before = Pt(4)
    return sh


def _pill(slide, x, y, w, text, fill, fg=WHITE, size=9.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(.34))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    tf = sh.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER; p.font.name = "Aptos"; p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = fg
    return sh


def _arrow(slide, x, y, w=.38, h=.24, color=MUTED):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background(); return sh


def _base(slide, unit: LectureUnit):
    _fill_bg(slide)
    c = PHASE_COLOR[unit.phase]
    _pill(slide, .55, .32, 1.05, f"UNIT {unit.number:02d}", c)
    _pill(slide, 1.72, .32, 1.05, unit.phase, c)
    _pill(slide, 11.78, .32, .95, f"{unit.planned_minutes} MIN", c)
    _text(slide, .58, .83, 12.0, .55, unit.title, size=24, bold=True)
    _text(slide, .60, 1.43, 12.05, .55, unit.engineering_question, size=12.2, color=MUTED, bold=True)
    # footer: action + source
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.0), Inches(7.08), Inches(13.333), Inches(.42))
    sh.fill.solid(); sh.fill.fore_color.rgb = INK; sh.line.fill.background()
    _text(slide, .58, 7.12, 8.9, .25, "YOU TRY · " + _short(unit.student_action, 105), size=8.5, color=WHITE, bold=True)
    _text(slide, 9.55, 7.12, 3.15, .25, _short(unit.source_anchor or "ISCARB pedagogy", 58), size=7.7, color=RGBColor(195,205,215), align=PP_ALIGN.RIGHT)


def _generic_cards(slide, unit: LectureUnit, labels: list[str] | None = None, max_cards=6):
    items = [_strip_prefix(x) for x in (unit.core_content or unit.pedagogy_content)[:max_cards]]
    if labels:
        cards = [(labels[i], items[i] if i < len(items) else "") for i in range(min(len(labels), max_cards))]
    else:
        cards = [(f"{i+1:02d}", item) for i, item in enumerate(items)]
    cols = 3 if len(cards) > 4 else 2
    rows = (len(cards) + cols - 1) // cols
    w = 3.82 if cols == 3 else 5.9; h = 1.35 if rows >= 2 else 2.4
    for i, (title, body) in enumerate(cards):
        r, cidx = divmod(i, cols)
        x = .62 + cidx * (w + .34); y = 2.16 + r * (h + .28)
        _box(slide, x, y, w, h, title, body, fill=PHASE_SOFT[unit.phase], line=PHASE_COLOR[unit.phase])


def _render_01(slide, bp, u):
    signals = [_strip_prefix(x) for x in u.core_content[:3]]
    _box(slide, .62, 2.12, 4.1, 3.75, "INCIDENT", bp.central_engineering_crisis, fill=SOFT_RED, line=RED, title_color=RED, title_size=16, body_size=12)
    for i, s in enumerate(signals):
        _box(slide, 5.05, 2.12 + i*1.18, 3.55, .95, f"SIGNAL {i+1}", s, fill=WHITE, line=LINE)
    _box(slide, 8.92, 2.12, 3.76, 3.75, "DECISION UNDER UNCERTAINTY", "Diagnose the system without revealing or assuming the root cause. What evidence would change your first decision?", fill=SOFT_BLUE, line=BLUE, title_color=BLUE, title_size=15, body_size=12)


def _render_02(slide, bp, u):
    fams = bp.source_topic_families[:6]
    for i, name in enumerate(fams):
        r, c = divmod(i, 3); x = .62 + c*4.08; y = 2.15 + r*1.75
        _box(slide, x, y, 3.75, 1.42, f"FAMILY {i+1}", name, fill=SOFT_BLUE if r==0 else WHITE, line=BLUE if r==0 else LINE, title_color=BLUE)
        if i < len(fams)-1 and c < 2: _arrow(slide, x+3.8, y+.6, .23, .18, BLUE)
    _box(slide, .62, 5.86, 12.06, .77, "DOMAIN SPINE", "Complete primary coverage: every major P1 family remains inside the 90-minute lecture; depth varies, coverage does not.", fill=INK, line=INK, title_color=WHITE, body_color=RGBColor(220,228,235), title_size=11, body_size=9.5)


def _render_03(slide, bp, u):
    clos = bp.clOs[:5]
    w = 2.28
    for i,c in enumerate(clos):
        _box(slide, .58+i*2.52, 2.15, w, 3.85, c.id, c.statement, fill=WHITE, line=BLUE, title_color=BLUE, title_size=15, body_size=11)
        _pill(slide, .73+i*2.52, 5.42, 1.98, "PROOF REQUIRED", BLUE, size=8)
    _text(slide, .62, 6.18, 12.0, .44, "CLOs are promises only when each ends in observable evidence.", size=12.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)


def _render_04(slide, bp, u):
    names = ["ANALYTICAL", "JUDGMENT", "EVIDENCE", "SOCIO-TECH", "RISK-AWARE", "ETHICAL"]
    bodies = [_strip_prefix(x) for x in u.pedagogy_content if ":" in x][:6]
    for i,n in enumerate(names):
        r,c=divmod(i,3); _box(slide,.62+c*4.08,2.15+r*1.82,3.75,1.5,n,bodies[i] if i<len(bodies) else "",fill=PHASE_SOFT[u.phase],line=PHASE_COLOR[u.phase],title_color=PHASE_COLOR[u.phase])
    _box(slide, 4.26, 5.92, 4.85, .7, "H-STACK", "Technical correctness + human consequences + evidence + ownership.", fill=INK, line=INK, title_color=WHITE, body_color=WHITE, title_size=11, body_size=9)


def _render_05(slide, bp, u):
    stages=["PREDICT","CONSTRAINT","DERIVE","NAME"]
    bodies=[_find_line(u,"predict",u.student_action),_find_line(u,"constraint","What cannot be violated?"),_find_line(u,"deriv","What principle follows?"),_find_line(u,"principle",u.takeaway)]
    for i,s in enumerate(stages):
        x=.62+i*3.02; _box(slide,x,2.35,2.62,2.85,s,bodies[i],fill=[SOFT_BLUE,SOFT_AMBER,SOFT_GREEN,SOFT_VIOLET][i],line=[BLUE,AMBER,GREEN,VIOLET][i],title_color=[BLUE,AMBER,GREEN,VIOLET][i],title_size=15,body_size=11)
        if i<3:_arrow(slide,x+2.69,3.6,.26,.22,INK)
    _box(slide,.62,5.55,11.95,.83,"ORDER MATTERS","Prediction must occur before explanatory content; otherwise the learner is recalling, not reasoning.",fill=INK,line=INK,title_color=WHITE,body_color=WHITE,title_size=10.5,body_size=9.5)


def _render_06(slide,bp,u):
    labels=["INPUT","MECHANISM","OUTPUT","ASSUMPTION","FAILURE"]
    items=[_strip_prefix(x) for x in [*u.pedagogy_content,*u.core_content][:5]]
    for i,l in enumerate(labels):
        x=.48+i*2.56; _box(slide,x,2.45,2.15,2.85,l,items[i] if i<len(items) else "",fill=WHITE,line=PHASE_COLOR[u.phase],title_color=PHASE_COLOR[u.phase],title_size=12.5,body_size=10)
        if i<4:_arrow(slide,x+2.18,3.68,.3,.2,PHASE_COLOR[u.phase])
    _text(slide,.62,5.68,12,.55,"A mechanism is understood only when its assumption and failure mode are explicit.",size=13,color=MUTED,bold=True,align=PP_ALIGN.CENTER)


def _render_07(slide,bp,u):
    items=[_strip_prefix(x) for x in u.core_content[:3]]
    labels=["OUTER LAYER","APPLICATION LAYER","ASSET / RECORD LAYER"]
    widths=[10.9,8.2,5.5]; xs=[1.22,2.57,3.92]; ys=[2.2,3.18,4.16]
    fills=[SOFT_BLUE,SOFT_GREEN,SOFT_VIOLET]; lines=[BLUE,GREEN,VIOLET]
    for i in range(3): _box(slide,xs[i],ys[i],widths[i],.82,labels[i],items[i] if i<len(items) else "",fill=fills[i],line=lines[i],title_color=lines[i],title_size=11,body_size=8.5)
    _box(slide,3.22,5.45,6.9,.9,"DESIGN TEST","If one layer is compromised, what prevents direct asset access?",fill=INK,line=INK,title_color=WHITE,body_color=WHITE,title_size=11,body_size=9.5)


def _render_08(slide,bp,u):
    assumptions=" | ".join(u.scenario_assumptions[:2])
    _box(slide,.62,2.15,4.6,3.72,"ALTERNATIVE A",_find_line(u,"alternative a",u.core_content[0] if u.core_content else "Centralized design"),fill=SOFT_BLUE,line=BLUE,title_color=BLUE,title_size=16,body_size=12)
    _box(slide,8.08,2.15,4.6,3.72,"ALTERNATIVE B",_find_line(u,"alternative b",u.core_content[1] if len(u.core_content)>1 else "Distributed design"),fill=SOFT_GREEN,line=GREEN,title_color=GREEN,title_size=16,body_size=12)
    _box(slide,5.5,2.48,2.3,2.95,"TRADE-OFF","Cost\nUsability\nResilience\nComplexity",fill=WHITE,line=AMBER,title_color=AMBER,title_size=13,body_size=11)
    if assumptions:_text(slide,.68,6.05,11.9,.44,"SYNTHETIC EXERCISE DATA · "+_short(assumptions,150),size=9,color=MUTED,align=PP_ALIGN.CENTER)


def _render_09(slide,bp,u):
    cards=[("CLAIM",u.takeaway),("MEASURE",u.evidence or u.student_action),("FALSIFY",_find_line(u,"falsif","What observation would make us abandon the decision?"))]
    colors=[BLUE,GREEN,RED]; fills=[SOFT_BLUE,SOFT_GREEN,SOFT_RED]
    for i,(t,b) in enumerate(cards):
        x=.62+i*4.08; _box(slide,x,2.3,3.75,3.55,t,b,fill=fills[i],line=colors[i],title_color=colors[i],title_size=16,body_size=12)
        if i<2:_arrow(slide,x+3.78,3.86,.25,.2,INK)
    _text(slide,.62,6.15,12,.4,"Passing a test is not the same as supporting an engineering claim.",size=12.5,color=MUTED,bold=True,align=PP_ALIGN.CENTER)


def _render_10(slide,bp,u):
    labels=["KNOWN","UNKNOWN","DECISION-SENSITIVE UNKNOWN","WHAT WE MONITOR"]
    keys=["known:","unknown:","decision-sensitive unknown","what we monitor"]
    fills=[SOFT_GREEN,SOFT_RED,SOFT_AMBER,SOFT_BLUE]; lines=[GREEN,RED,AMBER,BLUE]
    for i,(lab,key) in enumerate(zip(labels,keys)):
        r,c=divmod(i,2); body=_find_line(u,key,"Make the uncertainty explicit.")
        _box(slide,.72+c*6.08,2.18+r*2.12,5.72,1.78,lab,body,fill=fills[i],line=lines[i],title_color=lines[i],title_size=13,body_size=10.5)
    _pill(slide,4.55,6.37,4.25,"SENIOR DESIGN REVIEW",INK,size=10)


def _render_11(slide,bp,u):
    context=" | ".join(u.scenario_assumptions[:2]) or _find_line(u,"saudi","Hypothetical Saudi professional context")
    _box(slide,.62,2.2,4.2,3.95,"SAUDI CONTEXT",context,fill=SOFT_GREEN,line=GREEN,title_color=GREEN,title_size=15,body_size=12)
    impacts=[_strip_prefix(x) for x in [*u.pedagogy_content,*u.core_content][:3]]
    for i,b in enumerate(impacts): _box(slide,5.18,2.2+i*1.28,7.5,1.02,f"DESIGN IMPACT {i+1}",b,fill=WHITE,line=LINE,title_color=INK,title_size=10.5,body_size=9.5)
    _text(slide,5.2,6.18,7.35,.4,"Context must change a technical decision—not decorate it.",size=11.5,color=MUTED,bold=True,align=PP_ALIGN.CENTER)


def _render_12(slide,bp,u):
    labels=["EVENT","EVIDENCE","RESPONSIBILITY","AMANAH / OWNERSHIP"]
    bodies=[u.core_content[0] if u.core_content else "System action",u.evidence or "Traceable evidence",u.pedagogy_content[0] if u.pedagogy_content else "Named responsible role",u.takeaway]
    colors=[BLUE,GREEN,AMBER,VIOLET]; fills=[SOFT_BLUE,SOFT_GREEN,SOFT_AMBER,SOFT_VIOLET]
    for i,l in enumerate(labels):
        x=.62+i*3.02; _box(slide,x,2.55,2.62,2.82,l,bodies[i],fill=fills[i],line=colors[i],title_color=colors[i],title_size=11.5,body_size=10)
        if i<3:_arrow(slide,x+2.68,3.78,.27,.2,INK)


def _render_13(slide,bp,u):
    cards=[("ENDURING PRINCIPLE",u.core_content[0] if u.core_content else u.takeaway),("CURRENT PRACTICE",u.enrichment_content[0] if u.enrichment_content else "Apply the source principle in contemporary practice."),("NEXT QUESTION",u.student_action)]
    colors=[INK,BLUE,VIOLET]; fills=[WHITE,SOFT_BLUE,SOFT_VIOLET]
    for i,(t,b) in enumerate(cards):
        x=.65+i*4.04; _box(slide,x,2.35,3.7,3.55,t,b,fill=fills[i],line=colors[i],title_color=colors[i],title_size=13,body_size=11)
        if i<2:_arrow(slide,x+3.73,3.88,.24,.2,MUTED)
    _text(slide,.62,6.18,12,.38,"Trend = what changes; principle = what must still be true.",size=12,color=MUTED,bold=True,align=PP_ALIGN.CENTER)


def _render_14(slide,bp,u):
    labels=["DESIGN LOAD","ALERT / TASK LOAD","AUTOMATION / RUNBOOK","PRACTITIONER RESILIENCE"]
    bodies=[u.core_content[0] if u.core_content else "Operational mechanism",u.pedagogy_content[0] if u.pedagogy_content else "Cognitive demand",u.pedagogy_content[1] if len(u.pedagogy_content)>1 else "Reduce repetitive work",u.takeaway]
    for i,l in enumerate(labels):
        x=.62+i*3.02; _box(slide,x,2.55,2.62,2.65,l,bodies[i],fill=[SOFT_RED,SOFT_AMBER,SOFT_BLUE,SOFT_GREEN][i],line=[RED,AMBER,BLUE,GREEN][i],title_color=[RED,AMBER,BLUE,GREEN][i],title_size=11,body_size=9.7)
        if i<3:_arrow(slide,x+2.68,3.72,.27,.2,INK)
    _box(slide,3.55,5.6,6.35,.85,"ENGINEERING WELLBEING","Reduce avoidable cognitive load through architecture, observability, automation and clear recovery paths.",fill=INK,line=INK,title_color=WHITE,body_color=WHITE,title_size=10.5,body_size=9)


def _render_15(slide,bp,u):
    may=_find_line(u,"ai may assist","Draft alternatives, test ideas, summarize supplied evidence.")
    must=_find_line(u,"ai must not","Autonomously approve architecture or critical claims.")
    _box(slide,.62,2.15,5.72,2.25,"AI MAY ASSIST",may,fill=SOFT_GREEN,line=GREEN,title_color=GREEN,title_size=16,body_size=12)
    _box(slide,6.72,2.15,5.96,2.25,"AI MUST NOT BE TRUSTED AUTONOMOUSLY",must,fill=SOFT_RED,line=RED,title_color=RED,title_size=13,body_size=11.5)
    steps=["CLAIM","ASSUMPTION","SOURCE CHECK","TEST","FAILURE SEARCH","HUMAN SIGN-OFF"]
    for i,s in enumerate(steps):
        x=.62+i*2.02; _pill(slide,x,5.0,1.78,s,INK,size=7.7)
        if i<5:_arrow(slide,x+1.8,5.08,.17,.15,MUTED)
    _text(slide,.62,5.75,12,.55,u.student_action,size=11,color=MUTED,bold=True,align=PP_ALIGN.CENTER)


def _render_16(slide,bp,u):
    labels=["PROBLEM","THREAT / ANALYSIS","ARCHITECTURE","TRADE-OFF","EVIDENCE","ASSURANCE"]
    bodies=[u.pedagogy_content[i] if i<len(u.pedagogy_content) else "Required portfolio evidence" for i in range(6)]
    for i,l in enumerate(labels):
        r,c=divmod(i,3); _box(slide,.62+c*4.08,2.18+r*1.78,3.75,1.45,l,bodies[i],fill=WHITE,line=VIOLET,title_color=VIOLET,title_size=11,body_size=9.5)
    _pill(slide,4.42,5.95,4.5,"PROOF OF CAPABILITY",VIOLET,size=10)


def _render_17(slide,bp,u):
    items=[("BEFORE","Current architecture"),("MUTATION"," | ".join(u.scenario_assumptions[:2]) or "Constraint changes"),("ADAPT","Use only mechanisms already taught"),("PEER CRITIQUE",u.student_action)]
    colors=[BLUE,RED,GREEN,VIOLET]; fills=[SOFT_BLUE,SOFT_RED,SOFT_GREEN,SOFT_VIOLET]
    for i,(t,b) in enumerate(items):
        x=.62+i*3.02; _box(slide,x,2.45,2.62,3.05,t,b,fill=fills[i],line=colors[i],title_color=colors[i],title_size=12,body_size=10.5)
        if i<3:_arrow(slide,x+2.68,3.78,.27,.2,INK)
    _text(slide,.62,5.9,12,.48,"Robustness is demonstrated when the design survives a changed constraint without importing unvetted technology.",size=11.7,color=MUTED,bold=True,align=PP_ALIGN.CENTER)


def _render_18(slide,bp,u):
    labels=["CLAIM","EVIDENCE","WARRANT","COUNTER-EVIDENCE","RESIDUAL UNCERTAINTY"]
    keys=["claim","evidence","warrant","counter","residual"]
    colors=[BLUE,GREEN,AMBER,RED,VIOLET]
    fills=[SOFT_BLUE,SOFT_GREEN,SOFT_AMBER,SOFT_RED,SOFT_VIOLET]
    for i,(l,k) in enumerate(zip(labels,keys)):
        x=.42+i*2.55; body=_find_line(u,k,u.pedagogy_content[i] if i<len(u.pedagogy_content) else "")
        _box(slide,x,2.35,2.25,3.45,l,body,fill=fills[i],line=colors[i],title_color=colors[i],title_size=10.5,body_size=9.3)
        if i<4:_arrow(slide,x+2.28,3.88,.22,.18,INK)
    _text(slide,.62,6.18,12,.35,"Evidence earns confidence; it does not erase uncertainty.",size=12,color=MUTED,bold=True,align=PP_ALIGN.CENTER)


def _render_19(slide,bp,u):
    levels=[("4","DISTINGUISHED",GREEN),("3","READY",BLUE),("2","DEVELOPING",AMBER),("1","NOT YET",RED)]
    for i,(n,name,c) in enumerate(levels): _pill(slide,6.7+i*1.46,2.05,1.28,f"{n} · {name}",c,size=6.8)
    criteria=bp.rubric_criteria[:6]
    for r,crit in enumerate(criteria):
        y=2.55+r*.59
        _text(slide,.62,y,5.7,.42,_short(crit.criterion,62),size=9.5,bold=True)
        vals=[crit.distinguished,crit.ready,crit.developing,crit.not_yet_ready]
        for i,val in enumerate(vals):
            sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.7+i*1.46),Inches(y),Inches(1.28),Inches(.4))
            sh.fill.solid(); sh.fill.fore_color.rgb=[SOFT_GREEN,SOFT_BLUE,SOFT_AMBER,SOFT_RED][i]; sh.line.color.rgb=[GREEN,BLUE,AMBER,RED][i]
            tf=sh.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=_short(val,38); p.font.size=Pt(5.8); p.font.color.rgb=INK
    _text(slide,.62,6.2,12,.4,"The rubric measures engineering capability—not slide polish or keyword recall.",size=11.5,color=MUTED,bold=True,align=PP_ALIGN.CENTER)


def _render_20(slide,bp,u):
    _box(slide,2.25,2.05,8.85,.9,"TOP CLAIM",_find_line(u,"top claim",u.takeaway),fill=INK,line=INK,title_color=WHITE,body_color=WHITE,title_size=11,body_size=9.5)
    for i,c in enumerate(bp.clOs[:5]):
        x=.42+i*2.55; _box(slide,x,3.3,2.25,1.7,c.id,_short(c.evidence_expected,130),fill=PHASE_SOFT[u.phase],line=VIOLET,title_color=VIOLET,title_size=12,body_size=8.7)
    decisions=[("APPROVE",GREEN),("CONDITIONAL",BLUE),("REDESIGN",AMBER),("REJECT",RED)]
    for i,(d,c) in enumerate(decisions): _pill(slide,2.0+i*2.45,5.52,2.15,d,c,size=9.5)
    _text(slide,.62,6.17,12,.38,"Bounded authorization = evidence-proportionate permission to operate, never proof of absolute security.",size=11,color=MUTED,bold=True,align=PP_ALIGN.CENTER)


_RENDERERS={1:_render_01,2:_render_02,3:_render_03,4:_render_04,5:_render_05,6:_render_06,7:_render_07,8:_render_08,9:_render_09,10:_render_10,11:_render_11,12:_render_12,13:_render_13,14:_render_14,15:_render_15,16:_render_16,17:_render_17,18:_render_18,19:_render_19,20:_render_20}


def visual_type(unit_no: int) -> str:
    return {
        1:"crisis board",2:"domain spine",3:"CLO cards",4:"H-Stack grid",5:"derive sequence",6:"mechanism flow",7:"layered architecture",8:"trade-off matrix",9:"falsification chain",10:"uncertainty board",11:"Saudi context decision",12:"accountability chain",13:"trend timeline",14:"wellbeing loop",15:"AI permission gate",16:"portfolio canvas",17:"constraint mutation",18:"evidence chain",19:"rubric matrix",20:"assurance decision gate",
    }.get(unit_no,"visual cards")


def export_presenter_pptx(bp: Blueprint, out: Path) -> Path:
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    for u in bp.units:
        slide=_blank(prs); _base(slide,u)
        renderer=_RENDERERS.get(u.number)
        if renderer: renderer(slide,bp,u)
        else: _generic_cards(slide,u)
    prs.save(out); return out


def _preview_payload(bp: Blueprint) -> list[dict]:
    return [{
        "n":u.number,"phase":u.phase,"title":u.title,"q":u.engineering_question,"minutes":u.planned_minutes,
        "type":visual_type(u.number),"core":u.core_content[:4],"ped":u.pedagogy_content[:4],"enrich":u.enrichment_content[:2],
        "assume":u.scenario_assumptions[:2],"action":u.student_action,"takeaway":u.takeaway,"source":u.source_anchor,
    } for u in bp.units]


def render_presenter_preview(bp: Blueprint, release_state: str = "BLOCKED") -> str:
    data=json.dumps(_preview_payload(bp),ensure_ascii=False).replace("</","<\\/")
    title=html.escape(bp.lecture_title)
    thesis=html.escape(bp.engineering_thesis)
    css="""
    :root{--ink:#17202a;--muted:#6b7783;--paper:#eef2f5;--blue:#3568e8;--green:#16856b;--amber:#b77b1d;--violet:#7955dc;--red:#b64040}*{box-sizing:border-box}body{margin:0;font-family:Inter,system-ui,sans-serif;background:#0d141d;color:white}.deck{height:100vh;display:grid;grid-template-columns:260px 1fr}.rail{padding:18px;background:#111b26;border-right:1px solid #273444;overflow:auto}.rail h1{font-size:16px;margin:4px 0}.rail p{font-size:11px;color:#9fb0c0;line-height:1.4}.thumb{display:block;width:100%;text-align:left;background:#172330;color:#cdd7e0;border:1px solid #263647;border-radius:10px;padding:9px 10px;margin:7px 0;cursor:pointer}.thumb.active{border-color:#6f8ff7;background:#1c2d45}.thumb b{font-size:10px}.thumb span{display:block;font-size:9px;color:#90a3b5;margin-top:3px}.stage{display:grid;place-items:center;padding:28px}.slide{width:min(1180px,calc(100vw - 330px));aspect-ratio:16/9;background:#f7f9fb;color:var(--ink);border-radius:18px;box-shadow:0 30px 90px rgba(0,0,0,.32);padding:32px 40px;display:grid;grid-template-rows:auto auto 1fr auto}.meta{display:flex;gap:8px;align-items:center;font-size:10px;font-weight:800}.pill{padding:5px 9px;border-radius:999px;color:white}.title{font-size:30px;font-weight:850;letter-spacing:-.03em;margin:14px 0 5px}.q{font-size:15px;color:var(--muted);font-weight:700;max-width:1000px}.visual{align-self:stretch;display:grid;gap:12px;margin-top:20px}.cards2{grid-template-columns:1fr 1fr}.cards3{grid-template-columns:repeat(3,1fr)}.cards4{grid-template-columns:repeat(4,1fr)}.cards5{grid-template-columns:repeat(5,1fr)}.card{background:white;border:1px solid #dce3e8;border-radius:13px;padding:14px;min-height:90px}.card b{display:block;font-size:12px;margin-bottom:6px}.card p{font-size:11px;color:#5d6a75;line-height:1.38;margin:0}.accent{border-left:5px solid var(--blue)}.footer{background:#17202a;color:white;margin:18px -40px -32px;padding:10px 40px;font-size:10px;font-weight:700;display:flex;justify-content:space-between}.controls{position:fixed;right:28px;bottom:22px;display:flex;gap:8px}.controls button{background:#fff;color:#17202a;border:0;border-radius:999px;padding:9px 13px;font-weight:800;cursor:pointer}@media(max-width:900px){.deck{grid-template-columns:1fr}.rail{display:none}.slide{width:95vw}.stage{padding:10px}.cards5,.cards4,.cards3{grid-template-columns:1fr 1fr}}
    """
    script="""
    const units=PAYLOAD; const colors={IFHAM:'#3568e8',MARIS:'#16856b',ATQAN:'#b77b1d',MAYYIZ:'#7955dc'}; let idx=0;
    const esc=s=>String(s??'').replace(/[&<>\"']/g,c=>({'&':'&amp;','<':'&lt;','>':'&gt;','\"':'&quot;',"'":'&#39;'}[c]));
    function cards(u){let arr=[]; if(u.n===3){arr=u.ped.slice(0,5)} else if(u.n===10){arr=u.ped.slice(0,4)} else if(u.n===15){arr=u.ped.slice(0,4)} else if(u.n>=16){arr=u.ped.slice(0,6)} else arr=(u.core.length?u.core:u.ped).slice(0,6); if(u.assume.length&&arr.length<5)arr=arr.concat(u.assume.slice(0,1)); const n=Math.max(2,Math.min(5,arr.length||2)); return `<div class="visual cards${n}">${arr.map((x,i)=>`<div class="card accent"><b>${esc(u.type.toUpperCase())} · ${String(i+1).padStart(2,'0')}</b><p>${esc(x)}</p></div>`).join('')}</div>`}
    function render(){const u=units[idx],c=colors[u.phase]; document.getElementById('main').innerHTML=`<div class="slide"><div class="meta"><span class="pill" style="background:${c}">UNIT ${String(u.n).padStart(2,'0')}</span><span class="pill" style="background:${c}">${u.phase}</span><span style="color:#71808d">${u.minutes} MIN · ${esc(u.type)}</span></div><div><div class="title">${esc(u.title)}</div><div class="q">${esc(u.q)}</div></div>${cards(u)}<div class="footer"><span>YOU TRY · ${esc(u.action)}</span><span>${esc(u.source||'ISCARB pedagogy')}</span></div></div>`; document.querySelectorAll('.thumb').forEach((b,i)=>b.classList.toggle('active',i===idx));}
    function go(n){idx=(n+units.length)%units.length; render()}; document.addEventListener('keydown',e=>{if(e.key==='ArrowRight'||e.key==='PageDown')go(idx+1); if(e.key==='ArrowLeft'||e.key==='PageUp')go(idx-1)});
    window.addEventListener('DOMContentLoaded',()=>{document.getElementById('railList').innerHTML=units.map((u,i)=>`<button class="thumb" onclick="go(${i})"><b>${String(u.n).padStart(2,'0')} · ${u.phase}</b><span>${esc(u.title)}</span></button>`).join(''); render();});
    """.replace("PAYLOAD",data)
    return f"""<!doctype html><html><head><meta charset='utf-8'><meta name='viewport' content='width=device-width,initial-scale=1'><title>{title} · Presenter</title><style>{css}</style></head><body><div class='deck'><aside class='rail'><div style='font-size:10px;color:#8fa1b2;font-weight:900;letter-spacing:.12em'>ISCARB PRESENTER</div><h1>{title}</h1><p>{thesis}</p><p>Release state: <b>{html.escape(release_state)}</b><br>Arrow keys move through 20 slides.</p><div id='railList'></div></aside><main class='stage' id='main'></main></div><div class='controls'><button onclick='go(idx-1)'>← Previous</button><button onclick='go(idx+1)'>Next →</button></div><script>{script}</script></body></html>"""
