from __future__ import annotations

"""Deterministic P1 source profiling used only when the AI profiling call is unavailable.

The fallback is deliberately conservative: it does not invent subject matter. It
extracts page/slide/section headings and short source excerpts from the PRIMARY
file and turns those into a coverage contract. Gemini still receives the full
source during blueprint generation; this module only prevents the whole compile
from dying before generation when a profiling-model quota is exhausted.
"""

import re
from pathlib import Path

from .models import CoverageItem, SourceProfile, TopicFamily
from .source_bundle import SourceBundle
from .source_text import extract_source_text


_FURNITURE = {
    "software engineering", "chapter", "contents", "copyright", "all rights reserved",
    "ian sommerville", "page", "slide", "learning objectives", "objectives",
}


# PDF and PPTX bullet furniture that survives text extraction. Left in place it
# reaches the faculty deck as a literal glyph in a slide heading.
_BULLET_EDGE = " \t\r\n-|\u2022\u00b7\u25a0\u25aa\u25c6\u25cf\u25b6\u25fc\u25fe\u2023\u2043\u00a7\u2192*_"

# A contact block is authorship metadata, never lecture content. Matching the
# address itself keeps this rule source-agnostic instead of name-specific.
_CONTACT_RE = re.compile(r"[\w.+-]+@[\w-]+\.[\w.]+")


# "SLIDE 73: ..." is how the SlideShare extractor numbers what it scraped. It is
# scaffolding, not a lecture heading, and must not survive into a unit title.
_EXTRACTOR_PREFIX = re.compile(r"^\s*(?:slide|page|frame)\s*\d{1,4}\s*[:.\u2013\u2014-]\s*", re.I)
# U+00B7 is deliberately absent: it is the separator this module joins excerpt
# lines with, and stripping it collapsed every source slide into one run-on
# line, which silently broke both importance scoring and furniture detection.
_INLINE_BULLET = re.compile(r"[\u2022\u25a0\u25aa\u25c6\u25cf\u25b6\u25fc\u25fe\u2023\u2043]+")


def _display_label(text: str) -> str:
    """Strip the extractor's slide coordinate from text a human will read.

    Applied only after the coordinate has been captured into the anchor - the
    "SLIDE 73:" prefix is how the anchor is derived, so removing it any earlier
    silently downgrades every anchor to a section ordinal.
    """
    return _EXTRACTOR_PREFIX.sub("", str(text or "")).strip(_BULLET_EDGE).lstrip(":;,. ")


def _clean(text: str, limit: int = 160) -> str:
    text = re.sub(r"\s+", " ", str(text or ""))
    # Bullet glyphs survive extraction mid-string too, not only at the edges.
    text = _INLINE_BULLET.sub(" ", text)
    text = re.sub(r"\s+", " ", text).strip(_BULLET_EDGE)
    if len(text) > limit:
        # Cut on a word boundary so a heading never ends mid-word.
        head = text[:limit]
        cut = head.rfind(" ")
        text = head[:cut] if cut >= limit * 0.6 else head
    return text.rstrip(" ,;:-")


def _meaningful_lines(text: str) -> list[str]:
    lines: list[str] = []
    for raw in str(text or "").splitlines():
        line = _clean(raw, 4000)
        if len(line) < 3:
            continue
        low = line.lower()
        if low.isdigit() or re.fullmatch(r"\d+[./-]?\d*", low):
            continue
        if any(low == x for x in _FURNITURE):
            continue
        if line not in lines:
            lines.append(line)
    return lines


def _knowledge_type(text: str) -> str:
    low = text.lower()
    if re.match(r"(?:example\b|case study\b|ex\s*[;:.])", low.strip()):
        return "EXAMPLE"
    if any(x in low for x in ("algorithm", "pseudo", "procedure", "steps:")):
        return "ALGORITHM"
    if any(x in low for x in ("equation", "formula", "probability", "pofod", "rocof", "mtbf", "mttf", "=", "%")):
        return "EQUATION"
    if any(x in low for x in ("architecture", "layer", "component", "subsystem", "client-server", "distributed")):
        return "ARCHITECTURE"
    if any(x in low for x in ("process", "activities", "workflow", "review", "inspection", "testing")):
        return "PROCESS"
    if any(x in low for x in ("protocol", "request", "response", "message", "transaction")):
        return "PROTOCOL"
    if any(x in low for x in ("trade-off", "tradeoff", "cost", "versus", " vs ", "advantage", "disadvantage")):
        return "TRADE_OFF"
    if any(x in low for x in ("example", "case study", "scenario")) or re.search(r"\bex\s*[;:.]", low):
        return "EXAMPLE"
    if any(x in low for x in ("principle", "guideline", "rule", "design for")):
        return "DESIGN_PRINCIPLE"
    if any(x in low for x in ("failure", "fault", "behavior", "behaviour", "state", "recovery")):
        return "SYSTEM_BEHAVIOR"
    return "CONCEPT"


def _is_furniture_line(line: str) -> bool:
    low = _clean(line, 180).lower()
    if not low:
        return True
    exact = {
        "cpit-455 software engineering (ii)", "advanced software engineering",
        "advanced softwre engineering", "today lecture topic", "today lecture topic – part 2",
        "today lecture topic - part 2", "the cimt compass", "where academic meets business reality",
        "concept implementation measurement trend", "measurement implementation reality",
    }
    if low in exact:
        return True
    if any(low.startswith(x) for x in ("adeeb noor", "it department", "faculty of computing", "king abdulaziz university", "fall 2025")):
        return True
    # An instructor contact block is authorship metadata, not teachable content,
    # and must never become a lecture title or a slide heading.
    if _CONTACT_RE.search(low):
        return True
    return False


# A wrapped body line ends where the page ran out of width, so it trails off on
# a function word and cannot serve as a heading.
_DANGLING_TAIL = re.compile(
    r"\b(of|for|to|the|a|an|and|or|but|in|on|at|by|with|from|that|which|is|are|"
    r"was|were|be|been|as|it|its|their|this|these|those|than|then|so|if|when)$",
    re.IGNORECASE,
)


def _is_title_like(line: str) -> bool:
    """True when the line reads as a heading rather than a wrapped body line."""
    text = str(line or "").strip()
    if not (2 <= len(text.split()) <= 11) or len(text) > 100:
        return False
    if text.endswith((",", ";", ":", "-", "\u2013", "\u2014")):
        return False
    if _DANGLING_TAIL.search(text.rstrip(".")):
        return False
    first = text.split()[0]
    # A continuation of the sentence above starts lowercase; a heading does not.
    return not (first[:1].islower() and first.lower() not in {"e.g.", "i.e."})


# Ten teaching units need distinct material. Below this the deck is the same
# checkpoint retold, which the old profile still reported as scope_fit="FIT".
MIN_MAJOR_CHECKPOINTS_FOR_A_LECTURE = 3


def _thin_source_warning(coverage) -> list[str]:
    """Say plainly when the upload cannot support a 90-minute lecture.

    scope_fit only ever distinguished "too much" from "FIT", so a two-page
    syllabus header was reported as a fit for ninety minutes. Faculty saw a
    complete-looking twenty-unit deck and a scatter of coverage-rubric misses
    that never named the actual problem.
    """
    majors = len([x for x in coverage if x.importance == "major"])
    if majors >= MIN_MAJOR_CHECKPOINTS_FOR_A_LECTURE:
        return []
    return [
        f"SOURCE TOO THIN: this upload yields only {majors} major teaching "
        f"checkpoint{'s' if majors != 1 else ''}, below the {MIN_MAJOR_CHECKPOINTS_FOR_A_LECTURE} "
        "needed for a 90-minute lecture. The teaching units will revisit the same "
        "material under different questions rather than cover new ground. Upload the "
        "chapter or lecture slides themselves; a syllabus, cover page or course "
        "outline does not carry teachable content."
    ]


_SECTION_ORDINAL = re.compile(r"^\d{1,2}\.\d+(?:\.\d+)?\b")


def _is_heading_label(text: str) -> bool:
    """Whether a label reads as the page's own heading rather than its prose.

    A page that only continues the previous section has no heading of its own.
    Its text is still a coverage checkpoint, but promoting it to a topic family
    put whole sentences on the domain-spine slide, which then overflowed.
    """
    clean = str(text or "").strip()
    if not clean:
        return False
    return bool(_SECTION_ORDINAL.match(clean)) or _is_title_like(clean) or (
        3 <= len(clean.split()) <= 11 and len(clean) <= 100)


def _heading_line(lines: list[str]) -> str | None:
    """The page's own heading, or None when it carries body text only."""
    clean = [x for x in lines if not _is_furniture_line(x)]
    # Numbered chapter sections are the strongest deterministic checkpoint.
    for x in clean[:12]:
        if _SECTION_ORDINAL.match(x):
            return _clean(x, 120)
    # Prefer compact title-like lines over body sentences. The word floor used
    # to be three, which rejected real two-word headings ("Cost/dependability
    # curve") and then accepted the wrapped body line under them ("Because of
    # very high costs of") - a mid-sentence fragment became the unit title and
    # the unit question built on it read as broken English.
    for x in clean[:12]:
        if _is_title_like(x):
            return _clean(x, 120)
    for x in clean[:12]:
        if 3 <= len(x.split()) <= 11 and len(x) <= 100:
            return _clean(x, 120)
    return None


def _choose_label(lines: list[str], page_no: int) -> str:
    heading = _heading_line(lines)
    if heading:
        return heading
    clean = [x for x in lines if not _is_furniture_line(x)]
    # Nothing on the page reads as a heading, so build one from its opening
    # statement. Returning the raw line printed a wrapped body fragment as the
    # slide heading, which is what "ends mid thought" was catching downstream.
    return _headline_from(clean[:3]) if clean else f"Primary source page {page_no}"


MIN_HEADLINE_WORDS = 3


def _headline_from(lines) -> str:
    """A readable heading derived from body text, never a mid-sentence fragment.

    Takes the following lines as well as the first: a page wraps its sentences at
    the column width, so the opening statement usually finishes on the next line.
    Reading only the first line produced headings that stopped in the middle of
    the sentence they were built from.
    """
    if isinstance(lines, str):
        lines = [lines]
    text = _clean(" ".join(str(x) for x in list(lines)[:3]), 400)
    # A heading is the page's first claim, not a run of them.
    text = re.split(r"(?<=[.!?])\s+|\s+[·•]\s+", text, maxsplit=1)[0].strip()
    # A colon usually separates a topic from its expansion ("Security threats
    # fall into three types: ..."). Take the topic, unless the topic is a single
    # administrative word ("Exercises:") that names nothing teachable.
    head = text.split(":", 1)[0].strip()
    if len(head.split()) >= MIN_HEADLINE_WORDS:
        text = head
    text = _clean(text, 120)
    # A width-driven clip lands wherever the line ran out, so drop trailing
    # function words until the heading stands on its own.
    words = text.split()
    while len(words) > MIN_HEADLINE_WORDS and _DANGLING_TAIL.search(words[-1]):
        words.pop()
    return " ".join(words).rstrip(" ,;:-")



def _page_importance(label: str, lines: list[str], page_no: int) -> str:
    low = label.lower()
    support_markers = (
        "class", "take-home", "to master", "your challenge", "assignment",
        "today lecture topic", "the big picture", "big picture", "cimt compass",
        "setting the stage", "our roadmap", "the core idea", "the big why",
        "key clo", "availability: system is ready", "to master today's lesson",
        "ai era", "ai revolution", "aligning with industry trends", "keeping up with trends",
        "example", "ex;", "remember", "cautionary tale", "how to show this",
        "topics covered", "key points", "references", "exercises", "further reading",
    )
    slide_no = _embedded_slide_number(label)
    if page_no == 1 or slide_no == 1 or any(m in low for m in support_markers):
        return "supporting"
    blob = " ".join(lines[:18]).lower()
    if slide_no and len(blob.split()) < 12:
        return "supporting"
    major_markers = (
        "definition", "requirements", "architecture", "process", "model", "metrics",
        "measurement", "testing", "assurance", "design", "fault", "failure",
        "security", "reliability", "dependability", "resilience", "safety",
        "component", "composition", "reuse", "distributed", "client", "server",
        "middleware", "saas", "service-oriented", "risk assessment", "formal methods",
        "redundancy", "diversity", "sociotechnical", "programming for",
    )
    if re.match(r"^\d{1,2}\.\d+(?:\.\d+)?\b", label) or any(m in blob for m in major_markers):
        return "major"
    content_lines = [x for x in lines if not _is_furniture_line(x)]
    return "major" if len(content_lines) >= 4 and len(" ".join(content_lines).split()) >= 28 else "supporting"


class Chunk(tuple):
    """One coverage chunk: (idx, label, excerpt), plus the source pages it spans.

    A slide deck chunks one page at a time, so the index is the page. A book
    chapter chunks by section, and a section runs over several pages; the pages
    travel with the chunk so the anchor can cite all of them while the tuple
    shape every existing caller unpacks stays the same.
    """
    pages: tuple[int, ...]

    def __new__(cls, idx: int, label: str, excerpt: str, pages=()):
        obj = super().__new__(cls, (idx, label, excerpt))
        obj.pages = tuple(pages) or (idx,)
        return obj


def _pdf_chunks(path: Path) -> list[tuple[int, str, str]]:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages: list[tuple[int, list[str]]] = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        lines = _meaningful_lines(text)
        if lines:
            pages.append((i, lines))
    if _looks_like_book(pages):
        # pypdf breaks kerned words apart ("soft ware", "engi neering") on
        # typeset book pages; PyMuPDF reads the same pages cleanly. Decks keep
        # the pypdf lines their furniture rules were tuned on.
        return _section_chunks(_fitz_page_lines(path) or pages)
    out: list[tuple[int, str, str]] = []
    for position, (i, lines) in enumerate(pages):
        label = _choose_label(lines, i)
        body = [x for x in lines if not _is_furniture_line(x)]
        following = pages[position + 1][1] if position + 1 < len(pages) else []
        excerpt = " · ".join(_carry_dangling_tail(body, [x for x in following if not _is_furniture_line(x)]))
        out.append(Chunk(i, label, excerpt))
    return out


# The page ended, the sentence did not. Nothing on this page can complete it, so
# the statement reached the slide as "...the infrastructure is configured to".
MAX_CARRIED_LINES = 2


def _carry_dangling_tail(body: list[str], following: list[str]) -> list[str]:
    """Complete a page's last sentence from the page that continues it."""
    if not body or not following:
        return body
    tail = body[-1].rstrip()
    if not _DANGLING_TAIL.search(tail.rstrip(".")) or tail.endswith((".", "!", "?", ":")):
        return body
    carried = tail
    for line in following[:MAX_CARRIED_LINES]:
        carried = f"{carried} {line.strip()}"
        if carried.rstrip().endswith((".", "!", "?")):
            break
    return [*body[:-1], carried]


# A book chapter is not a slide deck: its unit of meaning is the section, not
# the page. Chunking it by page gave every teaching slot a heading that was the
# first wrapped line of prose on some page ("attack. / 14.3 I System
# survivability") and split each section across two or three slots. The
# instructor's own decks title each slide with the section it teaches; a chapter
# is chunked the same way.
MIN_BOOK_PAGES = 6
MIN_BOOK_WORDS_PER_PAGE = 200
BOOK_HEADER_SHARE = 0.5
MAX_SECTION_WORDS = 900
MAX_EXCERPT_CHARS = 3200
_RUNNING_HEADER_GLYPH = re.compile(r"^\d{1,2}(?:\.\d+)*\s+I\s+\S")
_SECTION_HEADING = re.compile(r"^(\d{1,2}\.\d+(?:\.\d+)?)\s+([A-Z][^.]{2,80})$")
_ORDINAL_ONLY = re.compile(r"^\d{1,2}\.\d+(?:\.\d+)?$")
_BACK_MATTER = ("references", "exercises", "further reading", "key points", "chapter summary", "summary")


def _fitz_page_lines(path: Path) -> list[tuple[int, list[str]]]:
    try:
        import fitz
        doc = fitz.open(str(path))
    except Exception:
        return []
    pages: list[tuple[int, list[str]]] = []
    try:
        for i, page in enumerate(doc, 1):
            lines = _meaningful_lines(page.get_text("text"))
            if lines:
                pages.append((i, lines))
    finally:
        doc.close()
    return pages


def _back_matter_name(line: str) -> str:
    """"K E Y  P O I N T S" and "References" both name a back-matter block."""
    text = str(line or "").strip().rstrip(":")
    if not text or len(text) > 40:
        return ""
    compact = re.sub(r"\s+", "", text).lower()
    names = {re.sub(r"\s+", "", x): x for x in _BACK_MATTER}
    if compact in names and (text.isupper() or len(text.split()) <= 3):
        return names[compact].title()
    return ""


def _is_running_header(line: str) -> bool:
    text = str(line or "").strip()
    if len(text.split()) < 2:
        return False
    edge_number = bool(_LEADING_PAGE_NO.search(text) or _TRAILING_PAGE_NO.search(text))
    return edge_number or bool(_RUNNING_HEADER_GLYPH.match(text))


def _looks_like_book(pages: list[tuple[int, list[str]]]) -> bool:
    if len(pages) < MIN_BOOK_PAGES:
        return False
    words = sorted(sum(len(x.split()) for x in lines) for _n, lines in pages)
    if words[len(words) // 2] < MIN_BOOK_WORDS_PER_PAGE:
        return False
    with_header = sum(1 for _n, lines in pages if any(_is_running_header(x) for x in lines[:3]))
    return with_header / len(pages) >= BOOK_HEADER_SHARE


def _join_wrapped(lines: list[str]) -> str:
    """Rejoin lines the column width broke, restoring hyphenated words."""
    text = ""
    for line in lines:
        if not text:
            text = line
        elif text.endswith("-") and line[:1].islower():
            text = text[:-1] + line
        else:
            text = f"{text} {line}"
    return text


def _sentences(text: str) -> list[str]:
    parts = re.split(r"(?<=[.!?])\s+(?=[A-Z(\"\u201c])", " ".join(text.split()))
    return [x.strip() for x in parts if x.strip()]


def _recurring_lines(pages: list[tuple[int, list[str]]]) -> set[str]:
    """Running headers by frequency, page-number-free, before any sectioning.

    The edge-number rule catches "368 Chapter 14 I Security engineering"; a
    reader that emits the header and its page number as separate lines leaves
    "Chapter 14 I Security engineering" bare, and it opened every section's
    first sentence. What makes it a header is that it is on half the pages.
    """
    counts: dict[str, int] = {}
    for _n, lines in pages:
        keys = set()
        for line in lines[:4]:
            keys |= _furniture_keys(line)
        for key in keys:
            counts[key] = counts.get(key, 0) + 1
    threshold = max(3, int(len(pages) * RECURRING_FURNITURE_SHARE))
    return {key for key, n in counts.items() if n >= threshold}


def _section_chunks(pages: list[tuple[int, list[str]]]) -> list[Chunk]:
    recurring = _recurring_lines(pages)

    def is_header(line: str) -> bool:
        return _is_running_header(line) or bool(_furniture_keys(line) & recurring)

    first_page, first_lines = pages[0]
    title = next((x for x in first_lines if not is_header(x) and not x.isdigit()), "Primary lecture")
    title = _clean(title, 120)
    # (label, [(line, page), ...]) in source order. The opener page is its own
    # chunk; sections start from the second page, so a contents list on the
    # opener never spawns empty sections.
    sections: list[tuple[str, list[tuple[str, int]]]] = [(title, [(x, first_page) for x in first_lines if not is_header(x)])]
    current_label = f"Introduction to {title}"
    current: list[tuple[str, int]] = []
    pending_ordinal = ""

    def close():
        nonlocal current
        if current:
            sections.append((current_label, current))
        current = []

    for page_no, lines in pages[1:]:
        for line in lines:
            if is_header(line):
                continue
            if _ORDINAL_ONLY.match(line):
                pending_ordinal = line
                continue
            if pending_ordinal:
                line = f"{pending_ordinal} {line}"
                pending_ordinal = ""
            heading = _SECTION_HEADING.match(line)
            if heading and len(line.split()) <= 9:
                close()
                current_label = _clean(line, 120)
                continue
            back_matter = _back_matter_name(line)
            if back_matter:
                close()
                current_label = back_matter
                continue
            current.append((line, page_no))
    close()

    out: list[Chunk] = []
    for label, entries in sections:
        for part_no, part in enumerate(_split_long_section(entries), 1):
            part_label = label if part_no == 1 and len(entries) == len(part) else f"{label} (part {part_no})"
            text = _join_wrapped([line for line, _p in part])
            excerpt = " · ".join(_sentences(text))
            if len(excerpt) > MAX_EXCERPT_CHARS:
                excerpt = excerpt[:MAX_EXCERPT_CHARS].rsplit(" · ", 1)[0]
            page_list = sorted({p for _l, p in part})
            out.append(Chunk(len(out) + 1, part_label, excerpt, page_list))
    return out


def _split_long_section(entries: list[tuple[str, int]]) -> list[list[tuple[str, int]]]:
    """Cut a section that outruns one teaching slot at a sentence boundary.

    The design-guidelines section of a chapter runs to five pages; as one chunk
    it would take one slot and lose most of itself to the canvas, while a thin
    section next to it took a slot of its own.
    """
    total = sum(len(line.split()) for line, _p in entries)
    if total <= MAX_SECTION_WORDS:
        return [entries]
    parts_wanted = -(-total // MAX_SECTION_WORDS)
    target = total / parts_wanted
    parts: list[list[tuple[str, int]]] = []
    current: list[tuple[str, int]] = []
    count = 0
    for line, page in entries:
        current.append((line, page))
        count += len(line.split())
        if count >= target and line.rstrip().endswith((".", "!", "?")) and len(parts) < parts_wanted - 1:
            parts.append(current)
            current, count = [], 0
    if current:
        parts.append(current)
    return parts


def _pptx_chunks(path: Path) -> list[tuple[int, str, str]]:
    from pptx import Presentation
    prs = Presentation(str(path))
    out: list[tuple[int, str, str]] = []
    for i, slide in enumerate(prs.slides, 1):
        raw = "\n".join(sh.text for sh in slide.shapes if hasattr(sh, "text") and sh.text)
        lines = _meaningful_lines(raw)
        if not lines:
            continue
        out.append((i, _choose_label(lines, i), " · ".join(x for x in lines if not _is_furniture_line(x))))
    return out



def _embedded_slide_number(text: str) -> int | None:
    m = re.match(r"^\s*SLIDE\s+(\d+)\s*:", str(text or ""), flags=re.I)
    return int(m.group(1)) if m else None


def _declared_slide_count(text: str) -> int | None:
    """Return the player-declared deck size before host recommendation chrome."""
    head = re.split(r"\bRecommended\b", str(text or ""), maxsplit=1, flags=re.I)[0]
    for pattern in (r"\b1\s*/\s*(\d{1,4})\b", r"\b(\d{1,4})\s+slides\b"):
        m = re.search(pattern, head, flags=re.I)
        if m:
            value = int(m.group(1))
            if 2 <= value <= 500:
                return value
    return None


def _anchor_slides(anchor: str) -> set[int]:
    """Extract explicit P1 slide coordinates from source anchors."""
    found: set[int] = set()
    for m in re.finditer(r"\bSLIDES?\s+(\d+)(?:\s*[-–—]\s*(\d+))?", str(anchor or ""), flags=re.I):
        start = int(m.group(1)); end = int(m.group(2) or start)
        if end < start:
            start, end = end, start
        if end - start <= 200:
            found.update(range(start, end + 1))
    return found



def _doc_chunks(path: Path) -> list[tuple[int, str, str]]:
    # Public presentation pages may expose `SLIDE n:` labels in extracted
    # text. Respect the deck boundary and never ingest Recommended cards.
    text = extract_source_text(path, limit=500_000)
    declared_slides = _declared_slide_count(text)
    paras = [_clean(x, 520) for x in re.split(r"\n{1,}", text) if _clean(x, 520)]
    out: list[tuple[int, str, str]] = []
    for i, para in enumerate(paras[:200], 1):
        if re.match(r"^recommended\b", para, flags=re.I):
            break
        words = para.split()
        label = _clean(" ".join(words[:12]), 110)
        slide_no = _embedded_slide_number(label) or _embedded_slide_number(para)
        if declared_slides and slide_no and slide_no > declared_slides:
            continue
        out.append((i, label, para))
    return out[:80]


def _chunks(path: Path) -> tuple[str, list[tuple[int, str, str]]]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return "PAGE", _pdf_chunks(path)
    if ext == ".pptx":
        return "SLIDE", _pptx_chunks(path)
    return "SECTION", _doc_chunks(path)


# "Class2: Dependable Systems" names the session; the lecture is about dependable
# systems. The ordinal travels into every derived sentence - the crisis, the five
# CLOs, the Unit 1 heading - and reads as "a decision about Class2: Dependable
# Systems", so it is dropped once, here, where the title is chosen.
_SESSION_PREFIX = re.compile(
    r"^\s*(?:class|lecture|lesson|session|week|lab)\s*"
    r"[-\u2013\u2014]?\s*\d{1,3}[a-z]?\s*(?:[:.)\-\u2013\u2014]|\b)\s*",
    re.IGNORECASE,
)


def _strip_session_prefix(title: str) -> str:
    text = str(title or "").strip()
    stripped = _SESSION_PREFIX.sub("", text, count=1).strip(" -\u2013\u2014:.")
    # Only accept the reduction when a real subject survives it: "Class 2" on its
    # own is the best title that page offers, and a remainder that continues the
    # phrase ("Part 2 of the design review") means the ordinal was not a prefix.
    if len(stripped.split()) < 2 or stripped.split()[0].lower() in _STOPWORDS:
        return text
    return stripped


def _title_from(primary_name: str, chunks: list[tuple[int, str, str]]) -> str:
    stem = Path(primary_name).stem.replace("_", " ").replace("-", " ")
    stem = re.sub(r"\s+", " ", stem).strip()
    # Prefer a plausible first source heading over a machine filename.
    for _page, label, _body in chunks[:6]:
        first = _strip_page_coordinate(_display_label(_clean(label, 120)))
        if not (5 <= len(first) <= 120):
            continue
        if re.fullmatch(r"[A-Z0-9_. -]+", first) or _is_furniture_line(first):
            continue
        return _strip_session_prefix(first)
    return _strip_session_prefix(stem) or "Primary lecture"


# A running header or footer appears on most slides; a real checkpoint does not.
RECURRING_FURNITURE_SHARE = 0.30
MIN_SLIDES_FOR_FREQUENCY_RULE = 8


def _recurring_furniture(chunks: list[tuple[int, str, str]]) -> set[str]:
    """Lines that repeat across the deck, which is what makes them furniture.

    The hardcoded furniture list only knows the CPIT decks. A Sommerville chapter
    carries its own running footer, and it was being picked up as a checkpoint and
    printed as a unit heading. Frequency identifies furniture in any deck without
    naming any of them: content appears on its own slide, chrome appears on all.
    """
    if len(chunks) < MIN_SLIDES_FOR_FREQUENCY_RULE:
        return set()
    counts: dict[str, int] = {}
    for _idx, _label, excerpt in chunks:
        keys = set()
        for x in _meaningful_lines(excerpt.replace(" \u00b7 ", "\n")):
            keys |= _furniture_keys(x)
        for line in keys:
            counts[line] = counts.get(line, 0) + 1
    threshold = max(3, int(len(chunks) * RECURRING_FURNITURE_SHARE))
    return {line for line, n in counts.items() if n >= threshold}


def _norm_key(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", str(text or "").lower()).strip()


# A page coordinate is not part of the heading a learner reads.
_LEADING_PAGE_NO = re.compile(r"^\s*\d{1,4}\s+(?=[^\d\s])")
_TRAILING_PAGE_NO = re.compile(r"\s+\d{1,4}\s*$")


def _strip_page_coordinate(text: str) -> str:
    """Drop a leading/trailing page number from a heading, keeping section numbers.

    "14.1 I Security risk management 369" is a real heading wearing a page number;
    the number belongs in the source anchor, never in the slide title. A section
    ordinal ("14.1", "14.2.2") is part of the heading and is left in place because
    it is followed by a dot rather than whitespace.
    """
    out = _TRAILING_PAGE_NO.sub("", _LEADING_PAGE_NO.sub("", str(text or ""))).strip()
    return out if len(out.split()) >= 2 else str(text or "").strip()


def _furniture_keys(text: str) -> set[str]:
    """Comparison keys for one line: exact, plus page-coordinate-free when it wears one.

    A book chapter prints the same header on every page and only the page number
    changes, so exact-string frequency never saw a repeat: "Chapter 14 I Security
    engineering 367" and "368 Chapter 14 I Security engineering" are two distinct
    strings, and both were surviving as teaching checkpoints - the faculty deck's
    domain spine then read as a list of page headers.

    The number-free form is only compared for a line that carries its number at the
    edge, where a page coordinate sits. A body line whose number is part of what it
    says ("Body sentence number 4") keeps its exact key, so a series of numbered
    content lines is never mistaken for chrome. A section heading is unaffected
    either way: it spans a few pages and stays far below the recurrence threshold.
    """
    keys = {_norm_key(text)}
    raw = str(text or "")
    if _LEADING_PAGE_NO.search(raw) or _TRAILING_PAGE_NO.search(raw):
        keys.add(_norm_key(_strip_page_coordinate(raw)))
    return {k for k in keys if k}


# A spine entry has to be projectable next to a dozen others. A source page that
# carries no heading of its own contributes a checkpoint, not a family name, and
# the family names that remain are compacted to heading length.
MAX_FAMILY_NAME_CHARS = 64


def _family_name(label: str) -> str:
    text = _clean(label, MAX_FAMILY_NAME_CHARS)
    words = text.split()
    while len(words) > MIN_HEADLINE_WORDS and _DANGLING_TAIL.search(words[-1]):
        words.pop()
    return " ".join(words).rstrip(" ,;:-") or str(label or "").strip()


def _anchor_for_pages(coordinate: str, pages) -> str:
    pages = sorted(set(int(x) for x in pages))
    if len(pages) == 1:
        return f"[P1] {coordinate} {pages[0]}"
    if pages == list(range(pages[0], pages[-1] + 1)):
        return f"[P1] {coordinate}S {pages[0]}\u2013{pages[-1]}"
    return "; ".join(f"[P1] {coordinate} {x}" for x in pages)


def build_deterministic_source_profile(bundle: SourceBundle, reason: str = "") -> SourceProfile:
    primary = bundle.primary
    coordinate, chunks = _chunks(primary.path)
    if not chunks:
        # Last-resort single source item. It is intentionally explicit rather than
        # pretending that a semantic profile exists when extraction yielded nothing.
        chunks = [(1, Path(primary.display_name).stem, "Primary source supplied by faculty")]

    # Preserve source order while suppressing exact duplicate headings caused by
    # recurring lecture furniture. We still keep a page-level coverage item for
    # every content-bearing source page/slide.
    recurring = _recurring_furniture(chunks)

    families: list[TopicFamily] = []
    seen_family: set[str] = set()
    coverage: list[CoverageItem] = []
    for chunk in chunks[:80]:
        idx, label, excerpt = chunk
        pages = tuple(getattr(chunk, "pages", ()) or (idx,))
        clean_label = _clean(label, 120) or f"Source {coordinate.title()} {idx}"
        if recurring and _furniture_keys(_display_label(clean_label)) & recurring:
            # A running header chosen as this slide's heading says nothing about
            # the slide. Take the first line that is actually specific to it.
            candidates = _meaningful_lines(excerpt.replace(" \u00b7 ", "\n"))
            for position, candidate in enumerate(candidates):
                if not (_furniture_keys(candidate) & recurring) and not _is_furniture_line(candidate):
                    clean_label = _headline_from(candidates[position:position + 3])
                    break
        embedded_slide = _embedded_slide_number(clean_label)
        anchor = f"[P1] SLIDE {embedded_slide}" if embedded_slide else _anchor_for_pages(coordinate, pages)
        # The coordinate now lives in the anchor; the heading must read as prose.
        clean_label = _strip_page_coordinate(_display_label(clean_label)) or f"Source {coordinate.title()} {idx}"
        # The parts of one long section are one family with several checkpoints.
        key = re.sub(r"[^a-z0-9]+", " ", re.sub(r"\s*\(part \d+\)$", "", clean_label).lower()).strip()
        source_lines = _meaningful_lines(excerpt.replace(" · ", "\n"))
        importance = _page_importance(clean_label, source_lines, idx)
        is_title_like = importance == "supporting"
        coverage.append(CoverageItem(
            id=f"P1-{coordinate[0]}{idx:02d}",
            label=clean_label,
            knowledge_type=_knowledge_type(excerpt),
            importance=importance,
            source_anchor=anchor,
            # Preserve the source payload. Display compaction is not extraction.
            why_important=excerpt,
        ))
        if key and key not in seen_family and not is_title_like and _is_heading_label(clean_label):
            seen_family.add(key)
            families.append(TopicFamily(
                name=_family_name(re.sub(r"\s*\(part \d+\)$", "", clean_label)),
                source_anchor=anchor,
                why_important=_clean(excerpt, 260),
            ))

    if not families:
        first = coverage[0]
        families = [TopicFamily(name=first.label, source_anchor=first.source_anchor, why_important=first.why_important)]
        coverage[0].importance = "major"

    title = _title_from(primary.display_name, chunks)
    focus = _clean(bundle.lecture_focus, 260) or title
    warning = (
        "Source profile was built deterministically because the AI profiling call was unavailable: "
        + _clean(reason, 300)
        + ". The full primary source is still supplied to blueprint generation; RELEASE still requires normal gates."
    ) if reason else "Source profile built locally without an AI call. This is a source inventory, not semantic approval; RELEASE still requires normal gates."
    return SourceProfile(
        lecture_title=title,
        course_or_level="Faculty-supplied lecture",
        weekly_focus=focus,
        topic_families=families[:40],
        coverage_items=coverage[:80],
        technical_boundaries=[
            "P1 remains the authority for technical claims and terminology.",
            "Deterministic fallback records source coordinates; it does not add external technical claims.",
        ],
        source_warnings=[warning, *_thin_source_warning(coverage)],
        session_minutes=90,
        scope_fit="COMPRESS" if len([x for x in coverage if x.importance == "major"]) > 16 else "FIT",
        in_scope_families=[x.name for x in families[:40]],
        deferred_topics=[],
        source_conflicts=[],
        source_manifest=bundle.manifest_lines(),
    )




MIN_ANCHORS_FOR_SEMANTIC_BOUNDARY = 3


def _semantic_slide_ceiling(profile: SourceProfile) -> int | None:
    # Honor explicit Source Lock bounds when flattened host text contains extra slides.
    candidates: list[int] = []
    for warning in profile.source_warnings or []:
        text = str(warning or "")
        for m in re.finditer(r"(?:core|primary|actual|lecture|chapter)[^.\n]{0,180}?slides?\s*1\s*[-–—]\s*(\d{1,4})", text, flags=re.I):
            candidates.append(int(m.group(1)))
        for m in re.finditer(r"slides?\s*1\s*[-–—]\s*(\d{1,4})[^.\n]{0,180}?(?:core|primary|actual|lecture|chapter)", text, flags=re.I):
            candidates.append(int(m.group(1)))
        for m in re.finditer(r"slides?\s*(\d{1,4})\s*(?:through|to|[-–—])\s*(\d{1,4})[^.\n]{0,220}?(?:external|unrelated|recommended|filter(?:ed)?\s*out|host[- ]page)", text, flags=re.I):
            start = int(m.group(1))
            if start > 1:
                candidates.append(start - 1)
    valid = [x for x in candidates if 2 <= x <= 500]

    # The semantic profile's own anchors are the authoritative chapter extent.
    # If the model mapped this chapter's families through slide 58, then slide 59
    # is outside the chapter no matter what the host page happens to serve after
    # it - recommendation cards, cross-chapter summaries, unrelated decks. A
    # warning phrase describing that boundary is a convenience, never the only
    # way to learn it.
    anchored: set[int] = set()
    for row in profile.coverage_items or []:
        anchored.update(_anchor_slides(row.source_anchor))
    if len(anchored) >= MIN_ANCHORS_FOR_SEMANTIC_BOUNDARY:
        valid.append(max(anchored))

    return min(valid) if valid else None


_STOPWORDS = {
    "the", "a", "an", "and", "or", "of", "for", "to", "in", "on", "with", "vs",
    "versus", "its", "their", "how", "what", "why", "is", "are", "be", "using",
}


def _significant_tokens(label: str) -> set[str]:
    words = re.findall(r"[a-z0-9]+", str(label or "").lower())
    return {w for w in words if len(w) > 2 and w not in _STOPWORDS}


def _restates_semantic_coverage(item_tokens: set[str], semantic_token_sets: list[set[str]]) -> bool:
    """Whether a raw checkpoint only repeats a family the model already mapped.

    Inside a semantically mapped range the model's own family label is the better
    teaching heading. Re-adding the raw slide line next to it costs a unit of the
    90-minute budget and teaches nothing new, which is how a deck ends up with a
    near-empty slide.
    """
    if not item_tokens:
        return False
    for tokens in semantic_token_sets:
        if not tokens:
            continue
        shared = len(item_tokens & tokens)
        if shared and shared / min(len(item_tokens), len(tokens)) >= 0.7:
            return True
    return False


def reconcile_source_profile(ai_profile: SourceProfile, bundle: SourceBundle) -> SourceProfile:
    """Add only uncovered real source slides to the semantic P1 contract."""
    deterministic = build_deterministic_source_profile(bundle, "chapter-completeness reconciliation")
    existing = {re.sub(r"[^a-z0-9]+", " ", x.label.lower()).strip() for x in ai_profile.coverage_items}
    semantic_slides: set[int] = set()
    for row in ai_profile.coverage_items:
        semantic_slides.update(_anchor_slides(row.source_anchor))

    semantic_ceiling = _semantic_slide_ceiling(ai_profile)
    semantic_token_sets = [_significant_tokens(x.label) for x in ai_profile.coverage_items]
    semantic_span = max(semantic_slides) if semantic_slides else 0

    merged = list(ai_profile.coverage_items)
    for item in deterministic.coverage_items:
        if item.importance != "major":
            continue
        item_slides = _anchor_slides(item.source_anchor)
        # A slide past the chapter's semantic extent is not this chapter's content.
        if semantic_ceiling and item_slides and max(item_slides) > semantic_ceiling:
            continue
        if item_slides and item_slides.issubset(semantic_slides):
            continue
        # Inside the mapped range, keep only what the model genuinely missed.
        if item_slides and semantic_span and min(item_slides) <= semantic_span:
            if _restates_semantic_coverage(_significant_tokens(item.label), semantic_token_sets):
                continue
        key = re.sub(r"[^a-z0-9]+", " ", item.label.lower()).strip()
        if key and key not in existing:
            merged.append(item)
            existing.add(key)
            semantic_slides.update(item_slides)
            semantic_token_sets.append(_significant_tokens(item.label))

    ai_profile.coverage_items = merged[:80]
    ai_profile.deferred_topics = []
    ai_profile.source_warnings = list(dict.fromkeys([
        *ai_profile.source_warnings,
        "Deterministic source-coordinate reconciliation excludes duplicate semantic slides and host-page recommendation cards from mandatory P1 coverage.",
    ]))[:20]
    return ai_profile
