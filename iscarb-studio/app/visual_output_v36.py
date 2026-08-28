from __future__ import annotations

import json
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pypdf import PdfReader, PdfWriter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from .models import Blueprint, LectureUnit
from .visual_provenance import classify_visual
from .faculty_visual import export_faculty_presenter_pptx as _export_pptx
from .faculty_visual import render_faculty_presenter_preview as _render_preview
from . import presenter_pdf as pp
from .source_visuals import plans_for_blueprint, local_asset, asset_data_uri, VisualPlan


MAGENTA = RGBColor(210, 11, 109)
TEAL = RGBColor(10, 53, 62)
WHITE = RGBColor(255, 255, 255)
SAND = RGBColor(215, 180, 135)
INK = RGBColor(24, 31, 29)
MUTED = RGBColor(102, 112, 107)
SOFT = RGBColor(247, 244, 237)
GREEN = RGBColor(32, 142, 95)


def _ppt_provenance(slide, label: str, citation: str) -> None:
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.42), Inches(.30), Inches(3.08), Inches(.36))
    sh.fill.solid(); sh.fill.fore_color.rgb = TEAL
    sh.line.color.rgb = MAGENTA; sh.line.width = Pt(.7)
    tf = sh.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Aptos"; p.font.size = Pt(7.2); p.font.bold = True; p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    tx = slide.shapes.add_textbox(Inches(8.38), Inches(6.80), Inches(4.35), Inches(.22))
    tf = tx.text_frame; tf.clear(); tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = citation[:92]
    p.font.name = "Aptos"; p.font.size = Pt(5.8); p.font.color.rgb = SAND
    p.alignment = PP_ALIGN.RIGHT


def _ppt_text(slide, x, y, w, h, text, size=11, color=INK, bold=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = str(text or "")
    p.font.name = "Aptos"; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return tx


def _ppt_source_visual(slide, unit: LectureUnit, plan: VisualPlan) -> bool:
    if not plan.asset:
        return False
    path = local_asset(plan.asset)
    if not path:
        return False

    # Cover the legacy infographic region only; preserve title/question/footer.
    cover = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.36), Inches(1.92), Inches(12.62), Inches(4.92))
    cover.fill.solid(); cover.fill.fore_color.rgb = RGBColor(250, 249, 246); cover.line.fill.background()

    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.52), Inches(2.06), Inches(8.35), Inches(4.50))
    frame.fill.solid(); frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(219, 224, 219); frame.line.width = Pt(1)
    slide.shapes.add_picture(str(path), Inches(.67), Inches(2.20), width=Inches(8.05), height=Inches(4.20))

    side = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.05), Inches(2.06), Inches(3.65), Inches(4.50))
    side.fill.solid(); side.fill.fore_color.rgb = TEAL; side.line.color.rgb = MAGENTA; side.line.width = Pt(.9)

    _ppt_text(slide, 9.32, 2.30, 3.10, .30, f"SOURCE VISUAL · P1 SLIDE {plan.source_slide}", 8.0, SAND, True)
    _ppt_text(slide, 9.32, 2.72, 3.06, .34, "WHY THIS MATTERS", 10.0, WHITE, True)
    _ppt_text(slide, 9.32, 3.08, 3.04, 1.02, plan.teaching_purpose, 11.0, WHITE, False)
    _ppt_text(slide, 9.32, 4.20, 3.04, .34, "YOU TRY", 10.0, RGBColor(91, 224, 178), True)
    _ppt_text(slide, 9.32, 4.56, 3.04, 1.14, unit.student_action, 10.2, WHITE, False)
    _ppt_text(slide, 9.32, 5.88, 3.04, .46, plan.citation, 6.6, SAND, False)
    return True


def export_presenter_pptx(bp: Blueprint, out: Path) -> Path:
    out = _export_pptx(bp, Path(out))
    prs = Presentation(str(out))
    plans = plans_for_blueprint(bp)
    for slide, unit, plan in zip(prs.slides, bp.units, plans):
        if plan.reuse_mode == "USE" and _ppt_source_visual(slide, unit, plan):
            continue
        prov = classify_visual(unit)
        _ppt_provenance(slide, prov.label, prov.citation)
    prs.save(str(out))
    return out


def render_presenter_preview(bp: Blueprint, release_state: str = "BLOCKED") -> str:
    html = _render_preview(bp, release_state)
    plans = plans_for_blueprint(bp)
    source_payload = []
    fallback_payload = []
    for unit, plan in zip(bp.units, plans):
        if plan.reuse_mode == "USE" and plan.asset:
            uri = asset_data_uri(plan.asset)
            if uri:
                source_payload.append({
                    "unit": unit.number,
                    "image": uri,
                    "slide": plan.source_slide,
                    "purpose": plan.teaching_purpose,
                    "action": unit.student_action,
                    "citation": plan.citation,
                })
                continue
        prov = classify_visual(unit)
        fallback_payload.append({"unit": unit.number, **prov.__dict__})

    src_json = json.dumps(source_payload, ensure_ascii=False).replace("</", "<\\/")
    fallback_json = json.dumps(fallback_payload, ensure_ascii=False).replace("</", "<\\/")
    inject = f'''
    <style>
      .vprov{{position:absolute;right:34px;top:22px;z-index:5;background:#0a353e;color:#fff;border:1px solid #d20b6d;border-radius:999px;padding:5px 9px;font-size:7px;font-weight:900;letter-spacing:.04em}}
      .vprov small{{display:block;color:#d7b487;font-size:5.8px;font-weight:650;margin-top:2px;max-width:220px;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}}
      .slide{{position:relative}}
      .sourceVisualStage{{height:100%;display:grid;grid-template-columns:minmax(0,2.45fr) minmax(205px,.85fr);gap:16px;align-items:stretch}}
      .sourceVisualFigure{{margin:0;background:#fff;border:1px solid #dbe2dd;border-radius:16px;padding:10px;display:flex;align-items:center;justify-content:center;min-height:0;overflow:hidden;box-shadow:0 10px 24px rgba(29,41,33,.08)}}
      .sourceVisualFigure img{{width:100%;height:100%;object-fit:contain;display:block}}
      .sourceVisualAside{{background:#0a353e;border:1px solid #d20b6d;border-radius:16px;color:#fff;padding:17px 16px;display:flex;flex-direction:column;min-height:0}}
      .sourceVisualAside .tag{{color:#d7b487;font-size:8px;font-weight:950;letter-spacing:.08em}}
      .sourceVisualAside h3{{font-size:13px;margin:18px 0 6px;color:#fff}}
      .sourceVisualAside p{{font-size:11px;line-height:1.45;margin:0;color:#eaf2ef;font-weight:650}}
      .sourceVisualAside .try{{color:#5be0b2}}
      .sourceVisualAside cite{{margin-top:auto;font-style:normal;color:#d7b487;font-size:7px;line-height:1.35;border-top:1px solid rgba(255,255,255,.15);padding-top:9px}}
    </style>
    <script>
      (function(){{
        const source={src_json};
        const fallback={fallback_json};
        const byUnit=new Map(source.map(x=>[x.unit,x]));
        document.querySelectorAll('.slide').forEach((s,i)=>{{
          const n=i+1, x=byUnit.get(n);
          if(x){{
            const v=s.querySelector('.visual');
            if(v){{
              v.innerHTML=`<div class="sourceVisualStage"><figure class="sourceVisualFigure"><img alt="Source visual from P1 slide ${{x.slide}}" src="${{x.image}}"></figure><aside class="sourceVisualAside"><div class="tag">SOURCE VISUAL · P1 SLIDE ${{x.slide}}</div><h3>WHY THIS MATTERS</h3><p>${{escHtml(x.purpose)}}</p><h3 class="try">YOU TRY</h3><p>${{escHtml(x.action)}}</p><cite>${{escHtml(x.citation)}}</cite></aside></div>`;
            }}
            return;
          }}
          const p=fallback.find(z=>z.unit===n); if(!p)return;
          const b=document.createElement('div');b.className='vprov';b.textContent=p.label;
          const sm=document.createElement('small');sm.textContent=p.citation;b.appendChild(sm);s.appendChild(b);
        }});
        function escHtml(s){{return String(s||'').replace(/[&<>\"']/g,m=>({{'&':'&amp;','<':'&lt;','>':'&gt;','\"':'&quot;',"'":'&#039;'}}[m]));}}
      }})();
    </script>
    '''
    return html.replace("</body>", inject + "</body>")


def _fit_image(c, path: Path, x: float, y: float, w: float, h: float) -> None:
    try:
        img = ImageReader(str(path)); iw, ih = img.getSize()
        scale = min(w / iw, h / ih)
        dw, dh = iw * scale, ih * scale
        c.drawImage(img, x + (w-dw)/2, y + (h-dh)/2, width=dw, height=dh, preserveAspectRatio=True, mask='auto')
    except Exception:
        return


def _pdf_overlay(unit: LectureUnit, plan: VisualPlan, image_path: Path) -> bytes:
    buf = BytesIO(); c = canvas.Canvas(buf, pagesize=(960, 540))
    c.setFillColorRGB(250/255,249/255,246/255); c.rect(24, 68, 912, 380, fill=1, stroke=0)
    c.setFillColorRGB(1,1,1); c.roundRect(34, 80, 625, 350, 12, fill=1, stroke=0)
    _fit_image(c, image_path, 44, 90, 605, 330)
    c.setFillColorRGB(10/255,53/255,62/255); c.roundRect(680, 80, 246, 350, 12, fill=1, stroke=0)
    c.setFillColorRGB(215/255,180/255,135/255); c.setFont("Helvetica-Bold", 7.5)
    c.drawString(698, 405, f"SOURCE VISUAL · P1 SLIDE {plan.source_slide}")
    c.setFillColorRGB(1,1,1); c.setFont("Helvetica-Bold", 10); c.drawString(698, 373, "WHY THIS MATTERS")
    text=c.beginText(698,355); text.setFont("Helvetica",9); text.setLeading(12); text.setFillColorRGB(.92,.96,.94)
    for line in _wrap(plan.teaching_purpose, 36): text.textLine(line)
    c.drawText(text)
    c.setFillColorRGB(91/255,224/255,178/255); c.setFont("Helvetica-Bold",10); c.drawString(698, 270, "YOU TRY")
    text=c.beginText(698,252); text.setFont("Helvetica",8.5); text.setLeading(11); text.setFillColorRGB(1,1,1)
    for line in _wrap(unit.student_action, 36)[:7]: text.textLine(line)
    c.drawText(text)
    c.setStrokeColorRGB(.35,.45,.42); c.line(698,125,908,125)
    c.setFillColorRGB(215/255,180/255,135/255); c.setFont("Helvetica",6)
    for j,line in enumerate(_wrap(plan.citation, 48)[:3]): c.drawString(698,110-j*8,line)
    c.save(); return buf.getvalue()


def _wrap(text: str, width: int) -> list[str]:
    words=str(text or '').split(); lines=[]; line=[]
    for word in words:
        trial=' '.join(line+[word])
        if len(trial)>width and line:
            lines.append(' '.join(line)); line=[word]
        else: line.append(word)
    if line: lines.append(' '.join(line))
    return lines


def export_presenter_pdf(bp: Blueprint, out: Path) -> Path:
    # First create the deterministic Presenter PDF exactly as before.
    original_base = pp._base
    def _base_with_provenance(c, unit):
        original_base(c, unit)
        prov = classify_visual(unit)
        c.setFillColor(pp.TEAL); c.roundRect(650, 493, 188, 22, 11, fill=1, stroke=0)
        c.setFillColor(pp.WHITE); c.setFont("Helvetica-Bold", 7); c.drawCentredString(744, 500, prov.label)
        c.setFillColor(pp.GOLD); c.setFont("Helvetica", 5.6); c.drawRightString(920, 38, prov.citation[:95])
    pp._base = _base_with_provenance
    try:
        pp.export_presenter_pdf(bp, Path(out))
    finally:
        pp._base = original_base

    plans = plans_for_blueprint(bp)
    replacements: dict[int, tuple[LectureUnit, VisualPlan, Path]] = {}
    for i,(unit,plan) in enumerate(zip(bp.units, plans)):
        if plan.reuse_mode == "USE" and plan.asset:
            path = local_asset(plan.asset)
            if path: replacements[i]=(unit,plan,path)
    if not replacements:
        return Path(out)

    reader=PdfReader(str(out)); writer=PdfWriter()
    for i,page in enumerate(reader.pages):
        if i in replacements:
            unit,plan,path=replacements[i]
            overlay=PdfReader(BytesIO(_pdf_overlay(unit,plan,path))).pages[0]
            page.merge_page(overlay)
        writer.add_page(page)
    tmp=Path(out).with_suffix('.source-aware.tmp.pdf')
    with tmp.open('wb') as f: writer.write(f)
    tmp.replace(out)
    return Path(out)
