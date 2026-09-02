from __future__ import annotations

"""One lossless content projection for PDF, editable PPTX and HTML preview.

No guessed fault trees, made-up curves, word-boundary chopping or hidden
annotation substitutions. An exact source page can carry the technical detail;
otherwise every core statement is typeset. Dense drafts remain inspectable and
the readable-fit check prevents them being released without repair.
"""

import base64
import html
import io
import re
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path

from PIL import Image, ImageFilter
from reportlab.lib import colors
from reportlab.pdfbase.pdfmetrics import stringWidth, registerFont
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from .models import Blueprint, LectureUnit
from .source_visuals import FIGURE_KIND, PICTURE_KIND, anchor_slides, local_asset
from .source_visuals_v42 import plans_for_blueprint_v42

W, H = 960, 540
_FONT_ROOT = Path("/usr/share/fonts/truetype/dejavu")
if (_FONT_ROOT / "DejaVuSans.ttf").exists():
    registerFont(TTFont("ISCARB", str(_FONT_ROOT / "DejaVuSans.ttf")))
    registerFont(TTFont("ISCARB-Bold", str(_FONT_ROOT / "DejaVuSans-Bold.ttf")))
else:
    # ReportLab's own bundled font keeps local/no-system-font installations
    # portable too; the production image always supplies DejaVu Sans.
    import reportlab
    _FONT_ROOT = Path(reportlab.__file__).parent / "fonts"
    registerFont(TTFont("ISCARB", str(_FONT_ROOT / "Vera.ttf")))
    registerFont(TTFont("ISCARB-Bold", str(_FONT_ROOT / "VeraBd.ttf")))
INK, GREEN, GOLD, MUTED = "#182B29", "#005B39", "#B78A36", "#526460"
PHASES = {"IFHAM": "UNDERSTAND", "MARIS": "PRACTISE", "ATQAN": "MASTER", "MAYYIZ": "DISTINGUISH"}
# The CIMT compass is the spine the archived CPIT-455 decks put on their own
# slides. The Blueprint records a lens per unit; printing it in the header tells
# a learner which kind of thinking the slide is asking for.
LENSES = {"C": "CONCEPT", "I": "IMPLEMENTATION", "M": "MEASUREMENT", "T": "TREND"}


def eyebrow(u) -> str:
    lenses = " · ".join(LENSES.get(x, x) for x in (u.cimtlens or []))
    phase = PHASES.get(u.phase, u.phase)
    return f"ISCARB / {phase}" + (f" · {lenses}" if lenses else "") + f" / {JOBS[u.number-1]}"
JOBS = [
    "Professional decision & crisis", "Domain spine", "Five measurable outcomes", "Six H-Stack capabilities",
    "Predict · Constrain · Derive · Name", "Mechanism from first principles", "Implementation structure",
    "Alternatives & trade-offs", "Measurement & falsification", "Known · Unknown · Monitor",
    "Contextual application", "Accountability", "Contemporary practice", "Practitioner consequences",
    "Critical AI literacy", "Portfolio challenge", "Constraint mutation", "Evidence policy",
    "Four-level capability rubric", "Bounded assurance & decision",
]


def clean(value) -> str:
    return re.sub(r"\s+", " ", str(value or "")).strip()


def split_item(value: str, fallback: str = "") -> tuple[str, str]:
    value = clean(value)
    parts = re.split(r"\s*(?::| — | – )\s*", value, maxsplit=1)
    if len(parts) == 2 and 1 <= len(parts[0].split()) <= 8:
        return parts[0], parts[1]
    return fallback, value


def contextual_items(u):
    return [("SCENARIO ASSUMPTION", clean(x)) for x in u.scenario_assumptions if clean(x)] + [
        ("CONTEXT / " + clean(u.enrichment_basis[i] if i < len(u.enrichment_basis) else "UNVERIFIED"), clean(x))
        for i,x in enumerate(u.enrichment_content) if clean(x)]


def compact_source_fragments(items):
    """Pack short extraction fragments without summarizing or removing words.

    PDF extraction often makes every cell/line a separate item. Giving each
    fragment a full paragraph gap wastes the page. Preserve labels and complete
    longer statements; join only consecutive short unlabeled fragments.
    """
    if len(items) < 12:
        return items
    result, pending = [], []
    def flush():
        if pending:
            result.append(("", "; ".join(pending)))
            pending.clear()
    for label, body in items:
        if label or len(body.split()) >= 12:
            flush()
            result.append((label, body))
        else:
            if sum(len(x.split()) for x in pending) + len(body.split()) > 40:
                flush()
            pending.append(body)
    flush()
    return result


def teaching_items(bp: Blueprint, u: LectureUnit) -> list[tuple[str, str]]:
    if u.number == 1:
        return [("ENGINEERING CRISIS", clean(bp.central_engineering_crisis)),
                ("PROFESSIONAL PURPOSE", clean(bp.named_ethical_purpose)),
                *[split_item(x, "PRIMARY SOURCE") for x in u.core_content],
                *[split_item(x, "FRAMING") for x in u.pedagogy_content]]
    if u.number == 3:
        return [(c.id, clean(c.statement)) for c in bp.clOs]
    if u.number == 4:
        return [split_item(x) for x in u.pedagogy_content]
    if u.number == 5:
        # Keep the full explanation after prediction; never silently rely on
        # another unit to contain this unit's source facts.
        return (contextual_items(u) + [split_item(x, "REASONING STEP") for x in u.pedagogy_content]
                + ([("SOURCE EXPLANATION", " ".join(clean(x) for x in u.core_content))] if u.core_content else []))
    if u.number == 19:
        return [(r.criterion, " | ".join([r.distinguished, r.ready, r.developing, r.not_yet_ready]))
                for r in bp.rubric_criteria]
    core = compact_source_fragments([split_item(x) for x in u.core_content if clean(x)])
    ped = [split_item(x) for x in u.pedagogy_content if clean(x)]
    if core and not core[0][0]:
        core[0] = ("PRIMARY SOURCE", core[0][1])
    if ped and not ped[0][0]:
        ped[0] = ("PRACTICE", ped[0][1])
    return core + ped + contextual_items(u)


@dataclass
class Text:
    x: float
    y: float
    width: float
    lines: list[str]
    size: float
    color: str = INK
    bold: bool = False


class PresenterLayoutError(ValueError):
    """The source is preserved, but this 20-page projection cannot fit safely."""


def wrap(value: str, width: float, size: float, bold=False) -> list[str]:
    """Measure every word; never silently discard a suffix or a final line."""
    # The column balancer revisits the same text/width/font at many splits.
    # Return a copy so callers cannot mutate the bounded shared cache.
    return list(_wrapped(clean(value), width, size, bold))


# Break points a reader already expects inside a long token: a URL breaks after
# its separators, not in the middle of a word.
_TOKEN_BREAKS = "/-_.,:;=&?"

# A link or an identifier is a real token a line can be broken inside; article
# URLs carrying a full headline run past 150 characters and still belong on the
# slide. A run far longer than that is not a word at all, and it keeps tripping
# the horizontal-overflow guard so the artifact is refused rather than typeset
# as a wall of one token.
MAX_BREAKABLE_TOKEN_CHARS = 300


def _split_long_token(word: str, font: str, size: float, width: float) -> list[str]:
    """Break a token that cannot fit on a line of its own.

    A source excerpt sometimes carries a URL or an identifier longer than the
    column. Left whole it is drawn straight off the page, and the layout shrinks
    the whole slide chasing a fit it can never reach - which is how one link in a
    source page dropped a teaching slide to 10pt. Every character is kept; only
    the line it sits on changes.
    """
    pieces, current = [], ""
    for char in word:
        if current and stringWidth(current + char, font, size) > width:
            cut = max((current.rfind(x) for x in _TOKEN_BREAKS), default=-1)
            # Break after a separator when one is close enough to the end that
            # the line still carries most of its content.
            if cut >= len(current) * .5:
                pieces.append(current[:cut + 1])
                current = current[cut + 1:]
            else:
                pieces.append(current)
                current = ""
        current += char
    if current:
        pieces.append(current)
    return pieces or [word]


@lru_cache(maxsize=4096)
def _wrapped(value: str, width: float, size: float, bold: bool) -> tuple[str, ...]:
    font = "ISCARB-Bold" if bold else "ISCARB"
    lines, current = [], ""
    for word in value.split():
        parts = [word]
        if len(word) <= MAX_BREAKABLE_TOKEN_CHARS and stringWidth(word, font, size) > width:
            parts = _split_long_token(word, font, size, width)
        for part in parts:
            candidate = (current + " " + part).strip()
            if current and stringWidth(candidate, font, size) > width:
                lines.append(current)
                current = part
            else:
                current = candidate
    if current:
        lines.append(current)
    return tuple(lines)


def _overflows(blocks) -> bool:
    """A word too long to break is as unprojectable as a column that runs off the page.

    The wrapper puts an unbreakable token (a URL, a long identifier) on a line of
    its own and reports nothing, so a vertical-only fit test called the layout
    good and the export preflight then refused to render it. Measuring the drawn
    line here keeps the gate's verdict and the exporter's verdict the same.
    """
    return any(
        stringWidth(line, "ISCARB-Bold" if block.bold else "ISCARB", block.size) > block.width + .1
        for block in blocks for line in block.lines
    )


def item_layout(items, x, y, width, height, preferred=21, minimum=10):
    """Fit complete statements. Return the actual size for release validation."""
    for size in range(preferred, 5, -1):
        blocks, cursor = [], y
        for label, body in items:
            label_size = max(6, size - 3)
            label_lines = wrap(label, width, label_size, True) if label else []
            body_lines = wrap(body, width, size)
            if label_lines:
                blocks.append(Text(x, cursor, width, label_lines, label_size, GREEN, True))
                cursor += len(label_lines) * label_size * 1.22 + max(1, size * .2)
            blocks.append(Text(x, cursor, width, body_lines, size))
            cursor += len(body_lines) * size * 1.22 + max(2, size * .4)
        fits = cursor <= y + height and not _overflows(blocks)
        if fits or size == 6:
            return blocks, size, fits


def text_layout(items, x=44, y=166, width=872, height=278):
    # Balanced columns are chosen by measured content, never arbitrary item
    # truncation. A single coherent list remains a single composition.
    single = item_layout(items, x, y, width, height)
    if single[1] >= 18 or len(items) < 4:
        return single
    best = single
    for split in range(1, len(items)):
        half = (width - 36) / 2
        a = item_layout(items[:split], x, y, half, height)
        b = item_layout(items[split:], x + half + 36, y, half, height)
        candidate = (a[0] + b[0], min(a[1], b[1]), a[2] and b[2])
        if (candidate[2], candidate[1]) > (best[2], best[1]):
            best = candidate
    return best


# The source page is the teaching surface; the text column is what a slide falls
# back to when the source has no page to show. A unit teaching two source pages
# used to show neither, because one picture beside a two-page anchor looked like
# a claim about both - so half a deck became prose about slides the learner
# could have simply been shown. The picture is shown and the caption says
# exactly which page it is, which is what the provenance rule was protecting.
SOURCE_VISUAL_UNITS = frozenset({5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 15})


def exact_source_path(u, plan):
    if u.number not in SOURCE_VISUAL_UNITS or plan is None or plan.reuse_mode != "USE":
        return None
    if plan.source_slide not in anchor_slides(u.source_anchor):
        return None
    path = local_asset(plan.asset) if plan.asset else None
    return path if path and path.exists() else None


def _is_figure(plan) -> bool:
    return bool(plan is not None and plan.asset is not None and getattr(plan.asset, "source_kind", "") == FIGURE_KIND)


def source_caption(u, plan) -> str:
    """Name the page on screen, and the span the unit teaches, without conflating them."""
    coordinates = anchor_slides(u.source_anchor)
    coordinate = "SLIDE" if "SLIDE" in (u.source_anchor or "").upper() else "PAGE"
    shown = f"[P1] {coordinate} {plan.source_slide}" if plan.source_slide else clean(u.source_anchor)
    what = "source figure" if _is_figure(plan) else "original source page"
    if len(coordinates) > 1:
        return f"{shown} shown · this unit teaches {clean(u.source_anchor)} · {what}; ISCARB practice at right"
    return f"{shown} · {what}; ISCARB practice at right"


# CPIT-455's own decks give a picture the whole body of the slide and lay the
# teaching text on it in small blocks, rather than shrinking the picture into a
# column beside a paragraph. A cropped source picture is projected the same way
# here: it fills the canvas, and each block is placed where the picture has the
# least going on, so the lecturer's own labels stay readable underneath.
OVERLAY_WIDTHS = (392, 336, 288)
OVERLAY_WIDTH = OVERLAY_WIDTHS[0]
OVERLAY_PAD = 12
OVERLAY_GAP = 14
OVERLAY_MARGIN = 16
PICTURE_TOP = 44
PICTURE_BOTTOM = 506
DETAIL_COLS, DETAIL_ROWS = 96, 54


def picture_box(size) -> tuple[float, float, float, float]:
    """(x, y, width, height) in PDF points: the whole picture, never cropped.

    The picture is letterboxed into the band between the running header and the
    footer. Cropping it to fill the canvas cut the lecturer's own label boxes
    off the edges of the slide, which is the opposite of showing the source.
    """
    iw, ih = size
    band = PICTURE_BOTTOM - PICTURE_TOP
    scale = min(W / iw, band / ih)
    dw, dh = iw * scale, ih * scale
    return (W - dw) / 2, H - PICTURE_TOP - (band + dh) / 2, dw, dh


def is_source_picture(plan) -> bool:
    asset = getattr(plan, "asset", None) if plan is not None else None
    return asset is not None and getattr(asset, "source_kind", "") == PICTURE_KIND


ISCARB_PRACTICE = "ISCARB PRACTICE"


def overlay_items(u) -> list[tuple[str, str]]:
    items = [("ENGINEERING QUESTION", clean(u.engineering_question))]
    items += [split_item(x) for x in u.pedagogy_content if clean(x)]
    items += [("YOUR TASK", clean(u.student_action))]
    # Every item the side-by-side layout would show has to appear here too: a
    # picture slide is a different arrangement of the unit, never less of it.
    # Consecutive unlabelled fragments are one teaching move written in two
    # sentences, so they share one block instead of repeating a made-up label.
    merged: list[list[str]] = []
    for label, body in items:
        if not body:
            continue
        if not label and merged and merged[-1][0] == ISCARB_PRACTICE:
            merged[-1][1] += " " + body
            continue
        merged.append([label or ISCARB_PRACTICE, body])
    return [(label, body) for label, body in merged]


def _overlay_blocks(items, size, width):
    """Label/body line sets and their heights, or None when one runs too long."""
    blocks = []
    inner = width - 2 * OVERLAY_PAD
    for label, body in items:
        label_lines = wrap(label, inner, size - 4, True)
        body_lines = wrap(body, inner, size)
        if len(label_lines) > 1 or len(body_lines) > 4:
            return None
        height = 2 * OVERLAY_PAD + (size - 4) * 1.24 + len(body_lines) * size * 1.24
        blocks.append((label_lines, body_lines, size, height))
    return blocks


@lru_cache(maxsize=32)
def _detail_grid(path: str, mtime: int):
    """Per-cell edge energy of the picture as it lands on the canvas.

    High where the picture (or the lecturer's own label boxes) carries detail,
    low where a teaching block can sit without hiding anything.
    """
    try:
        with Image.open(path) as im:
            box = picture_box(im.size)
            canvas_image = Image.new("L", (W, H), 0)
            canvas_image.paste(im.convert("L").resize((max(1, round(box[2])), max(1, round(box[3])))),
                               (round(box[0]), round(H - box[1] - box[3])))
            # Text strokes are thin: averaged straight into a cell they vanish
            # into the picture's own noise. Spreading each edge first keeps a
            # caption or a label box legible to the measurement.
            edges = canvas_image.filter(ImageFilter.FIND_EDGES).filter(ImageFilter.MaxFilter(9))
            cells = edges.resize((DETAIL_COLS, DETAIL_ROWS), Image.BOX)
            return cells.tobytes()
    except Exception:
        return None


# Photographs carry edges everywhere, so an average would rate a face and a
# paragraph alike. What must not be covered is concentrated detail - the
# lecturer's own label boxes, a caption, a face - which shows up as cells well
# above the picture's own noise floor.
BUSY_CELL = 70
# When the picture's own words are PDF text we know exactly where they are and
# simply never land on them. When they are baked into the image - a screenshot
# of a table, a quote card - we cannot see them, so only a picture that is
# quiet all over may carry teaching text, and only in its quietest places.
MAX_BUSY_ON_AN_UNREADABLE_PICTURE = .30
MAX_BUSY_UNDER_AN_UNREADABLE_BLOCK = .22
# A block belongs on the picture only where the picture is genuinely quiet.
# Past this it would hide the thing the slide exists to show, and the unit keeps
# the side-by-side layout, where the picture is whole and the text is beside it.
MAX_BUSY_UNDER_A_BLOCK = .25


def _cell_energy(grid, x, y, width, height) -> tuple[float, float]:
    cw, ch = W / DETAIL_COLS, H / DETAIL_ROWS
    c0, c1 = max(0, int(x // cw)), min(DETAIL_COLS, int(-(-(x + width) // cw)))
    r0, r1 = max(0, int(y // ch)), min(DETAIL_ROWS, int(-(-(y + height) // ch)))
    cells = [grid[r * DETAIL_COLS + c] for r in range(r0, r1) for c in range(c0, c1)]
    if not cells:
        return 1.0, 255.0
    busy = sum(1 for v in cells if v >= BUSY_CELL) / len(cells)
    return busy, sum(cells) / len(cells)


def source_text_rects(plan, source) -> list[tuple[float, float, float, float]]:
    """The lecturer's own text on this picture, in canvas points."""
    boxes = getattr(getattr(plan, "asset", None), "text_boxes", ()) or ()
    if not boxes:
        return []
    try:
        with Image.open(source) as im:
            x, y, dw, dh = picture_box(im.size)
    except Exception:
        return []
    top = H - y - dh
    return [(x + a * dw, top + b * dh, x + c * dw, top + d * dh) for a, b, c, d in boxes]


def _place_overlay(grid, blocks, bounds, width, forbidden=(), quiet=None):
    """Position each block over the quietest free part of the picture."""
    left, top, right, bottom = bounds
    placed, rects, floor = [], [], top
    for index, (label_lines, body_lines, size, height) in enumerate(blocks):
        # Whatever is still to be taught has to fit below this block, or the
        # last one - always the student's task - would have nowhere to go.
        rest = sum(h for *_, h in blocks[index + 1:]) + OVERLAY_GAP * len(blocks[index + 1:])
        best = None
        for y in range(int(floor), int(bottom - rest - height) + 1, 10):
            for x in range(int(left), int(right - width) + 1, 20):
                if any(x < rx + width + OVERLAY_GAP and rx < x + width + OVERLAY_GAP
                       and y < ry + rh + OVERLAY_GAP and ry < y + height + OVERLAY_GAP for rx, ry, rh in rects):
                    continue
                # The lecturer's own labels are the point of the slide; a
                # teaching block never lands on one.
                if any(x < fx1 and fx0 < x + width and y < fy1 and fy0 < y + height
                       for fx0, fy0, fx1, fy1 in forbidden):
                    continue
                busy, mean = _cell_energy(grid, x, y, width, height)
                if best is None or (busy, mean) < best[:2]:
                    best = (busy, mean, x, y)
        if best is None or best[0] > quiet:
            # Every position would bury something the lecturer put on the page.
            # This picture keeps its own surface; the unit falls back to showing
            # it beside the teaching text.
            return None
        _, _, x, y = best
        floor = y
        rects.append((x, y, height))
        placed.append((label_lines, body_lines, size, x, y, height, width))
    return placed


def overlay_bounds(source) -> tuple[float, float, float, float] | None:
    """The band a teaching block may occupy: the picture and its own margins.

    A picture that does not fill the band leaves an empty ground beside or below
    it, and a block is happier there than on the picture - which is what the
    quietest-place search finds on its own once the ground is in range.
    """
    try:
        with Image.open(source) as im:
            picture_box(im.size)
    except Exception:
        return None
    return (OVERLAY_MARGIN, PICTURE_TOP + OVERLAY_MARGIN,
            W - OVERLAY_MARGIN, PICTURE_BOTTOM - OVERLAY_MARGIN)


def overlay_layout(items, source, plan=None):
    """The teaching blocks laid on the picture, or None when they do not fit."""
    bounds = overlay_bounds(source) if items else None
    grid = _detail_grid(str(source), Path(source).stat().st_mtime_ns) if bounds else None
    if not grid:
        return None
    room = bounds[3] - bounds[1]
    forbidden = source_text_rects(plan, source)
    quiet = MAX_BUSY_UNDER_A_BLOCK
    if not forbidden:
        if sum(1 for v in grid if v >= BUSY_CELL) / len(grid) > MAX_BUSY_ON_AN_UNREADABLE_PICTURE:
            return None
        quiet = MAX_BUSY_UNDER_AN_UNREADABLE_BLOCK
    for width in OVERLAY_WIDTHS:
        if bounds[2] - bounds[0] < width:
            continue
        for size in (17, 16, 15, 14):
            blocks = _overlay_blocks(items, size, width)
            if not blocks or sum(h for *_, h in blocks) + OVERLAY_GAP * (len(blocks) - 1) > room:
                continue
            placed = _place_overlay(grid, blocks, bounds, width, forbidden, quiet)
            if placed:
                return placed
    return None


def unit_layout(bp, u, plan=None):
    source = exact_source_path(u, plan)
    if source:
        items = [("ENGINEERING QUESTION", clean(u.engineering_question)),
                 *[split_item(x) for x in u.pedagogy_content if clean(x)],
                 *contextual_items(u),
                 ("YOUR TASK", clean(u.student_action))]
        blocks, size, fits = item_layout(items, 636, 65, 284, 433, preferred=17)
        return blocks, size, fits, source
    blocks, size, fits = text_layout(teaching_items(bp, u))
    return blocks, size, fits, None


# The footer prints the anchor on one line at 8pt. A unit citing every page of a
# chapter overran it, and only the export preflight noticed: the deck cleared
# every gate and then refused to produce a file.
MAX_ANCHOR_WIDTH = 790


def readability_problems(bp: Blueprint) -> dict[int, str]:
    # Source availability is resolved at export too; the gate checks the
    # conservative text representation rather than trusting an unavailable URL.
    problems = {}
    for u in bp.units:
        if stringWidth(clean(u.source_anchor), "ISCARB", 8) > MAX_ANCHOR_WIDTH:
            problems[u.number] = "Cite the source span in one printable phrase; the anchor line does not fit the footer."
        if len(wrap(u.engineering_question,872,13))>2 or len(wrap("YOUR TASK "+u.student_action,872,13))>2:
            problems[u.number] = "Shorten the engineering question/task to two lines without changing its purpose."
        if u.number == 19:
            if rubric_layout(bp)[1] < 12:
                problems[u.number] = "Condense repeated rubric wording; preserve six criteria and all four descriptors at 12pt or larger."
            continue
        _, size, fits, _ = unit_layout(bp, u)
        if not fits or size < 16:
            problems[u.number] = "Condense duplicate wording and long scaffolds, preserving every source fact/list/example; body text must fit at 16pt or larger."
    return problems


def readable_text_contract(bp: Blueprint) -> bool:
    return not readability_problems(bp)


def preflight_layout(bp, source_root=None):
    """Shared PDF/PPTX/preview geometry check; never return a clipped artifact."""
    plans = list(plans_for_blueprint_v42(bp, source_root=source_root))
    failures = []
    for u, plan in zip(bp.units, plans):
        blocks, _, fits, source = unit_layout(bp, u, plan)
        bottom = 498 if source else 448
        if u.number == 19 and not source:
            blocks = rubric_layout(bp)[0]
            fits = True
        unsafe = not fits or any(
            b.y + len(b.lines) * b.size * 1.22 > bottom
            or any(stringWidth(line, "ISCARB-Bold" if b.bold else "ISCARB", b.size) > b.width + .1 for line in b.lines)
            for b in blocks)
        if not source:
            unsafe = unsafe or len(wrap(u.engineering_question, 872, 13)) > 2 or len(wrap("YOUR TASK " + u.student_action, 872, 13)) > 2
            frame_blocks = [title_block(bp.lecture_title if u.number == 1 else u.title),
                Text(44, 120, 872, wrap(u.engineering_question, 872, 13), 13),
                Text(44, 470, 872, wrap("YOUR TASK " + u.student_action, 872, 13), 13)]
            unsafe = unsafe or any(any(stringWidth(line, "ISCARB-Bold" if b.bold else "ISCARB", b.size) > b.width + .1 for line in b.lines) for b in frame_blocks)
        unsafe = unsafe or stringWidth(clean(u.source_anchor), "ISCARB", 8) > MAX_ANCHOR_WIDTH
        if unsafe:
            failures.append(u.number)
    if failures:
        raise PresenterLayoutError("Presenter cannot fit units " + ", ".join(map(str, failures)) + ". No clipped file was generated. Shorten/restructure these units; the original source and blueprint remain available.")
    return plans


def title_block(title):
    for size in range(34, 15, -1):
        lines = wrap(clean(title), 872, size, True)
        if len(lines) * size * 1.22 <= 65:
            return Text(44, 46, 872, lines, size, GREEN, True)
    return Text(44, 46, 872, lines, size, GREEN, True)


def _text(c, block: Text):
    c.setFont("ISCARB-Bold" if block.bold else "ISCARB", block.size)
    c.setFillColor(colors.HexColor(block.color))
    for i, line in enumerate(block.lines):
        c.drawString(block.x, H - block.y - block.size - i * block.size * 1.22, line)


def _line(c, y, color=GOLD, x=44, width=872):
    c.setStrokeColor(colors.HexColor(color))
    c.setLineWidth(.8)
    c.line(x, H-y, x+width, H-y)


def _frame(c, bp, u):
    _text(c, Text(44, 20, 760, [eyebrow(u)], 10, GREEN, True))
    _text(c, title_block(bp.lecture_title if u.number == 1 else u.title))
    _text(c, Text(44, 120, 872, wrap(u.engineering_question, 872, 13), 13, MUTED))
    _line(c, 111)
    _line(c, 460)
    task = clean(u.student_action)
    lines = wrap("YOUR TASK  " + task, 872, 13)
    _text(c, Text(44, 470, 872, lines, 13, INK))
    _text(c, Text(44, 518, 800, [clean(u.source_anchor)], 8, MUTED))
    _text(c, Text(865, 517, 60, [f"{u.number:02d} / 20"], 9, GREEN, True))


def rubric_layout(bp):
    widths = [180, 173, 173, 173, 173]
    headers = ["CAPABILITY", "DISTINGUISHED", "READY", "DEVELOPING", "NOT YET READY"]
    blocks=[]
    x = 44
    for title, width in zip(headers, widths):
        blocks.append(Text(x+5, 155, width-10, wrap(title, width-10, 10, True), 10, GREEN, True))
        x += width
    for size in range(14,5,-1):
        body=[]; rules=[]; y=181
        for row in bp.rubric_criteria:
            values=[row.criterion,row.distinguished,row.ready,row.developing,row.not_yet_ready]
            cells=[wrap(value,width-10,size) for value,width in zip(values,widths)]
            x=44
            for lines,width in zip(cells,widths):
                body.append(Text(x+5,y,width-10,lines,size))
                x+=width
            y+=max(max(1,len(lines))*size*1.22 for lines in cells)+10
            rules.append(y-4)
        if y<=448:
            break
    return blocks+body,size,rules


def _rubric(c,bp):
    blocks,_,rules=rubric_layout(bp)
    for block in blocks: _text(c,block)
    for y in rules: _line(c,y,"#DCE5DF")


def _draw_source_picture(c, u, plan, source, placed, release_state):
    """The source picture is the slide; the teaching text sits on it."""
    c.setFillColor(colors.HexColor("#0C1F1A"))
    c.rect(0, 0, W, H, stroke=0, fill=1)
    with Image.open(source) as im:
        x, y, dw, dh = picture_box(im.size)
        c.drawImage(ImageReader(im), x, y, width=dw, height=dh)
    for label_lines, body_lines, size, x, y, height, width in placed:
        c.setFillColor(colors.HexColor("#0B3B2E"))
        c.roundRect(x, H - y - height, width, height, 7, stroke=0, fill=1)
        _text(c, Text(x + OVERLAY_PAD, y + OVERLAY_PAD, width - 2 * OVERLAY_PAD, label_lines, size - 4, "#EBC77A", True))
        _text(c, Text(x + OVERLAY_PAD, y + OVERLAY_PAD + (size - 4) * 1.24 + 2,
                      width - 2 * OVERLAY_PAD, body_lines, size, "#FFFFFF"))
    _text(c, Text(36, 17, 760, [eyebrow(u)], 10, "#EBC77A", True))
    _text(c, Text(36, 520, 780, [source_caption(u, plan)], 8, "#D9CBB6"))
    _text(c, Text(865, 517, 60, [f"{u.number:02d} / 20"], 9, "#EBC77A", True))
    _text(c, Text(805, 18, 120, ["VERIFIED RELEASE" if release_state.upper() == "READY" else "REVIEW DRAFT"], 8, "#D9CBB6", True))


def _draw_page(c, bp, u, plan, release_state="REVIEW"):
    blocks,size,fits,source = unit_layout(bp,u,plan)
    if source and is_source_picture(plan):
        placed = overlay_layout(overlay_items(u), source, plan)
        if placed:
            _draw_source_picture(c, u, plan, source, placed, release_state)
            return
    if source:
        _text(c, Text(36, 17, 850, [eyebrow(u)], 10, GREEN, True))
        _line(c, 42, x=36, width=884)
        with Image.open(source) as im:
            iw,ih=im.size
            scale=min(576/iw,454/ih)
            dw,dh=iw*scale,ih*scale
            c.drawImage(ImageReader(im),36+(576-dw)/2,H-54-dh,width=dw,height=dh)
        for block in blocks:
            _text(c,block)
        _text(c,Text(36,520,800,[source_caption(u,plan)],8,MUTED))
        _text(c,Text(865,517,60,[f"{u.number:02d} / 20"],9,GREEN,True))
        _text(c,Text(805,18,120,["VERIFIED RELEASE" if release_state.upper()=="READY" else "REVIEW DRAFT"],8,MUTED,True))
        return
    _frame(c,bp,u)
    _text(c,Text(805,20,120,["VERIFIED RELEASE" if release_state.upper()=="READY" else "REVIEW DRAFT"],8,MUTED,True))
    if u.number == 19:
        _rubric(c,bp)
        return
    for block in blocks:
        _text(c,block)
    if not fits or size < 16:
        _text(c, Text(44, 502, 790, ["REVIEW: dense content — simplify complete statements before classroom release."], 8, "#9D442E"))


def export_presenter_pdf(bp: Blueprint, out: Path, source_root=None, release_state="REVIEW") -> Path:
    out=Path(out)
    plans = preflight_layout(bp, source_root)
    c=canvas.Canvas(str(out),pagesize=(W,H))
    c.setTitle(bp.lecture_title)
    c.setAuthor("ISCARB Faculty Studio")
    for u,plan in zip(bp.units,plans):
        _draw_page(c,bp,u,plan,release_state)
        c.showPage()
    c.save()
    return out


def render_presenter_preview(bp: Blueprint, release_state="BLOCKED", source_root=None) -> str:
    # The browser previews the actual PDF surface: one projection, no separate
    # HTML layout that can falsely pass while downloaded slides clip content.
    plans = preflight_layout(bp, source_root)
    buf=io.BytesIO()
    c=canvas.Canvas(buf,pagesize=(W,H))
    for u,plan in zip(bp.units,plans):
        _draw_page(c,bp,u,plan,release_state)
        c.showPage()
    c.save()
    import fitz
    doc=fitz.open(stream=buf.getvalue(),filetype="pdf")
    pages=[]
    for u,page in zip(bp.units,doc):
        png=page.get_pixmap(matrix=fitz.Matrix(1.5,1.5)).tobytes("png")
        pages.append(f'<section class="slide" id="unit-{u.number}"><img alt="{html.escape(u.title,quote=True)}" src="data:image/png;base64,{base64.b64encode(png).decode()}" /></section>')
    doc.close()
    return ('<!doctype html><html lang="en"><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">'
            '<title>'+html.escape(bp.lecture_title)+'</title><style>body{margin:0;background:#e9eeeb;font:15px system-ui;color:#183a30}'
            'header{padding:18px 4%;position:sticky;top:0;background:#fff;display:flex;justify-content:space-between;z-index:2}'
            '.slide{max-width:1200px;margin:24px auto;background:#fff;box-shadow:0 8px 32px #173f3020}.slide img{width:100%;display:block}'
            '@media print{header{display:none}.slide{margin:0;break-after:page;box-shadow:none}@page{size:16in 9in;margin:0}}</style>'
            '<header><strong>'+html.escape(bp.lecture_title)+'</strong><span>'+html.escape(release_state)+' · 20 units</span></header>'
            +''.join(pages)+'</html>')


def export_presenter_pptx(bp: Blueprint, out: Path, source_root=None, release_state="REVIEW") -> Path:
    # Native editable text and source images use the exact PDF layout plan.
    plans = preflight_layout(bp, source_root)
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    prs=Presentation()
    prs.slide_width=Inches(13.333333)
    prs.slide_height=Inches(7.5)

    def add(slide, block):
        shape=slide.shapes.add_textbox(Inches(block.x/72),Inches(block.y/72),Inches(block.width/72),Inches(max(20,len(block.lines)*block.size*1.22+4)/72))
        tf=shape.text_frame
        tf.clear()
        tf.margin_top=tf.margin_bottom=tf.margin_left=tf.margin_right=0
        for i,line in enumerate(block.lines):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.text=line
            p.font.size=Pt(block.size)
            p.font.name="Arial"
            p.font.bold=block.bold
            p.font.color.rgb=RGBColor.from_string(block.color.lstrip("#"))
            p.space_before=p.space_after=Pt(0)
            p.line_spacing=1.22

    for u,plan in zip(bp.units,plans):
        slide=prs.slides.add_slide(prs.slide_layouts[6])
        add(slide,Text(805,20,120,["VERIFIED RELEASE" if release_state.upper()=="READY" else "REVIEW DRAFT"],8,MUTED,True))
        blocks,_,_,source=unit_layout(bp,u,plan)
        placed = overlay_layout(overlay_items(u), source, plan) if source and is_source_picture(plan) else None
        if placed:
            # Same picture slide as the PDF: the source picture fills the band
            # and the teaching text sits in the places it leaves free.
            ground=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
            ground.fill.solid(); ground.fill.fore_color.rgb=RGBColor.from_string("0C1F1A"); ground.line.fill.background()
            with Image.open(source) as im: px,py,dw,dh=picture_box(im.size)
            slide.shapes.add_picture(str(source),Inches(px/72),Inches((H-py-dh)/72),width=Inches(dw/72),height=Inches(dh/72))
            for label_lines,body_lines,size,x,y,height,width in placed:
                panel=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x/72),Inches(y/72),Inches(width/72),Inches(height/72))
                panel.fill.solid(); panel.fill.fore_color.rgb=RGBColor.from_string("0B3B2E"); panel.line.fill.background()
                panel.adjustments[0]=.08
                add(slide,Text(x+OVERLAY_PAD,y+OVERLAY_PAD,width-2*OVERLAY_PAD,label_lines,size-4,"#EBC77A",True))
                add(slide,Text(x+OVERLAY_PAD,y+OVERLAY_PAD+(size-4)*1.24+2,width-2*OVERLAY_PAD,body_lines,size,"#FFFFFF"))
            add(slide,Text(36,17,760,[eyebrow(u)],10,"#EBC77A",True))
            add(slide,Text(36,520,780,[source_caption(u,plan)],8,"#D9CBB6"))
            add(slide,Text(865,517,60,[f"{u.number:02d} / 20"],9,"#EBC77A",True))
            slide.notes_slide.notes_text_frame.text = "[Sources]\n"+clean(u.source_anchor)+"\n[/Sources]\n"+"\n".join(u.core_content+u.pedagogy_content)
            continue
        if source:
            add(slide,Text(36,17,850,[eyebrow(u)],10,GREEN,True))
            with Image.open(source) as im: iw,ih=im.size
            scale=min(576/iw,454/ih)
            dw,dh=iw*scale,ih*scale
            slide.shapes.add_picture(str(source),Inches((36+(576-dw)/2)/72),Inches(54/72),width=Inches(dw/72),height=Inches(dh/72))
            for block in blocks: add(slide,block)
            add(slide,Text(36,520,800,[source_caption(u,plan)],8,MUTED))
            add(slide,Text(865,517,60,[f"{u.number:02d} / 20"],9,GREEN,True))
            slide.notes_slide.notes_text_frame.text = "[Sources]\n"+clean(u.source_anchor)+"\n[/Sources]\n"+"\n".join(u.core_content+u.pedagogy_content)
            continue
        add(slide,Text(44,20,872,[eyebrow(u)],10,GREEN,True))
        add(slide,title_block(bp.lecture_title if u.number==1 else u.title))
        add(slide,Text(44,120,872,wrap(u.engineering_question,872,13),13,MUTED))
        if u.number==19:
            for block in rubric_layout(bp)[0]: add(slide,block)
        else:
            blocks,_,_,source=unit_layout(bp,u,plan)
            if source:
                with Image.open(source) as im: iw,ih=im.size
                scale=min(610/iw,298/ih)
                dw,dh=iw*scale,ih*scale
                slide.shapes.add_picture(str(source),Inches((44+(610-dw)/2)/72),Inches(146/72),width=Inches(dw/72),height=Inches(dh/72))
            for block in blocks: add(slide,block)
        add(slide,Text(44,470,872,wrap("YOUR TASK  "+u.student_action,872,13),13))
        add(slide,Text(44,518,790,[clean(u.source_anchor)],8,MUTED))
        add(slide,Text(865,517,60,[f"{u.number:02d} / 20"],9,GREEN,True))
        slide.notes_slide.notes_text_frame.text = (
            "[Sources]\n"+clean(u.source_anchor)+"\n[/Sources]\nPRIMARY CORE\n"+"\n".join(u.core_content)
            +"\nISCARB PEDAGOGY\n"+"\n".join(u.pedagogy_content)+"\nTAKEAWAY\n"+u.takeaway)
    out=Path(out)
    prs.save(str(out))
    return out
